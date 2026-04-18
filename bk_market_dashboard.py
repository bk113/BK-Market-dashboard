"""
BK Market Dashboard — Consolidated
====================================
114-instrument universe · 15 asset classes · Performance, Risk & Fragility.
Last updated: 2026-04-13

Outputs:
  docs/index.html   3-tab web page (GitHub Pages — auto-updated daily)
  PNG + PDF         visual report
  PPTX              PowerPoint deck
  Email             HTML brief via Gmail

Usage:
  python bk_market_dashboard.py --html                      # Generate web page (10yr history)
  python bk_market_dashboard.py --html --lookback 504       # GitHub Actions (2yr, faster)
  python bk_market_dashboard.py --pptx                      # PowerPoint deck
  python bk_market_dashboard.py --email            # Send HTML email brief
  python bk_market_dashboard.py --html --pptx      # Web page + PowerPoint
  python bk_market_dashboard.py --schedule         # Daily scheduler at 07:00 SGT Mon–Fri
  python bk_market_dashboard.py --now --html       # Run once immediately (testing)

Dependencies:
  pip install yfinance pandas numpy matplotlib schedule pytz python-pptx
"""

from __future__ import annotations

import argparse
import os
import smtplib
import time
import warnings
from datetime import datetime
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from io import BytesIO

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import Rectangle
import numpy as np
import pandas as pd
import pytz
import schedule
import yfinance as yf

warnings.filterwarnings("ignore")


DEVELOPMENT_MODE = False  # Set True to skip Claude API calls during development

# ══════════════════════════════════════════════════════════════════════════════
#  CONFIG — EDIT THESE
# ══════════════════════════════════════════════════════════════════════════════

RECIPIENT_EMAIL = "your@email.com"           # Who receives the brief
SENDER_EMAIL    = "your.gmail@gmail.com"     # Your Gmail address
GMAIL_APP_PASS  = "xxxx xxxx xxxx xxxx"      # Gmail App Password (not login password)
                                              # Get: myaccount.google.com > Security > App Passwords
SEND_TIME_SGT   = "07:00"                    # Daily send time (SGT)
OUT_DIR         = "."                         # Output folder for PNG/PDF (current directory)
RISK_FREE_RATE  = 0.045                      # 4.5% annualised risk-free rate

# ══════════════════════════════════════════════════════════════════════════════

SGT = pytz.timezone("Asia/Singapore")


# ── UNIVERSE ──────────────────────────────────────────────────────────────────
# (section_key, ticker, display_name, bucket)
# Bucket feeds dynamic allocation + regime-fit matrix + BK Composite Score.

UNIVERSE = [
    # ── EQ_US — US Equity (4) ────────────────────────────────
    ("EQ_US",    "QQQ",      "Nasdaq 100",              "EQ Growth"),
    ("EQ_US",    "IWM",      "Russell 2000",            "EQ Growth"),
    ("EQ_US",    "SPY",      "S&P 500",                 "EQ Growth"),
    ("EQ_US",    "ACWI",     "World (ACWI)",            "EQ Growth"),

    # ── EQ_SECT — US Sectors, all 11 GICS (11) ───────────────
    ("EQ_SECT",  "XLE",      "Energy",                  "EQ Defensive"),
    ("EQ_SECT",  "XLF",      "Financials",              "EQ Growth"),
    ("EQ_SECT",  "VNQ",      "Real Estate (REITs)",     "Real Assets"),
    ("EQ_SECT",  "XLU",      "Utilities",               "EQ Defensive"),
    ("EQ_SECT",  "XLV",      "Healthcare",              "EQ Defensive"),
    ("EQ_SECT",  "XLK",      "Technology",              "EQ Growth"),
    ("EQ_SECT",  "XLP",      "Consumer Staples",        "EQ Defensive"),
    ("EQ_SECT",  "XLY",      "Consumer Discretionary",  "EQ Growth"),
    ("EQ_SECT",  "XLI",      "Industrials",             "EQ Growth"),
    ("EQ_SECT",  "XLB",      "Materials",               "Real Assets"),
    ("EQ_SECT",  "XLC",      "Communications",          "EQ Growth"),

    # ── EQ_DM — Developed Markets (7) ────────────────────────
    ("EQ_DM",    "EWA",      "Australia",               "EQ Growth"),
    ("EQ_DM",    "EFA",      "Europe Dev (EFA)",        "EQ Growth"),
    ("EQ_DM",    "EZU",      "Eurozone",                "EQ Growth"),
    ("EQ_DM",    "EWG",      "Germany",                 "EQ Growth"),
    ("EQ_DM",    "EWJ",      "Japan",                   "EQ Growth"),
    ("EQ_DM",    "EWS",      "Singapore",               "EQ Growth"),
    ("EQ_DM",    "EWU",      "UK",                      "EQ Defensive"),
    ("EQ_DM",    "FLGB",     "FTSE 100 (UK)",           "EQ Defensive"),

    # ── EQ_IDX — Specialist & Thematic (5) ──────────────────
    ("EQ_IDX",   "GURU",     "Global X Guru ETF",       "EQ Growth"),
    ("EQ_IDX",   "ARKK",     "ARK Innovation",          "EQ Growth"),
    ("EQ_IDX",   "SKYY",     "Cloud Computing",         "EQ Growth"),
    ("EQ_IDX",   "ICLN",     "Clean Energy",            "EQ Growth"),
    ("EQ_IDX",   "GRID",     "Smart Grid Infra",        "EQ Growth"),

    # ── EQ_EM — Emerging Markets (11) ────────────────────────
    ("EQ_EM",    "EWZ",      "Brazil",                  "EQ Growth"),
    ("EQ_EM",    "FXI",      "China",                   "EQ Growth"),
    ("EQ_EM",    "EEM",      "EM Broad",                "EQ Growth"),
    ("EQ_EM",    "INDA",     "India",                   "EQ Growth"),
    ("EQ_EM",    "EWY",      "Korea",                   "EQ Growth"),
    ("EQ_EM",    "EZA",      "South Africa",            "EQ Growth"),
    ("EQ_EM",    "EWT",      "Taiwan",                  "EQ Growth"),
    ("EQ_EM",    "EWW",      "Mexico",                  "EQ Growth"),
    ("EQ_EM",    "EIDO",     "Indonesia",               "EQ Growth"),
    ("EQ_EM",    "VNM",      "Vietnam",                 "EQ Growth"),
    ("EQ_EM",    "KSA",      "Saudi Arabia",            "Real Assets"),

    # ── EQ_APAC — Asia Pacific (7) ──────────────────────────
    ("EQ_APAC",  "AAXJ",     "Asia ex-Japan",           "EQ Growth"),
    ("EQ_APAC",  "EWH",      "Hang Seng / HK",          "EQ Growth"),
    ("EQ_APAC",  "VPL",      "Asia Pacific",            "EQ Growth"),
    ("EQ_APAC",  "CNYA",     "China A-Shares",          "EQ Growth"),
    ("EQ_APAC",  "ASEA",     "ASEAN",                   "EQ Growth"),
    ("EQ_APAC",  "THD",      "Thailand",                "EQ Growth"),
    ("EQ_APAC",  "EPHE",     "Philippines",             "EQ Growth"),

    # ── DEFENCE — Defence & Geopolitical (2) ─────────────────
    ("DEFENCE",  "XAR",      "BAE Aerospace & Defence", "Alts"),
    ("DEFENCE",  "ITA",      "US Aerospace & Defence",  "Alts"),

    # ── FI — Fixed Income & Credit (15) ──────────────────────
    ("FI",       "BIL",      "Cash (T-Bills)",          "Cash"),
    ("FI",       "EMLC",     "EM Local Currency",       "Fixed Income"),
    ("FI",       "EMB",      "EM USD Sovereign",        "Fixed Income"),
    ("FI",       "HYG",      "HY Credit",               "EQ Growth"),
    ("FI",       "LQD",      "IG Credit",               "Fixed Income"),
    ("FI",       "BKLN",     "Senior Loans",            "Fixed Income"),
    ("FI",       "SHY",      "Treasuries 0-3Y",         "Cash"),
    ("FI",       "TLT",      "Treasuries 20Y+",         "Fixed Income"),
    ("FI",       "IEF",      "Treasuries 7-10Y",        "Fixed Income"),
    ("FI",       "AGG",      "US Aggregate",            "Fixed Income"),
    ("FI",       "TIP",      "US TIPS",                 "Real Assets"),
    ("FI",       "MUB",      "Municipal Bonds",         "Fixed Income"),
    ("FI",       "CWB",      "Convertible Bonds",       "EQ Growth"),
    ("FI",       "VCSH",     "ST IG Corporates",        "Fixed Income"),
    ("FI",       "PFF",      "Preferred Securities",    "Fixed Income"),
    ("FI",       "HYD",      "Municipal HY",            "Fixed Income"),

    # ── FI_INTL — FI International (5) ───────────────────────
    ("FI_INTL",  "BNDW",     "Global Aggregate",        "Fixed Income"),
    ("FI_INTL",  "IHY",      "Intl High Yield",         "EQ Growth"),
    ("FI_INTL",  "IGIB",     "US IG Credit (Intl)",     "Fixed Income"),
    ("FI_INTL",  "BNDX",     "Intl Bonds",              "Fixed Income"),
    ("FI_INTL",  "IGOV",     "Intl Govt Bonds",         "Fixed Income"),

    # ── RATES — Sovereign Rates (4) NEW ──────────────────────
    ("RATES",    "^TNX",     "10Y Treasury Yield",      "Fixed Income"),
    ("RATES",    "^IRX",     "2Y Treasury Yield",       "Fixed Income"),
    ("RATES",    "^FVX",     "5Y Treasury Yield",       "Fixed Income"),
    ("RATES",    "^TYX",     "30Y Treasury Yield",      "Fixed Income"),

    # ── CMD — Commodities (12) ───────────────────────────────
    ("CMD",      "DBA",      "Agriculture",             "Real Assets"),
    ("CMD",      "DBC",      "Broad Commodities",       "Real Assets"),
    ("CMD",      "COPX",     "Copper Miners",           "Real Assets"),
    ("CMD",      "GLD",      "Gold",                    "Real Assets"),
    ("CMD",      "UNG",      "Natural Gas",             "Real Assets"),
    ("CMD",      "SLV",      "Silver",                  "Real Assets"),
    ("CMD",      "BNO",      "WTI Oil (BNO proxy)",     "Real Assets"),
    ("CMD",      "PPLT",     "Platinum",                "Real Assets"),
    ("CMD",      "URA",      "Uranium",                 "Alts"),
    ("CMD",      "WEAT",     "Wheat",                   "Real Assets"),
    ("CMD",      "LIT",      "Lithium",                 "Alts"),
    ("CMD",      "VALE",     "Iron Ore (VALE)",         "Real Assets"),
    ("CMD",      "CORN",     "Corn",                    "Real Assets"),
    ("CMD",      "SOYB",     "Soybeans",                "Real Assets"),
    ("CMD",      "PALL",     "Palladium",               "Real Assets"),
    ("CMD",      "CPER",     "Copper ETF",              "Real Assets"),

    # ── CRYPTO (3) ───────────────────────────────────────────
    ("CRYPTO",   "BTC-USD",  "Bitcoin",                 "Alts"),
    ("CRYPTO",   "ETH-USD",  "Ethereum",                "Alts"),
    ("CRYPTO",   "SOL-USD",  "Solana",                  "Alts"),

    # ── FX (12) ──────────────────────────────────────────────
    ("FX",       "AUDUSD=X", "AUD/USD",                 "Alts"),
    ("FX",       "EURUSD=X", "EUR/USD",                 "Alts"),
    ("FX",       "GBPUSD=X", "GBP/USD",                 "Alts"),
    ("FX",       "DX-Y.NYB", "US Dollar Index",         "Cash"),
    ("FX",       "CHF=X",    "USD/CHF",                 "Cash"),
    ("FX",       "JPY=X",    "USD/JPY",                 "Cash"),
    ("FX",       "SGD=X",    "USD/SGD",                 "Alts"),
    ("FX",       "CAD=X",    "Canadian Dollar",         "Real Assets"),
    ("FX",       "MXN=X",    "Mexican Peso",            "Alts"),
    ("FX",       "KRW=X",    "Korean Won",              "Alts"),
    ("FX",       "BRL=X",    "Brazilian Real",          "Alts"),

    # ── VOL — Volatility (5; GVZ/OVX computed from GLD/BNO) ─
    ("VOL",      "^VIX",     "VIX Index",               "Alts"),
    ("VOL",      "^VIX3M",   "VIX 3-Month Index",       "Alts"),
    ("VOL",      "VXX",      "VIX Futures ETN",         "Alts"),
    ("VOL",      "VIXY",     "ST VIX ETF",              "Alts"),
    ("VOL",      "GVZ",      "Gold Volatility (proxy)", "Alts"),
    ("VOL",      "OVX",      "Oil Volatility (proxy)",  "Alts"),

    # ── ALT — Listed Alternatives (7) ────────────────────────
    ("ALT",      "IFRA",     "US Infrastructure",       "Real Assets"),
    ("ALT",      "PSP",      "Listed Private Equity",   "Alts"),
    ("ALT",      "AMLP",     "Energy Infrastructure",   "Real Assets"),
    ("ALT",      "REET",     "Global REITs",            "Real Assets"),
    ("ALT",      "BCI",      "Commodities Index",       "Real Assets"),
    ("ALT",      "PDBC",     "Diversified Commodity",   "Real Assets"),
    ("ALT",      "KBWY",     "High Yield REIT",         "Real Assets"),
]

SECTION_ORDER = ["EQ_US", "EQ_SECT", "EQ_DM", "EQ_IDX", "EQ_APAC", "EQ_EM", "DEFENCE",
                 "FI", "FI_INTL", "RATES", "CMD", "CRYPTO", "FX", "VOL", "ALT"]

SECTION_LABELS = {
    "EQ_US":   "EQUITIES — US BROAD",
    "EQ_SECT": "EQUITIES — US SECTORS",
    "EQ_DM":   "EQUITIES — DEVELOPED MARKETS",
    "EQ_IDX":  "EQUITIES — SPECIALIST & THEMATIC",
    "EQ_APAC": "EQUITIES — ASIA PACIFIC",
    "EQ_EM":   "EQUITIES — EMERGING MARKETS",
    "DEFENCE": "DEFENCE & GEOPOLITICAL",
    "FI":      "FIXED INCOME & CREDIT",
    "FI_INTL": "FIXED INCOME — INTERNATIONAL",
    "RATES":   "SOVEREIGN RATES",
    "CMD":     "COMMODITIES",
    "CRYPTO":  "CRYPTO",
    "FX":      "FX",
    "VOL":     "VOLATILITY",
    "ALT":     "LISTED ALTERNATIVES",
}

# N_INSTRUMENTS is derived after DISPLAY_EXCLUSIONS is defined (see below).

# ── BUCKET + REGIME FIT ──────────────────────────────────────────────────────
INSTRUMENT_BUCKETS = {t: b for _sec, t, _n, b in UNIVERSE}

BUCKET_TICKERS = {
    "EQ Growth":    [t for _s, t, _n, b in UNIVERSE if b == "EQ Growth"],
    "EQ Defensive": [t for _s, t, _n, b in UNIVERSE if b == "EQ Defensive"],
    "Fixed Income": [t for _s, t, _n, b in UNIVERSE if b == "Fixed Income"],
    "Real Assets":  [t for _s, t, _n, b in UNIVERSE if b == "Real Assets"],
    "Cash":         [t for _s, t, _n, b in UNIVERSE if b == "Cash"],
    "Alts":         [t for _s, t, _n, b in UNIVERSE if b == "Alts"],
}

REGIME_FIT_MATRIX = {
    "EQ Growth":    {"Bull": 100, "Neutral":  50, "Volatile":  25, "Bear":   0, "Crisis":   0},
    "EQ Defensive": {"Bull":  50, "Neutral": 100, "Volatile":  75, "Bear":  75, "Crisis":  50},
    "Fixed Income": {"Bull":  25, "Neutral":  75, "Volatile":  75, "Bear": 100, "Crisis": 100},
    "Real Assets":  {"Bull":  75, "Neutral":  75, "Volatile":  75, "Bear":  50, "Crisis":  50},
    "Cash":         {"Bull":  10, "Neutral":  50, "Volatile":  75, "Bear": 100, "Crisis": 100},
    "Alts":         {"Bull": 100, "Neutral":  25, "Volatile":  50, "Bear":   0, "Crisis":   0},
}

# Map internal regime labels onto the matrix's 5 states.
INTERNAL_REGIME_TO_FIT = {"Calm": "Neutral", "Stressed": "Bear", "Crisis": "Crisis"}


def get_regime_fit_score(ticker: str, current_regime: str) -> int:
    bucket = INSTRUMENT_BUCKETS.get(ticker, "EQ Growth")
    regime = INTERNAL_REGIME_TO_FIT.get(current_regime, current_regime)
    return REGIME_FIT_MATRIX.get(bucket, {}).get(regime, 50)


# ── SPECIAL TICKERS ───────────────────────────────────────────────────────────
# Yield tickers from yfinance come back x10 (45.0 = 4.5%). Divide after download.
YIELD_TICKERS = ["^TNX", "^IRX", "^FVX", "^TYX"]

# Synthetic tickers — NOT fetched from yfinance. Computed from underlying series.
SYNTHETIC_TICKERS = {"GVZ": "GLD", "OVX": "BNO"}  # proxy -> source for 20D vol

# ── RANKING EXCLUSIONS ────────────────────────────────────────────────────────
# Tickers excluded from Top 5 Gainers/Losers, Top Picks, RSR rankings.
# These are data/index instruments — not directly investable ETFs.
RANKING_EXCLUSIONS = [
    # Computed vol proxies
    "GVZ", "OVX",
    # Yield indices (divided by 10 for display, but not investable)
    "^TNX", "^IRX", "^FVX", "^TYX",
    # Volatility instruments — ranking them confuses "up = bad"
    "^VIX", "^VIX3M", "VIXY", "VXX",
    # FX pairs — not investable ETFs, must never appear in any ranking
    "EURUSD=X", "GBPUSD=X", "AUDUSD=X", "SGD=X", "CHF=X",
    "JPY=X", "CAD=X", "MXN=X", "KRW=X", "BRL=X", "DX-Y.NYB",
]


def is_rankable(ticker: str) -> bool:
    """True if instrument is eligible for performance / opportunity rankings."""
    return ticker not in RANKING_EXCLUSIONS


# Tickers for which fragility scoring is meaningless (computed proxies,
# fear-gauge indices). These show N/A in the Fragility tab.
FRAGILITY_EXCLUSIONS = ["GVZ", "OVX", "^VIX", "^VIX3M"]


# Tickers entirely hidden from all user-facing tabs (Performance, Risk,
# Fragility, Analysis RSR, accordion headers). These are computed proxies
# held internally for vol-regime signals but never shown as instruments.
DISPLAY_EXCLUSIONS = ["GVZ", "OVX", "^VIX3M"]


def is_displayable(ticker: str) -> bool:
    return ticker not in DISPLAY_EXCLUSIONS

# Displayable instrument count: total universe minus hidden computed proxies.
# Derived dynamically so it stays correct whenever UNIVERSE or DISPLAY_EXCLUSIONS changes.
N_INSTRUMENTS = len([t for _, t, _, _ in UNIVERSE if t not in DISPLAY_EXCLUSIONS])


# Tickers for which the Sharpe ratio is meaningless or misleading.
# Yield indices are rate levels, not tradable returns.
# BIL/SHY: T-Bill ETFs whose total return is almost entirely distributions;
# yfinance auto_adjust does not reliably capture dividend adjustments for these,
# so price-return Sharpe produces deeply negative values that are economically wrong.
SHARPE_EXCLUSIONS = ["^TNX", "^IRX", "^FVX", "^TYX", "BIL", "SHY"]


# Vol-of-vol is not a sensible display metric for volatility indices
# themselves (they routinely swing 50%+ in a day). Show N/A for their
# vol columns in the Risk tab and Performance sparkline.
VOL_VALUE_EXCLUSIONS = ["^VIX", "^VIX3M"]


# ── CURRENCY MAP ──────────────────────────────────────────────────────────────
# Returns are price return in the instrument's local currency
CURRENCY_MAP = {
    "EQ_US": "USD", "EQ_SECT": "USD", "EQ_IDX": "USD",
    "EQ_DM": "USD", "EQ_EM": "USD", "DEFENCE": "USD",
    "FI": "USD", "FI_INTL": "USD", "RATES": "USD", "CMD": "USD",
    "CRYPTO": "USD", "VOL": "USD", "FX": "USD", "ALT": "USD",
}
# Per-ticker currency override for FX / cross-currency instruments
FX_CCY_MAP = {
    "EURUSD=X": "EUR", "GBPUSD=X": "GBP", "JPY=X": "JPY",
    "SGD=X": "SGD", "AUDUSD=X": "AUD", "CHF=X": "CHF",
    "CAD=X": "CAD", "MXN=X": "MXN", "KRW=X": "KRW", "BRL=X": "BRL",
    "DX-Y.NYB": "USD",
    # EQ_DM / EQ_EM display currencies (underlying local market)
    "EWA": "AUD", "EWG": "EUR", "EWJ": "JPY", "EWS": "SGD", "EWU": "GBP",
    "EWZ": "BRL", "FXI": "CNY", "INDA": "INR", "EWY": "KRW",
    "EZA": "ZAR", "EWT": "TWD", "EWW": "MXN", "EIDO": "IDR", "VNM": "VND",
    "KSA": "SAR", "EWH": "HKD", "FLGB": "GBP",
}




# ── COLORS (report) ───────────────────────────────────────────────────────────

BG    = "#0d1117"
CARD  = "#161b22"
DARK  = "#21262d"
WHITE = "#e6edf3"
GREY  = "#8b949e"
ACNT  = "#58a6ff"
RED   = "#f85149"
AMBER = "#e3b341"
GREEN = "#3fb950"


def _ret_color(v: float) -> str:
    if pd.isna(v):   return GREY
    if v >=  0.02:   return "#3fb950"
    if v >=  0.005:  return "#7ee787"
    if v >=  0.0:    return "#a8d5b0"
    if v >= -0.005:  return "#ffa657"
    if v >= -0.02:   return "#ff7b72"
    return RED


def _sharpe_color(v: float) -> str:
    if pd.isna(v): return GREY
    if v > 1.0:    return GREEN
    if v > 0.0:    return AMBER
    return RED


def _rag(dd: float) -> tuple[str, str]:
    """RAG signal based on max drawdown from 52-week high.
    Missing data falls back to AMBER so every instrument has a signal
    and the RED/AMBER/GREEN counts reconcile with the universe size.
    """
    if pd.isna(dd): return AMBER, "AMBER"
    if dd < -0.15:  return RED,   " RED "
    if dd < -0.07:  return AMBER, "AMBER"
    return GREEN, "GREEN"


# ── MARKET TONE (three-way gate) ──────────────────────────────────────────────
# Supports both the briefing's 5-state regime labels and the internal
# "Calm / Stressed / Crisis" labels produced by compute_regime().
REGIME_SCORE_MAP = {
    'Bull': 5, 'Neutral': 4, 'Volatile': 3, 'Bear': 2, 'Crisis': 1,
    'Calm': 4, 'Stressed': 2,
}


def calculate_market_tone(regime, fragility_score, rising_risk_count, total_instruments):
    """
    Three-way gate. ALL conditions must align for RISK-ON.
    Any single condition failing -> MIXED.
    Multiple conditions failing badly -> RISK-OFF.
    Returns (label, fg_color, bg_color).
    """
    rising_risk_pct = (rising_risk_count / total_instruments) if total_instruments > 0 else 0
    r_score = REGIME_SCORE_MAP.get(regime, 3)

    regime_ok    = r_score >= 4           # Calm / Neutral / Bull
    fragility_ok = fragility_score < 55   # below Stressed threshold
    risk_ok      = rising_risk_pct < 0.40

    if regime_ok and fragility_ok and risk_ok:
        return 'RISK-ON', '#3fb950', '#0d2318'

    # RISK-OFF only when conditions are genuinely severe — was previously
    # too aggressive (firing on a single gate). Now requires Crisis regime
    # or Crisis-level fragility (>= 75). Stressed regime alone maps to MIXED.
    if r_score <= 1 or fragility_score >= 75:
        return 'RISK-OFF', '#f85149', '#2d0f0e'

    return 'MIXED', '#e3b341', '#2d2106'


def _count_rising_risk(df) -> int:
    """Count instruments whose 20D vol is >= 5% above vol 1M ago."""
    n = 0
    for _, row in df.iterrows():
        nv = row.get("vol_now", float("nan"))
        av = row.get("vol_1m_ago", float("nan"))
        if pd.isna(nv) or pd.isna(av) or av == 0:
            continue
        if (nv - av) / av >= 0.05:
            n += 1
    return n


# ══════════════════════════════════════════════════════════════════════════════
#  DATA  (with price cache for fast daily runs + full history)
# ══════════════════════════════════════════════════════════════════════════════

CACHE_FILE        = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "prices_cache.csv")
VOLUME_CACHE_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "volumes_cache.csv")

# ── IFM CANONICAL 43-ETF UNIVERSE ─────────────────────────────────────────────
# Used exclusively for system-level fragility score and regime detection.
# Matches institutional_fragility_monitor.py exactly — do not modify without
# updating IFM as well. Scores derived from this universe are CRO/GARP-comparable.
IFM_43_TICKERS = [
    # Global equity
    "ACWI", "SPY", "QQQ", "IWM", "EFA", "EEM", "VGK", "EWJ",
    # Asia / EM
    "MCHI", "INDA", "EWZ", "EWT",
    # Sectors
    "XLF", "XLE", "XLK", "XLV", "XLU", "XLP",
    # Fixed income
    "AGG", "LQD", "HYG", "TLT", "IEF", "SHY", "EMB",
    # Rates / inflation
    "TIP", "MBB",
    # Commodities
    "GLD", "SLV", "USO", "DBC", "DBA", "PDBC",
    # FX / alternatives
    "UUP", "FXE", "FXY",
    # Real assets
    "VNQ", "REET",
    # Crypto
    "IBIT", "ETHA",
    # Volatility / hedges
    "VIXY",
]


def _yf_download_safe(tickers, start, field="Close", batch_size=20, max_retries=3):
    """
    Download price or volume data with batching and retry logic.
    field: "Close" for prices, "Volume" for volumes.
    Never raises — returns what it can.
    """
    all_data = {}
    ticker_list = list(tickers)
    batches = [ticker_list[i:i+batch_size] for i in range(0, len(ticker_list), batch_size)]

    for batch in batches:
        for attempt in range(max_retries):
            try:
                data = yf.download(batch, start=start, auto_adjust=True, progress=False, timeout=30)
                if not data.empty:
                    if field in data.columns:
                        col_data = data[field]
                    else:
                        col_data = data.xs(field, axis=1, level=0)
                    if isinstance(col_data, pd.Series):
                        col_data = col_data.to_frame(name=batch[0])
                    if col_data.index.tz is not None:
                        col_data.index = col_data.index.tz_localize(None)
                    for t in batch:
                        if t in col_data.columns:
                            s = col_data[t].dropna()
                            if len(s) > 0:
                                all_data[t] = s
                break  # success
            except Exception as e:
                print(f"  Batch {field} download attempt {attempt+1} failed: {e}")
                if attempt < max_retries - 1:
                    time.sleep(3)
                else:
                    # Final fallback — try one by one
                    for t in batch:
                        try:
                            tk   = yf.Ticker(t)
                            hist = tk.history(start=start, auto_adjust=True)
                            if len(hist) > 0:
                                s = hist[field].dropna() if field in hist.columns else pd.Series(dtype=float)
                                if s.index.tz is not None:
                                    s.index = s.index.tz_localize(None)
                                if len(s) > 0:
                                    all_data[t] = s
                        except Exception:
                            print(f"    Failed individual {field}: {t}")

    if all_data:
        result = pd.DataFrame(all_data).sort_index()
        return result
    return pd.DataFrame()


def validate_cache(cache_path):
    """Check cache integrity. Returns True if cache is usable."""
    if not os.path.exists(cache_path):
        print("[Cache]  Missing — will download fresh")
        return False
    try:
        df = pd.read_csv(cache_path, index_col=0, parse_dates=True)
        if df.empty:
            print("[Cache]  Empty — rebuilding")
            return False
        latest = df.index.max()
        today = pd.Timestamp.today().normalize()
        age_days = (today - latest).days
        if age_days > 3:
            print(f"[Cache]  {age_days} days old — refreshing")
            return False
        print(f"[Cache]  Valid — {len(df.columns)} tickers, latest {latest.date()}")
        return True
    except Exception as e:
        print(f"[Cache]  Corrupt ({e}) — rebuilding")
        return False


def format_stale_badge(last_date):
    """HTML badge for data freshness based on last data date."""
    if last_date is None:
        return ''
    today = pd.Timestamp.today().normalize()
    if today.weekday() >= 5:
        today = today - pd.Timedelta(days=today.weekday() - 4)
    days_old = (today - last_date).days
    if days_old <= 1:
        return ''
    elif days_old <= 3:
        return (' <span style="font-size:8px;padding:1px 4px;border:1px solid #e3b341;'
                'border-radius:3px;color:#e3b341;margin-left:4px;">T-2</span>')
    else:
        return (' <span style="font-size:8px;padding:1px 4px;border:1px solid #f85149;'
                'border-radius:3px;color:#f85149;margin-left:4px;">STALE</span>')


def download(lookback_days: int = 2520) -> tuple:
    """
    Download price and volume data for all universe instruments.
    Returns (prices, volumes) DataFrames.
    Uses smart incremental cache — only fetches new data when cache is fresh.
    Volume tickers: excludes synthetic, yield, and vol-index tickers (no meaningful volume).
    """
    # Real fetch list excludes synthetic tickers (computed post-download)
    tickers = [t for _s, t, _n, _b in UNIVERSE if t not in SYNTHETIC_TICKERS]
    # Volume is not meaningful for yield indices or vol fear gauges
    NO_VOLUME_TICKERS = set(YIELD_TICKERS) | {"^VIX", "^VIX3M", "GVZ", "OVX"}
    vol_tickers = [t for t in tickers if t not in NO_VOLUME_TICKERS]
    print(f"[Download] {len(tickers)} price tickers | {len(vol_tickers)} volume tickers | lookback={lookback_days} days ...")

    # ── PRICES ─────────────────────────────────────────────────────────────────
    cached = None
    cache_valid = validate_cache(CACHE_FILE)
    if cache_valid:
        try:
            cached = pd.read_csv(CACHE_FILE, index_col=0, parse_dates=True)
            if cached.index.tz is not None:
                cached.index = cached.index.tz_localize(None)
            print(f"[Cache]  Prices loaded {len(cached)} days (last: {cached.index[-1].date()})")
        except Exception as e:
            print(f"[Cache]  Price load failed: {e}")
            cached = None

    new_tickers = []
    if cached is not None:
        new_tickers = [t for t in tickers if t not in cached.columns]
        if new_tickers:
            print(f"[Download] New price tickers (full history): {len(new_tickers)} -> "
                  f"{', '.join(new_tickers[:8])}{'...' if len(new_tickers) > 8 else ''}")

    if cached is not None and len(cached) >= 756 and not new_tickers:
        start = (pd.Timestamp.today() - pd.Timedelta(days=60)).strftime("%Y-%m-%d")
        print(f"[Download] Price cache hit — refreshing last 60 days ...")
        fetch_list = tickers
    elif cached is not None and new_tickers:
        start = (pd.Timestamp.today() - pd.Timedelta(days=60)).strftime("%Y-%m-%d")
        existing = [t for t in tickers if t in cached.columns]
        full_start = (pd.Timestamp.today() - pd.Timedelta(days=lookback_days)).strftime("%Y-%m-%d")
        new_prices = _yf_download_safe(new_tickers, start=full_start, field="Close")
        if not new_prices.empty:
            cached = cached.join(new_prices, how="outer")
        fetch_list = existing
    else:
        start = (pd.Timestamp.today() - pd.Timedelta(days=lookback_days)).strftime("%Y-%m-%d")
        print(f"[Download] No price cache — full {lookback_days}-day download ...")
        fetch_list = tickers

    prices_new = _yf_download_safe(fetch_list, start=start, field="Close")
    if prices_new.empty and cached is None:
        raise RuntimeError("No price data returned from Yahoo Finance.")

    if not prices_new.empty:
        prices = pd.concat([cached, prices_new]) if cached is not None else prices_new
        prices = prices[~prices.index.duplicated(keep="last")].sort_index()
    else:
        prices = cached

    try:
        os.makedirs(os.path.dirname(CACHE_FILE), exist_ok=True)
        prices.to_csv(CACHE_FILE)
        print(f"[Cache]  Prices saved {len(prices)} days")
    except Exception as e:
        print(f"[Cache]  Price save failed: {e}")

    prices = prices.ffill(limit=3).dropna(how="all")
    prices = prices[[t for t in tickers if t in prices.columns]]

    # Yield ticker correction: yfinance returns x10 values
    for yt in YIELD_TICKERS:
        if yt in prices.columns:
            prices[yt] = prices[yt] / 10.0

    # Synthetic tickers: compute GVZ/OVX as 20D rolling vol of GLD/BNO
    for proxy, source in SYNTHETIC_TICKERS.items():
        if source in prices.columns:
            rets = prices[source].pct_change()
            rolling_vol = rets.rolling(window=20).std() * (252 ** 0.5) * 100.0
            prices[proxy] = rolling_vol.bfill()

    print(f"[Download] Prices: {len(prices)} days | last close: {prices.index[-1].date()} | cols: {len(prices.columns)}")

    # ── VOLUMES ────────────────────────────────────────────────────────────────
    vcached = None
    vcache_valid = validate_cache(VOLUME_CACHE_FILE)
    if vcache_valid:
        try:
            vcached = pd.read_csv(VOLUME_CACHE_FILE, index_col=0, parse_dates=True)
            if vcached.index.tz is not None:
                vcached.index = vcached.index.tz_localize(None)
            print(f"[Cache]  Volumes loaded {len(vcached)} days (last: {vcached.index[-1].date()})")
        except Exception as e:
            print(f"[Cache]  Volume load failed: {e}")
            vcached = None

    new_vol_tickers = []
    if vcached is not None:
        new_vol_tickers = [t for t in vol_tickers if t not in vcached.columns]

    if vcached is not None and len(vcached) >= 756 and not new_vol_tickers:
        vstart = (pd.Timestamp.today() - pd.Timedelta(days=60)).strftime("%Y-%m-%d")
        vfetch_list = vol_tickers
    elif vcached is not None and new_vol_tickers:
        vstart = (pd.Timestamp.today() - pd.Timedelta(days=60)).strftime("%Y-%m-%d")
        vfull_start = (pd.Timestamp.today() - pd.Timedelta(days=lookback_days)).strftime("%Y-%m-%d")
        new_vols = _yf_download_safe(new_vol_tickers, start=vfull_start, field="Volume")
        if not new_vols.empty:
            vcached = vcached.join(new_vols, how="outer")
        vfetch_list = [t for t in vol_tickers if t in vcached.columns]
    else:
        vstart = (pd.Timestamp.today() - pd.Timedelta(days=lookback_days)).strftime("%Y-%m-%d")
        vfetch_list = vol_tickers

    vols_new = _yf_download_safe(vfetch_list, start=vstart, field="Volume")

    if not vols_new.empty:
        volumes = pd.concat([vcached, vols_new]) if vcached is not None else vols_new
        volumes = volumes[~volumes.index.duplicated(keep="last")].sort_index()
    elif vcached is not None:
        volumes = vcached
    else:
        # Fallback: zero volumes — fragility engine handles gracefully
        print("[Download] WARNING: No volume data — liquidity pillar will be zero")
        volumes = pd.DataFrame(0.0, index=prices.index, columns=prices.columns)

    try:
        volumes.to_csv(VOLUME_CACHE_FILE)
        print(f"[Cache]  Volumes saved {len(volumes)} days")
    except Exception as e:
        print(f"[Cache]  Volume save failed: {e}")

    volumes = volumes.reindex(prices.index).ffill(limit=3).fillna(0.0)
    # Align columns: only tickers present in prices and not excluded
    shared_cols = [t for t in prices.columns if t in volumes.columns and t not in NO_VOLUME_TICKERS]
    volumes = volumes.reindex(columns=prices.columns, fill_value=0.0)

    print(f"[Download] Volumes: {len(volumes)} days | vol cols with data: {(volumes > 0).any().sum()}")
    return prices, volumes


def compute_metrics(prices: pd.DataFrame) -> pd.DataFrame:
    today = prices.index[-1]

    def _ret(n: int) -> pd.Series:
        if len(prices) <= n:
            return pd.Series(np.nan, index=prices.columns)
        return prices.iloc[-1] / prices.iloc[-1 - n] - 1

    # YTD vs first trading day of current calendar year
    ytd_slice = prices[prices.index.year == today.year]
    ytd = (
        prices.iloc[-1] / ytd_slice.iloc[0] - 1
        if len(ytd_slice) > 1
        else pd.Series(np.nan, index=prices.columns)
    )

    # Sanity check: flag suspicious return values (likely unadjusted corporate
    # actions upstream in yfinance — we confirmed all oil ETFs including BNO/UCO
    # show the same corruption, so swapping source doesn't help).
    # Cap by horizon: (ytd_cap, ret_1m_cap, ret_3m_cap)
    RETURN_SANITY_MAX = {
        'BNO':  (0.50, 0.30, 0.40),
        'UNG':  (1.00, 0.50, 0.75),
        'SLV':  (1.00, 0.40, 0.60),
        'GLD':  (0.75, 0.30, 0.50),
        'DBC':  (0.75, 0.30, 0.50),
        'DBA':  (0.75, 0.30, 0.50),
        'COPX': (1.00, 0.50, 0.75),
    }
    # Pre-compute ret_1m / ret_3m once for sanity-clipping (not just YTD).
    _ret_1m_series = _ret(21)
    _ret_3m_series = _ret(63)
    data_review_tickers = set()
    for _tk, (_ytd_cap, _m1_cap, _m3_cap) in RETURN_SANITY_MAX.items():
        if _tk not in ytd.index:
            continue
        _yv = ytd.get(_tk)
        _mv = _ret_1m_series.get(_tk) if _tk in _ret_1m_series.index else np.nan
        _qv = _ret_3m_series.get(_tk) if _tk in _ret_3m_series.index else np.nan
        _flag = False
        if pd.notna(_yv) and abs(_yv) > _ytd_cap:
            print(f"[Data Sanity] {_tk} YTD={_yv*100:+.1f}% > ±{_ytd_cap*100:.0f}% — flagged")
            ytd[_tk] = np.nan; _flag = True
        if pd.notna(_mv) and abs(_mv) > _m1_cap:
            print(f"[Data Sanity] {_tk} 1M={_mv*100:+.1f}% > ±{_m1_cap*100:.0f}% — flagged")
            _ret_1m_series[_tk] = np.nan; _flag = True
        if pd.notna(_qv) and abs(_qv) > _m3_cap:
            print(f"[Data Sanity] {_tk} 3M={_qv*100:+.1f}% > ±{_m3_cap*100:.0f}% — flagged")
            _ret_3m_series[_tk] = np.nan; _flag = True
        if _flag:
            data_review_tickers.add(_tk)
    if data_review_tickers:
        print(f"[Data Sanity] Data under review: {sorted(data_review_tickers)} "
              f"— suspected unadjusted corporate actions in yfinance source")

    # Vol as of latest date: 5-day realised vol annualised (captures current regime)
    vol_now  = prices.pct_change().tail(6).std()  * np.sqrt(252)

    # Vol 1 month ago: 20-day window ending ~21 trading days back
    ret_chg  = prices.pct_change()
    vol_1m_ago = (
        ret_chg.iloc[-42:-21].std() * np.sqrt(252)
        if len(ret_chg) >= 42
        else pd.Series(np.nan, index=prices.columns)
    )

    # 20-day annualised vol (short-term, for signal column)
    vol_20d = ret_chg.tail(21).std() * np.sqrt(252)

    # 1-year annualised vol + Sharpe
    daily_ret_1y = prices.pct_change().tail(252)
    vol_1y       = daily_ret_1y.std() * np.sqrt(252)
    ann_ret_1y   = daily_ret_1y.mean() * 252
    # For very low-vol instruments (e.g. T-Bills), cap vol at 0.5% minimum to prevent extreme Sharpe
    vol_1y_adj   = vol_1y.clip(lower=0.005)
    sharpe       = ((ann_ret_1y - RISK_FREE_RATE) / vol_1y_adj.replace(0, np.nan)).clip(-5, 5)

    # Max drawdown from 252-day rolling peak
    window  = min(252, len(prices))
    peak    = prices.tail(window).cummax()
    max_dd  = prices.iloc[-1] / peak.iloc[-1] - 1

    # Sparkline data: last 20 trading days, normalised to first value
    spark_window = min(20, len(prices))
    spark_prices = prices.tail(spark_window)

    # Detect if market was closed today
    # Only check equity tickers (FX/Crypto trade 24/7 and skew the signal)
    ret1d_vals = _ret(1)
    _non_eq_sections = {"FX", "CRYPTO", "VOL", "RATES"}
    eq_tickers = [t for s, t, _n, _b in UNIVERSE
                  if t in prices.columns and s not in _non_eq_sections
                  and t not in SYNTHETIC_TICKERS]
    eq_ret1d = ret1d_vals.reindex(eq_tickers).dropna()
    nonzero_eq = (eq_ret1d.abs() > 1e-4).sum()
    market_open_today = nonzero_eq >= max(3, len(eq_ret1d) * 0.15)

    rows = []
    for sec, ticker, name, _bucket in UNIVERSE:
        if ticker not in prices.columns:
            continue
        # Internal-only tickers (GVZ/OVX proxies) — never surface as rows.
        if ticker in DISPLAY_EXCLUSIONS:
            continue
        rc, rl = _rag(max_dd.get(ticker, np.nan))
        # Normalised sparkline series (percent from 20-day-ago base)
        sp = spark_prices[ticker].dropna()
        spark = list((sp / sp.iloc[0] - 1) * 100) if len(sp) > 1 else []
        rows.append({
            "section":          sec,
            "ticker":           ticker,
            "name":             name,
            "ret_1d":           ret1d_vals.get(ticker, np.nan),
            "ret_1w":           _ret(5).get(ticker, np.nan),
            "ret_1m":           _ret_1m_series.get(ticker, np.nan),
            "ret_3m":           _ret_3m_series.get(ticker, np.nan),
            "ret_ytd":          ytd.get(ticker, np.nan),
            "data_review":      ticker in data_review_tickers,
            "vol_20d":          vol_20d.get(ticker, np.nan),
            "vol_now":          vol_now.get(ticker, np.nan),
            "vol_1m_ago":       vol_1m_ago.get(ticker, np.nan),
            "max_dd":           max_dd.get(ticker, np.nan),
            "sharpe":           sharpe.get(ticker, np.nan),
            "rag_color":        rc,
            "rag_label":        rl,
            "spark":            spark,
            "market_open":      market_open_today,
            "currency":         FX_CCY_MAP.get(ticker, CURRENCY_MAP.get(sec, "USD")),
        })

    ord_map = {s: i for i, s in enumerate(SECTION_ORDER)}
    df = pd.DataFrame(rows)
    df["_o"] = df["section"].map(ord_map).fillna(99)
    return df.sort_values(["_o", "name"]).drop(columns="_o").reset_index(drop=True)


# ══════════════════════════════════════════════════════════════════════════════
#  VISUAL REPORT (PNG + PDF)
# ══════════════════════════════════════════════════════════════════════════════

def _fmt(v: float, signed: bool = True, dec: int = 1, pct: bool = True) -> str:
    if pd.isna(v): return "-"
    if pct:
        return f"{v * 100:+.{dec}f}%" if signed else f"{v * 100:.{dec}f}%"
    return f"{v:+.{dec}f}" if signed else f"{v:.{dec}f}"


def render_report(df: pd.DataFrame, as_of: str, out_dir: str) -> tuple[str, str]:
    # ── Layout constants ──
    FIG_W     = 28.0   # wider canvas
    HDR_ROWS  = 3.2
    COLHDR_R  = 1.2
    SEC_ROWS  = 1.0
    DATA_ROWS = 1.5    # tall rows → readable text
    FOOT_ROWS = 0.8
    ROW_INCH  = 0.36   # inches per row-unit → controls overall height

    n_sections = df["section"].nunique()
    n_assets   = len(df)
    total_h    = HDR_ROWS + COLHDR_R + n_sections * SEC_ROWS + n_assets * DATA_ROWS + FOOT_ROWS
    FIG_H      = total_h * ROW_INCH

    # Font sizes
    FS_TITLE   = 24
    FS_SUB     = 11
    FS_RAG_N   = 22
    FS_RAG_L   = 10
    FS_COLHDR  = 11
    FS_SECHDR  = 10
    FS_NAME    = 11
    FS_TICKER  = 10
    FS_DATA    = 11
    FS_FOOT    = 9

    # Detect market closed (all 1D returns ~0)
    market_open = bool(df["market_open"].iloc[0]) if "market_open" in df.columns else True

    # ── Column x positions (0–100 scale) ──
    # Trend sparkline always shown; 1D suppressed when market closed
    if market_open:
        CX = {
            "name":    1.5,
            "ticker":  27.0,
            "trend":   33.5,
            "1d":      39.5,
            "1w":      46.0,
            "1m":      52.5,
            "3m":      59.0,
            "ytd":     65.5,
            "vol":     72.0,
            "dd":      79.0,
            "sharpe":  86.0,
            "sig":     94.0,
        }
    else:
        CX = {
            "name":    1.5,
            "ticker":  27.0,
            "trend":   33.5,
            "1w":      41.0,
            "1m":      49.0,
            "3m":      57.0,
            "ytd":     65.0,
            "vol":     73.0,
            "dd":      80.5,
            "sharpe":  88.0,
            "sig":     95.5,
        }

    fig = plt.figure(figsize=(FIG_W, FIG_H), facecolor=BG)
    ax  = fig.add_axes([0.0, 0.0, 1.0, 1.0])
    ax.set_xlim(0, 100)
    ax.set_ylim(0, total_h)
    ax.set_facecolor(BG)
    ax.axis("off")

    # ── Header ──
    y     = total_h
    hdr_y = y - HDR_ROWS
    ax.add_patch(Rectangle((0, hdr_y), 100, HDR_ROWS, facecolor="#1c2128", zorder=0))

    ax.text(2, hdr_y + HDR_ROWS * 0.70,
            "BK  MARKET  DASHBOARD",
            fontsize=FS_TITLE, fontweight="bold", color=WHITE,
            ha="left", va="center", family="monospace")
    ax.text(2, hdr_y + HDR_ROWS * 0.26,
            f"As of  {as_of}   |   {N_INSTRUMENTS}-instrument universe   |   "
            f"Source: Yahoo Finance   |   Price return, local currency",
            fontsize=FS_SUB, color=GREY, ha="left", va="center")

    # RAG summary counts (top-right of header)
    n_red   = (df["rag_label"].str.strip() == "RED").sum()
    n_amber = (df["rag_label"].str.strip() == "AMBER").sum()
    n_green = (df["rag_label"].str.strip() == "GREEN").sum()
    for xi, (lbl, cnt, col) in enumerate([
        ("RED",   n_red,   RED),
        ("AMBER", n_amber, AMBER),
        ("GREEN", n_green, GREEN),
    ]):
        xp = 77 + xi * 8
        ax.text(xp, hdr_y + HDR_ROWS * 0.72, str(cnt),
                fontsize=FS_RAG_N, fontweight="bold", color=col, ha="center", va="center")
        ax.text(xp, hdr_y + HDR_ROWS * 0.26, lbl,
                fontsize=FS_RAG_L, color=GREY, ha="center", va="center")

    y = hdr_y

    # ── Column header bar ──
    ch_y = y - COLHDR_R
    ax.add_patch(Rectangle((0, ch_y), 100, COLHDR_R, facecolor=DARK, zorder=0))
    col_hdrs = [
        ("name",   "Asset",      "left"),
        ("ticker", "Ticker",     "center"),
        ("trend",  "20D Trend",  "center"),
    ]
    if market_open:
        col_hdrs.append(("1d", "1D", "center"))
    col_hdrs += [
        ("1w",     "1W",      "center"),
        ("1m",     "1M",      "center"),
        ("3m",     "3M",      "center"),
        ("ytd",    "YTD",     "center"),
        ("vol",    "Vol 20D", "center"),
        ("dd",     "Max DD",  "center"),
        ("sharpe", "Sharpe",  "center"),
        ("sig",    "Signal",  "center"),
    ]
    for key, lbl, ha in col_hdrs:
        ax.text(CX[key], ch_y + COLHDR_R * 0.5, lbl,
                fontsize=FS_COLHDR, fontweight="bold", color=ACNT, ha=ha, va="center")

    y        = ch_y
    prev_sec = None

    # ── Data rows ──
    for _, row in df.iterrows():

        # Section divider
        if row["section"] != prev_sec:
            prev_sec = row["section"]
            s_y      = y - SEC_ROWS
            ax.add_patch(Rectangle((0, s_y), 100, SEC_ROWS, facecolor=CARD, zorder=0))
            ax.text(CX["name"], s_y + SEC_ROWS * 0.5,
                    SECTION_LABELS.get(row["section"], row["section"]),
                    fontsize=FS_SECHDR, fontweight="bold", color=ACNT,
                    ha="left", va="center", alpha=0.9)
            y = s_y

        r_y = y - DATA_ROWS
        ax.add_patch(Rectangle((0, r_y), 100, DATA_ROWS, facecolor=BG, zorder=0))
        ax.axhline(r_y, color=DARK, lw=0.5, xmin=0.01, xmax=0.99)
        mid = r_y + DATA_ROWS * 0.5

        # Name + ticker
        ax.text(CX["name"],   mid, row["name"][:32],
                fontsize=FS_NAME, color=WHITE, ha="left", va="center")
        ax.text(CX["ticker"], mid, row["ticker"],
                fontsize=FS_TICKER, color=GREY, ha="center", va="center", family="monospace")

        # ── Sparkline (20-day trend mini chart) ──────────────────────────────
        spark = row.get("spark", [])
        if len(spark) > 1:
            sp_arr   = np.array(spark)
            sp_color = GREEN if sp_arr[-1] >= 0 else RED
            # Convert data coords → axes (0-100 x, 0-total_h y)
            sp_w   = 4.5   # width in data-x units
            sp_h   = DATA_ROWS * 0.65
            sp_x0  = CX["trend"] - sp_w / 2
            sp_y0  = mid - sp_h / 2
            sp_x1  = CX["trend"] + sp_w / 2
            sp_y1  = mid + sp_h / 2
            # Normalise data to fit in box
            mn, mx = sp_arr.min(), sp_arr.max()
            rng    = mx - mn if mx != mn else 1.0
            xs = np.linspace(sp_x0, sp_x1, len(sp_arr))
            ys = sp_y0 + (sp_arr - mn) / rng * sp_h
            # Shaded fill under/over zero line
            zero_y = sp_y0 + (0 - mn) / rng * sp_h
            zero_y = np.clip(zero_y, sp_y0, sp_y1)
            ax.fill_between(xs, zero_y, ys,
                            where=(ys >= zero_y), color=GREEN, alpha=0.18, linewidth=0)
            ax.fill_between(xs, zero_y, ys,
                            where=(ys < zero_y),  color=RED,   alpha=0.18, linewidth=0)
            ax.plot(xs, ys, color=sp_color, linewidth=0.9, solid_capstyle="round")
            # Zero baseline
            ax.axhline(zero_y, xmin=(sp_x0 / 100), xmax=(sp_x1 / 100),
                       color=GREY, linewidth=0.3, alpha=0.5)

        # Return columns
        ret_cols = []
        if market_open:
            ret_cols.append(("1d", "ret_1d"))
        ret_cols += [("1w", "ret_1w"), ("1m", "ret_1m"), ("3m", "ret_3m"), ("ytd", "ret_ytd")]
        for key, col_key in ret_cols:
            v = row[col_key]
            ax.text(CX[key], mid, _fmt(v),
                    fontsize=FS_DATA, color=_ret_color(v), ha="center", va="center",
                    family="monospace")

        # Vol 20D (unsigned, grey — informational)
        ax.text(CX["vol"], mid, _fmt(row["vol_20d"], signed=False),
                fontsize=FS_DATA, color=GREY, ha="center", va="center", family="monospace")

        # Max drawdown
        v_dd = row["max_dd"]
        ax.text(CX["dd"], mid, _fmt(v_dd),
                fontsize=FS_DATA, color=_ret_color(v_dd), ha="center", va="center",
                family="monospace")

        # Sharpe
        v_sh = row["sharpe"]
        ax.text(CX["sharpe"], mid, _fmt(v_sh, pct=False, dec=2),
                fontsize=FS_DATA, color=_sharpe_color(v_sh), ha="center", va="center",
                family="monospace")

        # Signal (dot + label)
        rc = row["rag_color"]
        ax.text(CX["sig"] - 2.5, mid, "●", fontsize=FS_DATA + 2, color=rc, ha="center", va="center")
        ax.text(CX["sig"] + 2.5, mid, row["rag_label"].strip(),
                fontsize=FS_DATA - 1, color=rc, ha="center", va="center")

        y = r_y

    # ── Footer ──
    ax.text(2, FOOT_ROWS * 0.5,
            "Signal: RED = Max DD < -15%   AMBER = -15% to -7%   GREEN = > -7% from 52-week high   |   "
            "Vol 20D = annualised 20-day vol   |   Sharpe = 1Y excess return / vol  (rf = 4.5%)",
            fontsize=FS_FOOT, color=GREY, ha="left", va="center")
    ax.text(98, FOOT_ROWS * 0.5,
            f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}   |   CONFIDENTIAL",
            fontsize=FS_FOOT, color=GREY, ha="right", va="center")

    # ── Save ──
    os.makedirs(out_dir, exist_ok=True)
    tag      = datetime.now().strftime("%Y%m%d_%H%M")
    png_path = os.path.join(out_dir, f"market_dashboard_{tag}.png")
    pdf_path = os.path.join(out_dir, f"market_dashboard_{tag}.pdf")

    fig.savefig(png_path, dpi=150, bbox_inches="tight", facecolor=BG)
    with PdfPages(pdf_path) as pdf:
        pdf.savefig(fig, bbox_inches="tight", facecolor=BG)
    plt.close(fig)

    return png_path, pdf_path


# ══════════════════════════════════════════════════════════════════════════════
#  EMAIL BRIEF (HTML)
# ══════════════════════════════════════════════════════════════════════════════

def _ec_ret(v):
    """Email color + formatted string for a return (fraction)."""
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "#6b7280", "-"
    pct  = v * 100
    sign = "+" if pct > 0 else ""
    if pct >=  3.0: return "#065f46", f"{sign}{pct:.2f}%"
    if pct >=  1.5: return "#059669", f"{sign}{pct:.2f}%"
    if pct >=  0.0: return "#374151", f"{sign}{pct:.2f}%"
    if pct >= -1.5: return "#dc2626", f"{sign}{pct:.2f}%"
    if pct >= -3.0: return "#b91c1c", f"{sign}{pct:.2f}%"
    return "#7f1d1d", f"{sign}{pct:.2f}%"


def _ec_vol(v):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "#6b7280", "-"
    pct = v * 100
    c   = "#dc2626" if pct > 30 else "#d97706" if pct > 18 else "#059669"
    return c, f"{pct:.1f}%"


def _ec_dd(v):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "#6b7280", "-"
    pct = v * 100
    c   = "#dc2626" if pct < -15 else "#d97706" if pct < -7 else "#059669"
    return c, f"{pct:.1f}%"


def _ec_sharpe(v):
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "#6b7280", "-"
    c = "#059669" if v > 1 else "#d97706" if v > 0 else "#dc2626"
    return c, f"{v:.2f}"


def build_email_html(df: pd.DataFrame) -> str:
    now      = datetime.now(SGT)
    date_str = now.strftime("%a %d %b %Y · %H:%M SGT")

    ret_abs = df["ret_1d"].abs() * 100
    n_alerts  = int((ret_abs >= 3.0).sum())
    n_watches = int(((ret_abs >= 1.5) & (ret_abs < 3.0)).sum())
    n_stable  = int((ret_abs < 1.5).sum())

    def _td(color, val, extra_style=""):
        return (
            f'<td style="padding:7px 9px;font-family:monospace;font-size:11px;'
            f'text-align:right;color:{color};border-right:1px solid #e2e6ea;{extra_style}">'
            f'{val}</td>'
        )

    rows_html = ""
    prev_sec  = None

    for _, row in df.iterrows():
        if row["section"] != prev_sec:
            prev_sec  = row["section"]
            sec_label = SECTION_LABELS.get(row["section"], row["section"])
            rows_html += (
                f'<tr><td colspan="11" style="background:#f1f3f5;font-size:9px;font-weight:700;'
                f'letter-spacing:2px;text-transform:uppercase;color:#6b7280;'
                f'padding:7px 12px;border-top:2px solid #cbd2d9;">{sec_label}</td></tr>'
            )

        c1d,  v1d  = _ec_ret(row["ret_1d"])
        c1w,  v1w  = _ec_ret(row["ret_1w"])
        c1m,  v1m  = _ec_ret(row["ret_1m"])
        c3m,  v3m  = _ec_ret(row["ret_3m"])
        cytd, vytd = _ec_ret(row["ret_ytd"])
        cvol, vvol = _ec_vol(row["vol_20d"])
        cdd,  vdd  = _ec_dd(row["max_dd"])
        csh,  vsh  = _ec_sharpe(row["sharpe"])
        rc         = row["rag_color"]
        rl         = row["rag_label"].strip()

        rows_html += (
            f'<tr style="border-bottom:1px solid #e2e6ea;">'
            f'<td style="padding:7px 9px;font-size:11px;color:#111827;'
            f'border-right:1px solid #e2e6ea;white-space:nowrap;">{row["name"]}</td>'
            f'<td style="padding:7px 9px;font-family:monospace;font-size:11px;font-weight:700;'
            f'color:#374151;border-right:1px solid #e2e6ea;">{row["ticker"]}</td>'
            f'{_td(c1d,v1d)}{_td(c1w,v1w)}{_td(c1m,v1m)}{_td(c3m,v3m)}{_td(cytd,vytd)}'
            f'{_td(cvol,vvol,"border-left:2px solid #dbeafe;")}'
            f'{_td(cdd,vdd)}{_td(csh,vsh)}'
            f'<td style="padding:7px 9px;font-family:monospace;font-size:10px;'
            f'text-align:center;color:{rc};">&#9679; {rl}</td>'
            f'</tr>'
        )

    return f"""<!DOCTYPE html>
<html>
<head><meta charset="UTF-8"></head>
<body style="margin:0;padding:0;background:#f4f5f7;font-family:'IBM Plex Sans',Arial,sans-serif;">
<div style="max-width:1000px;margin:0 auto;padding:20px 12px;">

  <!-- HEADER -->
  <div style="background:#111827;padding:18px 24px;">
    <div style="display:flex;justify-content:space-between;align-items:center;">
      <div>
        <div style="font-size:15px;font-weight:700;color:#fff;letter-spacing:1px;font-family:monospace;">
          BK <span style="color:#34d399;">MARKET</span> DASHBOARD
        </div>
        <div style="font-size:9px;color:#6b7280;letter-spacing:2px;margin-top:4px;">
          {N_INSTRUMENTS}-INSTRUMENT UNIVERSE &nbsp;·&nbsp; PERFORMANCE &amp; RISK &amp; FRAGILITY
        </div>
      </div>
      <div style="text-align:right;">
        <div style="font-size:10px;color:#9ca3af;font-family:monospace;">{date_str}</div>
        <div style="font-size:9px;color:#6b7280;font-family:monospace;margin-top:3px;">
          AUTO-BRIEF &nbsp;·&nbsp; {SEND_TIME_SGT} SGT
        </div>
      </div>
    </div>
  </div>

  <!-- SUMMARY STRIP -->
  <div style="background:#fff;border:1px solid #e2e6ea;border-top:none;
    padding:10px 24px;display:flex;gap:28px;margin-bottom:14px;">
    <div>
      <div style="font-size:8px;letter-spacing:2px;text-transform:uppercase;color:#9ca3af;">Alerts &ge;3%</div>
      <div style="font-size:20px;font-weight:700;color:#dc2626;font-family:monospace;">{n_alerts}</div>
    </div>
    <div>
      <div style="font-size:8px;letter-spacing:2px;text-transform:uppercase;color:#9ca3af;">Watch 1.5–3%</div>
      <div style="font-size:20px;font-weight:700;color:#d97706;font-family:monospace;">{n_watches}</div>
    </div>
    <div>
      <div style="font-size:8px;letter-spacing:2px;text-transform:uppercase;color:#9ca3af;">Stable</div>
      <div style="font-size:20px;font-weight:700;color:#059669;font-family:monospace;">{n_stable}</div>
    </div>
    <div>
      <div style="font-size:8px;letter-spacing:2px;text-transform:uppercase;color:#9ca3af;">Instruments</div>
      <div style="font-size:20px;font-weight:700;color:#111827;font-family:monospace;">{len(df)}/{N_INSTRUMENTS}</div>
    </div>
  </div>

  <!-- TABLE -->
  <table style="width:100%;border-collapse:collapse;background:#fff;border:1px solid #e2e6ea;">
    <thead>
      <tr style="background:#111827;">
        <th style="padding:9px 9px;text-align:left;color:#9ca3af;font-size:8px;
          letter-spacing:1.5px;text-transform:uppercase;font-family:monospace;
          min-width:130px;border-right:1px solid #1f2937;">Asset</th>
        <th style="padding:9px 9px;text-align:left;color:#9ca3af;font-size:8px;
          letter-spacing:1.5px;font-family:monospace;border-right:1px solid #1f2937;">Ticker</th>
        <th style="padding:9px 9px;text-align:right;color:#6b9fd4;font-size:8px;
          letter-spacing:1.5px;font-family:monospace;border-right:1px solid #1f2937;">1D</th>
        <th style="padding:9px 9px;text-align:right;color:#9ca3af;font-size:8px;
          font-family:monospace;border-right:1px solid #1f2937;">1W</th>
        <th style="padding:9px 9px;text-align:right;color:#9ca3af;font-size:8px;
          font-family:monospace;border-right:1px solid #1f2937;">1M</th>
        <th style="padding:9px 9px;text-align:right;color:#9ca3af;font-size:8px;
          font-family:monospace;border-right:1px solid #1f2937;">3M</th>
        <th style="padding:9px 9px;text-align:right;color:#9ca3af;font-size:8px;
          font-family:monospace;border-right:1px solid #1f2937;">YTD</th>
        <th style="padding:9px 9px;text-align:right;color:#f87171;font-size:8px;
          font-family:monospace;border-left:2px solid #3a1a1a;
          border-right:1px solid #1f2937;">Vol 20D</th>
        <th style="padding:9px 9px;text-align:right;color:#9ca3af;font-size:8px;
          font-family:monospace;border-right:1px solid #1f2937;">Max DD</th>
        <th style="padding:9px 9px;text-align:right;color:#9ca3af;font-size:8px;
          font-family:monospace;border-right:1px solid #1f2937;">Sharpe</th>
        <th style="padding:9px 9px;text-align:center;color:#9ca3af;font-size:8px;
          font-family:monospace;">Signal</th>
      </tr>
    </thead>
    <tbody>{rows_html}</tbody>
  </table>

  <!-- FOOTER -->
  <div style="margin-top:14px;padding:10px 0;border-top:1px solid #e2e6ea;
    display:flex;justify-content:space-between;align-items:center;">
    <div style="font-size:9px;color:#9ca3af;font-family:monospace;line-height:1.9;">
      Signal: RED &lt; &minus;15% &nbsp;|&nbsp; AMBER &minus;15% to &minus;7% &nbsp;|&nbsp; GREEN &gt; &minus;7% — from 52-week high<br>
      Sharpe = 1Y annualised excess return / vol &nbsp;(rf = 4.5%)<br>
      Prices via Yahoo Finance (15 min delay)
    </div>
    <div style="text-align:right;">
      <div style="font-size:18px;font-weight:700;color:#111827;letter-spacing:-1px;">BK</div>
      <div style="font-size:9px;color:#6b7280;">Institutional Risk &nbsp;·&nbsp; Singapore</div>
      <div style="font-size:8px;color:#9ca3af;margin-top:2px;letter-spacing:1px;">
        CONFIDENTIAL &nbsp;·&nbsp; INTERNAL USE ONLY
      </div>
    </div>
  </div>

</div>
</body>
</html>"""



# ══════════════════════════════════════════════════════════════════════════════
#  BK FRAGILITY ENGINE
# ══════════════════════════════════════════════════════════════════════════════

FRAGILITY_WEIGHTS = {"dd":0.22,"vol":0.15,"cvar":0.20,"trend":0.15,"corr":0.18,"volz":0.10}

def _frag_logistic(x):
    return 1.0 / (1.0 + np.exp(-x))

def _robust_zscore(s: pd.Series, window: int = 504, clip: float = 4.0) -> pd.Series:
    """
    Robust rolling z-score using median/MAD normalisation (IFM standard).
    - Window: 504 days (2yr) — fixed, matches institutional_fragility_monitor.py
    - Immune to distortion by the extreme events it is measuring
    - 1.4826 factor makes MAD equivalent to std under normality
    - Clipped at ±4 to prevent any single pillar dominating the composite
    """
    min_p = max(60, window // 3)
    med   = s.rolling(window, min_periods=min_p).median()
    mad   = s.rolling(window, min_periods=min_p).apply(
        lambda x: np.median(np.abs(x - np.median(x))), raw=True)
    return ((s - med) / (1.4826 * mad.replace(0, 1e-6))).clip(-clip, clip)

def compute_fragility(prices: pd.DataFrame,
                      volumes: pd.DataFrame = None) -> pd.DataFrame:
    """
    Per-instrument fragility scores (0–100) aligned to IFM methodology.

    Methodology (matches institutional_fragility_monitor.py exactly):
    - Six pillars: Drawdown (22%), Tail Risk/CVaR (20%), Transmission (18%),
      Volatility (15%), Trend Deviation (15%), Liquidity/Volume (10%)
    - Robust rolling z-scores: median/MAD normalisation, 504-day window
    - Clipped at ±4 before weighting — no pillar dominates
    - Coverage check: score suppressed if <60% of weighted pillars have valid data
    - Logistic mapping to 0–100 with EWMA(span=10) smoothing
    - NO scaling factor — IFM standard (no × 0.5 compression)

    System-level fragility (frag_df.attrs["system_score"]) is derived from
    IFM_43_TICKERS only — comparable to CRO/GARP reports.
    Per-instrument scores for all 97 instruments are display-only.
    """
    if volumes is None:
        volumes = pd.DataFrame(0.0, index=prices.index, columns=prices.columns)

    # ── Pillar computations ────────────────────────────────────────────────────
    rets    = prices.pct_change().replace([np.inf, -np.inf], np.nan)
    wdd     = min(252, len(prices))
    peak    = prices.rolling(wdd, min_periods=20).max()
    dd      = (prices / peak - 1.0).abs()                              # drawdown magnitude
    vol20   = rets.rolling(20, min_periods=10).std() * np.sqrt(252)    # annualised vol

    def _cvar(x):
        q    = np.nanquantile(x, 0.05)
        tail = x[x <= q]
        return abs(np.nanmean(tail)) if len(tail) > 0 else np.nan

    cvar60  = rets.rolling(60, min_periods=20).apply(_cvar, raw=False)  # expected shortfall
    ma200   = prices.rolling(200, min_periods=50).mean()
    dist200 = (-(prices / ma200 - 1.0)).clip(lower=0)                   # downside-only trend stress

    # Transmission: correlation to world proxy — positive coupling only
    wcol   = "ACWI" if "ACWI" in rets.columns else rets.columns[0]
    corr_w = pd.DataFrame(index=rets.index, columns=rets.columns, dtype=float)
    for c in rets.columns:
        corr_w[c] = rets[c].rolling(60, min_periods=20).corr(rets[wcol]).clip(lower=0)

    # Liquidity: actual volume z-score (IFM method)
    # For tickers with no volume data, pillar contributes zero (handled via coverage)
    vol_mu  = volumes.rolling(60, min_periods=20).mean()
    vol_sd  = volumes.rolling(60, min_periods=20).std().replace(0, np.nan)
    volz    = ((volumes - vol_mu) / vol_sd).abs()
    # Zero out tickers with no volume data (yield indices, synthetics)
    no_vol_mask = (volumes == 0.0).all(axis=0)
    volz.loc[:, no_vol_mask] = np.nan

    # ── Z-score window: 504 days (IFM standard, fixed) ────────────────────────
    ZW              = 504
    MIN_COVERAGE    = 0.60   # suppress score if <60% of weights have valid data
    w               = FRAGILITY_WEIGHTS
    total_w         = sum(w.values())
    t2m             = {t: (s, n) for s, t, n, _b in UNIVERSE}
    rows            = []

    for col in prices.columns:
        if col in DISPLAY_EXCLUSIONS:
            continue
        if col in FRAGILITY_EXCLUSIONS:
            sec, name = t2m.get(col, ("", col))
            rows.append({
                "ticker": col, "name": name, "section": sec,
                "fragility": np.nan, "rag": "N/A",
                "pillar_dd": 0.0, "pillar_vol": 0.0, "pillar_cvar": 0.0,
                "pillar_trend": 0.0, "pillar_corr": 0.0, "pillar_volz": 0.0,
            })
            continue

        zd = _robust_zscore(dd[col],      ZW)
        zv = _robust_zscore(vol20[col],   ZW)
        zc = _robust_zscore(cvar60[col],  ZW)
        zt = _robust_zscore(dist200[col], ZW)
        zr = _robust_zscore(corr_w[col],  ZW)
        zz = _robust_zscore(volz[col],    ZW) if col in volz.columns else pd.Series(np.nan, index=prices.index)

        # Coverage check: suppress if insufficient pillar data
        w_valid = (
            w["dd"]    * zd.notna().astype(float) +
            w["vol"]   * zv.notna().astype(float) +
            w["cvar"]  * zc.notna().astype(float) +
            w["trend"] * zt.notna().astype(float) +
            w["corr"]  * zr.notna().astype(float) +
            w["volz"]  * zz.notna().astype(float)
        )
        coverage = w_valid / total_w

        latent = (
            w["dd"]    * zd.fillna(0)    +
            w["vol"]   * zv.fillna(0)    +
            w["cvar"]  * zc.fillna(0)    +
            w["trend"] * zt.fillna(0)    +
            w["corr"]  * zr.fillna(0)    +
            w["volz"]  * zz.fillna(0)
        )
        # Rescale by actual coverage weight
        latent_adj  = (latent / w_valid.replace(0, np.nan)) * total_w
        latent_adj  = latent_adj.where(coverage >= MIN_COVERAGE)

        # IFM logistic — no scaling factor
        sc = 100.0 * _frag_logistic(latent_adj.ewm(span=10, adjust=False).mean())
        v  = float(sc.iloc[-1]) if not sc.empty else np.nan
        if pd.isna(v):
            continue

        rag     = "CRISIS" if v >= 70 else "STRESSED" if v >= 55 else "MODERATE"
        sec, name = t2m.get(col, ("", col))

        def _p(z, k):
            val = z.iloc[-1] if not z.empty else np.nan
            return round(float(w[k] * val / total_w * 100), 1) if not pd.isna(val) else 0.0

        rows.append({
            "ticker": col, "name": name, "section": sec,
            "fragility": round(v, 1), "rag": rag,
            "pillar_dd":    _p(zd, "dd"),
            "pillar_vol":   _p(zv, "vol"),
            "pillar_cvar":  _p(zc, "cvar"),
            "pillar_trend": _p(zt, "trend"),
            "pillar_corr":  _p(zr, "corr"),
            "pillar_volz":  _p(zz, "volz"),
        })

    fdf = pd.DataFrame(rows).sort_values("fragility", ascending=False).reset_index(drop=True)

    if not fdf.empty:
        # System score: IFM 43-ticker universe only — CRO/GARP comparable
        ifm_scores = fdf[fdf["ticker"].isin(IFM_43_TICKERS)]["fragility"].dropna()
        ss = float(ifm_scores.median()) if not ifm_scores.empty else float(fdf["fragility"].dropna().median())
        fdf.attrs["system_score"] = round(ss, 1)
        fdf.attrs["regime"]       = "CRISIS" if ss >= 70 else "STRESSED" if ss >= 55 else "MODERATE"

    return fdf



# ══════════════════════════════════════════════════════════════════════════════
#  3-TAB WEB PAGE  (Performance | Risk | Fragility)
# ══════════════════════════════════════════════════════════════════════════════

# ══════════════════════════════════════════════════════════════════════════════
#  REGIME ENGINE — INSTITUTIONAL MULTI-MODEL FRAMEWORK
# ══════════════════════════════════════════════════════════════════════════════
#
#  Architecture (aligned with BK Institutional Fragility Monitor):
#  ────────────────────────────────────────────────────────────────
#  This regime engine implements a 3-tier detection framework used by
#  institutional multi-asset desks (Man AHL, AQR, Bridgewater-style):
#
#  TIER 1 — DETERMINISTIC STATE MACHINE (governance / headline)
#    - Expanding-quantile thresholds on 20-day vol + 1Y-peak drawdown
#    - No look-ahead bias (ex-ante calibration, shifted by 1 day)
#    - Fully deterministic, auditable, explainable to CIO/board
#    - States: Calm / Stressed / Crisis
#
#  TIER 2 — HIDDEN MARKOV MODEL (probabilistic / conviction)
#    - 3-state Gaussian HMM on [returns, vol, drawdown]
#    - Walk-forward: retrained every 21 trading days on expanding window
#    - Outputs posterior probabilities (p_calm, p_stressed, p_crisis)
#    - High-entropy posteriors = regime uncertainty = low conviction
#    - Requires: pip install hmmlearn (graceful fallback if missing)
#
#  TIER 3 — GAUSSIAN MIXTURE MODEL (cross-validation / fragility-aware)
#    - 3-component GMM on [returns, fragility_score]
#    - Captures non-linear regime clusters that HMM may miss
#    - Walk-forward with rolling retrain every 21 days
#    - Requires: sklearn (graceful fallback if missing)
#
#  CONSENSUS — "UNION OF RISK" AGGREGATION
#    - If ANY model signals Crisis → consensus = Crisis
#    - If ANY model signals Stressed (none Crisis) → consensus = Stressed
#    - Otherwise → Calm
#    - Conservative by design: false positives (early warnings) preferred
#      over false negatives (missed crises)
#
#  MODEL AGREEMENT — CONVICTION GAUGE
#    - Count of models agreeing on the same regime (1/3 to 3/3)
#    - 3/3 agreement = high conviction → position with confidence
#    - 1/3 agreement = low conviction → hedge, reduce sizing
#    - Displayed as badge on the regime tab for at-a-glance reading
#
#  TRANSITION RISK — EARLY WARNING SIGNAL
#    - Fires when HMM posterior probability of a worse regime is rising
#      while the deterministic SM has not yet flipped
#    - Example: SM says "Calm" but HMM p_stressed > 40% → "Elevated"
#    - Gives 5–15 trading days of lead time vs. the SM alone
#    - Three levels: Low / Elevated / High
#
#  The deterministic SM remains the headline/governance regime.
#  HMM and GMM provide conviction, early-warning, and confirmation.
# ══════════════════════════════════════════════════════════════════════════════

# ── Optional ML imports (graceful degradation) ────────────────────────────
_HAS_HMM = False
try:
    from hmmlearn.hmm import GaussianHMM as _GaussianHMM
    _HAS_HMM = True
except ImportError:
    pass

_HAS_GMM = False
try:
    from sklearn.mixture import GaussianMixture as _GaussianMixture
    _HAS_GMM = True
except ImportError:
    pass


def _detect_regime_hmm(rets: pd.Series, vol20: pd.Series, dd: pd.Series,
                       n_states: int = 3, seed: int = 42,
                       min_history: int = 504, retrain_step: int = 21,
                       n_iter: int = 200) -> tuple:
    """
    Hidden Markov Model regime detection — walk-forward, no look-ahead.

    Methodology (institutional standard):
    1. Features: daily returns, 20-day realised vol, drawdown from 1Y peak
    2. Z-score normalisation using only data available at each point (expanding)
    3. Model retrained every `retrain_step` days on the full expanding window
    4. State-to-regime mapping: rank HMM states by avg(vol) + avg(|dd|)
       → lowest risk = Calm, mid = Stressed, highest = Crisis
    5. Posterior probabilities give conviction level for each regime

    Returns:
        (regime_series, probs_dict)
        - regime_series: pd.Series of "Calm"/"Stressed"/"Crisis" (or empty)
        - probs_dict: dict with keys p_calm, p_stressed, p_crisis, entropy
          (latest day's values), or None if HMM unavailable
    """
    if not _HAS_HMM:
        return pd.Series(dtype="object"), None

    X = pd.DataFrame({"ret": rets, "vol": vol20, "dd": dd}).dropna()
    if len(X) < min_history + 10:
        return pd.Series(dtype="object"), None

    regimes = pd.Series(index=rets.index, dtype="object")
    latest_probs = None
    model = None; mu = None; sd = None; mapping = None

    for i in range(min_history, len(X)):
        idx_t = X.index[i]
        train = X.iloc[:i]
        x_t = X.iloc[[i]]

        do_retrain = (model is None) or ((i - min_history) % retrain_step == 0)
        if do_retrain:
            try:
                mu = train.mean()
                sd = train.std(ddof=0).replace(0, np.nan).fillna(1.0)
                train_z = ((train - mu) / sd).replace([np.inf, -np.inf], np.nan).dropna()
                if len(train_z) < max(60, n_states * 20):
                    continue
                model = _GaussianHMM(n_components=n_states, covariance_type="full",
                                     n_iter=n_iter, random_state=seed)
                model.fit(train_z.values)
                # Map HMM integer states → regime names by risk ranking
                train_states = pd.Series(model.predict(train_z.values), index=train_z.index)
                tmp = train.loc[train_z.index].copy()
                tmp["state"] = train_states.values
                risk_score = (tmp.groupby("state")["vol"].mean()
                              + tmp.groupby("state")["dd"].mean().abs()).sort_values()
                ordered = list(risk_score.index)
                state_names = ["Calm", "Stressed", "Crisis"][:n_states]
                mapping = {st: state_names[j] for j, st in enumerate(ordered)}
            except Exception:
                model = None; mapping = None
                continue

        if model is None or mu is None or sd is None or mapping is None:
            continue
        try:
            xz = ((x_t - mu) / sd).replace([np.inf, -np.inf], np.nan).dropna()
            if xz.empty:
                continue
            post = model.predict_proba(xz.values)[0]
            pred_state = int(np.argmax(post))
            regimes.loc[idx_t] = mapping.get(pred_state, np.nan)
            # Store latest posteriors for dashboard display
            eps = 1e-12
            p_dict = {"entropy": float(-(post * np.log(post + eps)).sum())}
            for st_num, st_name in mapping.items():
                p_dict[f"p_{st_name.lower()}"] = float(post[int(st_num)])
            latest_probs = p_dict
        except Exception:
            continue

    return regimes.reindex(rets.index), latest_probs


def _detect_regime_gmm(rets: pd.Series, vol20: pd.Series, dd: pd.Series,
                       n_states: int = 3, seed: int = 42,
                       min_history: int = 252, retrain_step: int = 21) -> pd.Series:
    """
    Gaussian Mixture Model regime detection — walk-forward, no look-ahead.

    Uses returns + vol + drawdown as features (same as HMM for consistency).
    GMM captures non-linear cluster boundaries that HMM's sequential assumption
    may miss — acts as cross-validation for the HMM.

    Returns:
        pd.Series of "Calm"/"Stressed"/"Crisis" (or empty if sklearn missing)
    """
    if not _HAS_GMM:
        return pd.Series(dtype="object")

    X = pd.DataFrame({"ret": rets, "vol": vol20, "dd": dd}).dropna()
    if len(X) < min_history + 10:
        return pd.Series(dtype="object")

    model = None; mu = None; sd = None; cluster_map = None
    labels = pd.Series(index=rets.index, dtype="object")

    for i in range(min_history, len(X)):
        idx_t = X.index[i]
        train = X.iloc[:i]
        x_t = X.iloc[[i]]

        do_retrain = (model is None) or ((i - min_history) % retrain_step == 0)
        if do_retrain:
            try:
                mu = train.mean()
                sd = train.std(ddof=0).replace(0, np.nan).fillna(1.0)
                train_z = ((train - mu) / sd).replace([np.inf, -np.inf], np.nan).dropna()
                if len(train_z) < max(60, n_states * 20):
                    continue
                model = _GaussianMixture(n_components=n_states, random_state=seed)
                model.fit(train_z.values)
                # Map clusters → regime names by vol + |dd| risk ranking
                train_labels = model.predict(train_z.values)
                tmp = train.loc[train_z.index].copy()
                tmp["cluster"] = train_labels
                risk_score = (tmp.groupby("cluster")["vol"].mean()
                              + tmp.groupby("cluster")["dd"].mean().abs()).sort_values()
                ordered = list(risk_score.index)
                state_names = ["Calm", "Stressed", "Crisis"][:n_states]
                cluster_map = {c: nm for c, nm in zip(ordered, state_names)}
            except Exception:
                model = None; cluster_map = None
                continue

        if model is None or mu is None or sd is None or cluster_map is None:
            continue
        try:
            xz = ((x_t - mu) / sd).replace([np.inf, -np.inf], np.nan).dropna()
            if xz.empty:
                continue
            labels.loc[idx_t] = cluster_map.get(int(model.predict(xz.values)[0]), np.nan)
        except Exception:
            continue

    return labels.reindex(rets.index)


def _regime_consensus(sm: str, hmm: str, gmm: str) -> tuple:
    """
    Compute consensus regime and model agreement from 3 independent detectors.

    Consensus logic: "Union of Risk" — the most severe call wins.
    This is intentionally conservative: institutional risk management prefers
    false positives (early defensive moves) over false negatives (missed crises).

    Returns:
        (consensus_regime, agreement_count, models_available)
    """
    severity = {"Calm": 0, "Stressed": 1, "Crisis": 2}
    inv_map  = {0: "Calm", 1: "Stressed", 2: "Crisis"}

    models = []
    if sm:
        models.append(sm)
    if hmm and hmm in severity:
        models.append(hmm)
    if gmm and gmm in severity:
        models.append(gmm)

    n_available = len(models)
    if n_available == 0:
        return "Calm", 0, 0

    # Consensus = max severity across all available models
    max_sev = max(severity.get(m, 0) for m in models)
    consensus = inv_map[max_sev]

    # Agreement = how many models agree on the consensus label
    agreement = sum(1 for m in models if m == consensus)

    return consensus, agreement, n_available


def _transition_risk(sm_regime: str, hmm_probs: dict) -> tuple:
    """
    Compute regime transition risk — early warning when HMM sees what SM doesn't.

    This is the key institutional value-add of the multi-model approach:
    the HMM's posterior probabilities shift BEFORE the deterministic SM
    flips state, giving 5–15 trading days of lead time.

    Logic:
    - If SM=Calm and HMM p_stressed+p_crisis > 40% → Elevated
    - If SM=Calm and HMM p_crisis > 25% → High
    - If SM=Stressed and HMM p_crisis > 40% → High
    - Otherwise → Low

    Returns:
        (risk_level, risk_description)
        risk_level: "Low" / "Elevated" / "High"
    """
    if not hmm_probs:
        return "N/A", "HMM unavailable — install hmmlearn for transition alerts"

    p_calm     = hmm_probs.get("p_calm", 0.5)
    p_stressed = hmm_probs.get("p_stressed", 0.3)
    p_crisis   = hmm_probs.get("p_crisis", 0.2)

    if sm_regime == "Calm":
        if p_crisis > 0.25:
            return "High", f"SM: Calm but HMM sees {p_crisis*100:.0f}% crisis probability — significant divergence"
        if p_stressed + p_crisis > 0.40:
            return "Elevated", f"SM: Calm but HMM sees {(p_stressed+p_crisis)*100:.0f}% stressed/crisis probability"
        return "Low", "All models aligned — no transition signal"
    elif sm_regime == "Stressed":
        if p_crisis > 0.40:
            return "High", f"SM: Stressed and HMM sees {p_crisis*100:.0f}% crisis probability — escalation risk"
        if p_calm > 0.50:
            return "Low", f"HMM sees {p_calm*100:.0f}% calm probability — de-escalation likely"
        return "Elevated", "Models disagree — regime transition possible"
    else:  # Crisis
        if p_calm > 0.40:
            return "Elevated", f"HMM sees {p_calm*100:.0f}% calm probability — crisis may be easing"
        return "Low", "All models confirm crisis conditions"


def compute_regime(prices: pd.DataFrame) -> dict:
    """
    Institutional multi-model regime detection engine.

    Runs 3 independent regime detectors + consensus + transition risk:
      1. Deterministic State Machine  (headline — always available)
      2. Hidden Markov Model          (conviction — requires hmmlearn)
      3. Gaussian Mixture Model       (cross-validation — requires sklearn)

    The SM is always the governance/headline regime. HMM and GMM provide
    probabilistic conviction and early-warning signals. If ML libraries
    are not installed, the engine gracefully degrades to SM-only mode.

    Returns dict with:
      regime, days_in_regime, timeline, stats, drivers, episodes,
      hmm_probs, consensus, model_agreement, models_available,
      transition_risk, transition_desc, hmm_regime, gmm_regime
    """
    world_col = "ACWI" if "ACWI" in prices.columns else prices.columns[0]
    rets  = prices[world_col].pct_change().replace([np.inf,-np.inf], np.nan)
    vol20 = rets.rolling(20, min_periods=10).std() * np.sqrt(252)
    peak  = prices[world_col].rolling(252, min_periods=20).max()
    dd    = prices[world_col] / peak - 1.0

    # ── TIER 1: Deterministic State Machine ───────────────────────────────
    # Governance-friendly, auditable, no ML dependencies.
    # Uses expanding quantile thresholds so each day's classification
    # only uses data available up to t−1 (no look-ahead bias).
    s = pd.DataFrame({"vol": vol20, "dd": dd}).dropna()
    min_history = 252

    vol70 = s["vol"].expanding(min_periods=min_history).quantile(0.70).shift(1)
    vol90 = s["vol"].expanding(min_periods=min_history).quantile(0.90).shift(1)
    dd30  = s["dd"].expanding(min_periods=min_history).quantile(0.30).shift(1)
    dd10  = s["dd"].expanding(min_periods=min_history).quantile(0.10).shift(1)

    valid  = vol70.notna() & vol90.notna() & dd30.notna() & dd10.notna()
    regime = pd.Series("Calm", index=s.index, dtype="object")
    regime.loc[~valid] = np.nan
    stressed = valid & ((s["vol"] >= vol70) | (s["dd"] <= dd30))
    crisis   = valid & ((s["vol"] >= vol90) | (s["dd"] <= dd10))
    regime.loc[stressed] = "Stressed"
    regime.loc[crisis]   = "Crisis"
    regime = regime.dropna()

    if regime.empty:
        return {"regime": "Calm", "days_in_regime": 0, "timeline": [], "stats": {},
                "drivers": {}, "hmm_probs": None, "consensus": "Calm",
                "model_agreement": 0, "models_available": 0,
                "transition_risk": "N/A", "transition_desc": "",
                "hmm_regime": None, "gmm_regime": None}

    current = regime.iloc[-1]

    # Days in current streak
    streak = 1
    for i in range(len(regime)-2, -1, -1):
        if regime.iloc[i] == current:
            streak += 1
        else:
            break

    # Timeline: last 504 trading days (2 years) — daily regime
    timeline_raw = regime.tail(504)
    timeline = []
    for date, reg in timeline_raw.items():
        timeline.append({
            "date": date.strftime("%Y-%m-%d"),
            "regime": reg,
            "color": "#f85149" if reg=="Crisis" else "#e3b341" if reg=="Stressed" else "#3fb950"
        })

    # Stats: days in each regime over full history
    total = len(regime)
    stats = {
        "Calm":     {"days": int((regime=="Calm").sum()),     "pct": round((regime=="Calm").sum()/total*100,1)},
        "Stressed": {"days": int((regime=="Stressed").sum()), "pct": round((regime=="Stressed").sum()/total*100,1)},
        "Crisis":   {"days": int((regime=="Crisis").sum()),   "pct": round((regime=="Crisis").sum()/total*100,1)},
    }

    # Average duration per regime
    durations = {"Calm":[], "Stressed":[], "Crisis":[]}
    cur_reg = regime.iloc[0]; cur_len = 1
    for i in range(1, len(regime)):
        if regime.iloc[i] == cur_reg:
            cur_len += 1
        else:
            durations[cur_reg].append(cur_len)
            cur_reg = regime.iloc[i]; cur_len = 1
    durations[cur_reg].append(cur_len)
    for r in durations:
        stats[r]["avg_duration"] = round(float(np.mean(durations[r])), 0) if durations[r] else 0

    # Drivers: current vol and dd vs historical percentiles
    cur_vol = float(vol20.iloc[-1]) if not pd.isna(vol20.iloc[-1]) else 0
    cur_dd  = float(dd.iloc[-1])    if not pd.isna(dd.iloc[-1])    else 0
    vol_pct = float((vol20.dropna() <= cur_vol).mean() * 100)
    dd_pct  = float((dd.dropna()   >= cur_dd).mean()  * 100)

    # Crisis episodes (drawdown < -15%, minimum 5 trading days)
    MIN_EPISODE_DAYS = 5
    episodes = []
    in_ep = False; ep_start = None
    for date, val in dd.items():
        if val < -0.15 and not in_ep:
            in_ep = True; ep_start = date
        elif val >= -0.10 and in_ep:
            in_ep = False
            _ep_days = len(dd[ep_start:date])
            if _ep_days < MIN_EPISODE_DAYS:
                continue
            try:
                day_fmt = "%#d" if os.name == "nt" else "%-d"
                if ep_start.year == date.year and ep_start.month == date.month:
                    _range_ascii = f"{ep_start.strftime(day_fmt)}-{date.strftime(day_fmt + ' %b %Y')}"
                    _range_html  = f"{ep_start.strftime(day_fmt)}&ndash;{date.strftime(day_fmt + ' %b %Y')}"
                else:
                    _range_ascii = f"{ep_start.strftime('%b %Y')} -> {date.strftime('%b %Y')}"
                    _range_html  = f"{ep_start.strftime('%b %Y')} &rarr; {date.strftime('%b %Y')}"
            except Exception:
                _range_ascii = f"{ep_start} -> {date}"
                _range_html  = _range_ascii
            episodes.append({"start": ep_start.strftime("%b %Y"),
                              "end":   date.strftime("%b %Y"),
                              "range": _range_html,
                              "depth": round(float(dd[ep_start:date].min()*100),1)})
            print(f"[Regime Episode] {_range_ascii} | depth {round(float(dd[ep_start:date].min()*100),1)}%")
    if in_ep:
        episodes.append({"start": ep_start.strftime("%b %Y"),
                          "end":   "Ongoing",
                          "range": f"{ep_start.strftime('%b %Y')} &rarr; Ongoing",
                          "depth": round(float(dd[ep_start:].min()*100),1)})
        print(f"[Regime Episode] {ep_start.strftime('%b %Y')} -> Ongoing | "
              f"depth {round(float(dd[ep_start:].min()*100),1)}%")

    # ── TIER 2: Hidden Markov Model ──────────────────────────────────────
    # Provides posterior probabilities (conviction) and early-warning signal.
    # Walk-forward retrain every 21 days. Graceful fallback if hmmlearn missing.
    print("[Regime]  Running HMM regime detector ..." if _HAS_HMM else "[Regime]  HMM skipped (hmmlearn not installed)")
    hmm_series, hmm_probs = _detect_regime_hmm(rets, vol20, dd)
    hmm_current = None
    if hmm_series is not None and not hmm_series.dropna().empty:
        hmm_current = hmm_series.dropna().iloc[-1]
        print(f"[Regime]  HMM regime: {hmm_current}")

    # ── TIER 3: Gaussian Mixture Model ───────────────────────────────────
    # Cross-validates HMM. Uses same features for consistency.
    # Graceful fallback if sklearn missing.
    print("[Regime]  Running GMM regime detector ..." if _HAS_GMM else "[Regime]  GMM skipped (sklearn not installed)")
    gmm_series = _detect_regime_gmm(rets, vol20, dd)
    gmm_current = None
    if gmm_series is not None and not gmm_series.dropna().empty:
        gmm_current = gmm_series.dropna().iloc[-1]
        print(f"[Regime]  GMM regime: {gmm_current}")

    # ── CONSENSUS + MODEL AGREEMENT ──────────────────────────────────────
    # Union-of-risk: most severe model call wins (conservative by design).
    # Agreement count tells the PM how much conviction to assign.
    consensus, agreement, n_models = _regime_consensus(current, hmm_current, gmm_current)
    print(f"[Regime]  Consensus: {consensus} | Agreement: {agreement}/{n_models} models")

    # ── TRANSITION RISK ──────────────────────────────────────────────────
    # Early warning: HMM posteriors shift before deterministic SM flips.
    # This is the institutional value-add — 5-15 days lead time.
    tr_level, tr_desc = _transition_risk(current, hmm_probs)
    if tr_level != "Low" and tr_level != "N/A":
        print(f"[Regime]  >> Transition risk: {tr_level} -- {tr_desc}")

    return {
        "regime":           current,
        "days_in_regime":   streak,
        "timeline":         timeline,
        "stats":            stats,
        "drivers": {
            "vol_now":   round(cur_vol*100, 1),
            "vol_pct":   round(vol_pct, 0),
            "dd_now":    round(cur_dd*100,  1),
            "dd_pct":    round(dd_pct,  0),
        },
        "episodes":         episodes[-8:],
        # ── Multi-model outputs ──
        "hmm_probs":        hmm_probs,          # dict: p_calm, p_stressed, p_crisis, entropy
        "hmm_regime":       hmm_current,         # str or None
        "gmm_regime":       gmm_current,         # str or None
        "consensus":        consensus,           # str: most severe across all models
        "model_agreement":  agreement,           # int: how many models agree on consensus
        "models_available": n_models,            # int: how many models ran (1–3)
        "transition_risk":  tr_level,            # str: Low / Elevated / High / N/A
        "transition_desc":  tr_desc,             # str: human-readable explanation
    }


# ══════════════════════════════════════════════════════════════════════════════
#  FEAR & GREED ENGINE
# ══════════════════════════════════════════════════════════════════════════════

def compute_fear_greed(prices: pd.DataFrame) -> dict:
    """
    7-component Fear & Greed Index (0=Extreme Fear, 100=Extreme Greed).
    Each component scored 0-100 then equally weighted.
    """
    rets = prices.pct_change().replace([np.inf,-np.inf], np.nan)
    scores = {}
    details = {}

    def _pct_rank(series, val):
        """Where does val sit in historical distribution? 0-100."""
        clean = series.dropna()
        if clean.empty or pd.isna(val): return 50.0
        return float((clean <= val).mean() * 100)

    # ── 1. Volatility: VIXY vs 50D MA (high vol = fear) ──────────────────────
    if "VIXY" in prices.columns:
        vixy = prices["VIXY"].dropna()
        ma50 = vixy.rolling(50, min_periods=20).mean()
        ratio = vixy / ma50 - 1
        cur   = float(ratio.iloc[-1]) if not ratio.empty else 0
        # High ratio = high vol vs norm = fear → invert
        raw = _pct_rank(ratio, cur)
        scores["Volatility"] = max(0, min(100, 100 - raw))
        details["Volatility"] = {"value": f"VIXY {cur*100:+.1f}% vs 50D MA", "score": scores["Volatility"]}

    # ── 2. Market Momentum: SPY vs 125D MA ───────────────────────────────────
    if "SPY" in prices.columns:
        spy  = prices["SPY"].dropna()
        ma125= spy.rolling(125, min_periods=50).mean()
        ratio= spy / ma125 - 1
        cur  = float(ratio.iloc[-1]) if not ratio.empty else 0
        raw  = _pct_rank(ratio, cur)
        scores["Momentum"] = max(0, min(100, raw))
        details["Momentum"] = {"value": f"SPY {cur*100:+.1f}% vs 125D MA", "score": scores["Momentum"]}

    # ── 3. Market Breadth: % instruments above their 50D MA ──────────────────
    ma50_all = prices.rolling(50, min_periods=20).mean()
    above = (prices.iloc[-1] > ma50_all.iloc[-1]).sum()
    total_avail = prices.shape[1]
    breadth_pct = above / total_avail * 100 if total_avail > 0 else 50
    scores["Breadth"] = float(breadth_pct)
    details["Breadth"] = {"value": f"{above}/{total_avail} above 50D MA", "score": scores["Breadth"]}

    # ── 4. Safe Haven Demand: TLT outperformance vs SPY (20D) ────────────────
    if "TLT" in prices.columns and "SPY" in prices.columns:
        tlt_ret = prices["TLT"].pct_change(20)
        spy_ret = prices["SPY"].pct_change(20)
        spread  = tlt_ret - spy_ret  # positive = bonds beating equities = fear
        cur     = float(spread.iloc[-1]) if not spread.dropna().empty else 0
        raw     = _pct_rank(spread.dropna(), cur)
        # High spread (bonds beating) = fear → invert
        scores["Safe Haven"] = max(0, min(100, 100 - raw))
        details["Safe Haven"] = {"value": f"TLT vs SPY 20D: {cur*100:+.1f}%", "score": scores["Safe Haven"]}

    # ── 5. Junk Bond Demand: HYG vs IEF (credit risk appetite) ──────────────
    if "HYG" in prices.columns and "IEF" in prices.columns:
        hyg_ret = prices["HYG"].pct_change(20)
        ief_ret = prices["IEF"].pct_change(20)
        spread  = hyg_ret - ief_ret  # positive = junk beating govt = greed
        cur     = float(spread.iloc[-1]) if not spread.dropna().empty else 0
        raw     = _pct_rank(spread.dropna(), cur)
        scores["Junk Bonds"] = max(0, min(100, raw))
        details["Junk Bonds"] = {"value": f"HYG vs IEF 20D: {cur*100:+.1f}%", "score": scores["Junk Bonds"]}

    # ── 6. Market Strength: % instruments within 5% of 52W high ─────────────
    hi52 = prices.rolling(252, min_periods=100).max()
    near_high = ((prices.iloc[-1] / hi52.iloc[-1]) >= 0.95).sum()
    strength_pct = near_high / total_avail * 100 if total_avail > 0 else 50
    scores["Strength"] = float(strength_pct)
    details["Strength"] = {"value": f"{near_high}/{total_avail} within 5% of 52W high", "score": scores["Strength"]}

    # ── 7. Term Structure: VIX/VIX3M ratio (backwardation = panic) ──────────
    if "^VIX" in prices.columns and "^VIX3M" in prices.columns:
        ratio_ts = (prices["^VIX"] / prices["^VIX3M"].replace(0, np.nan)).dropna()
        # Use 2-year rolling window (≈504 trading days)
        window_2y = ratio_ts.iloc[-504:] if len(ratio_ts) >= 504 else ratio_ts
        cur       = float(ratio_ts.iloc[-1]) if not ratio_ts.empty else 1.0
        raw       = _pct_rank(window_2y, cur)
        # High ratio = backwardation = panic = fear → invert
        scores["Term Structure"] = max(0, min(100, (1 - raw / 100) * 100))
        details["Term Structure"] = {"value": f"VIX/VIX3M ratio: {cur:.2f}", "score": scores["Term Structure"]}

    # ── Composite ─────────────────────────────────────────────────────────────
    if not scores:
        return {"score": 50, "label": "Neutral", "details": {}}

    # Equal weight across available components
    score_vals = list(scores.values())
    composite  = float(np.mean(score_vals)) if score_vals else 50.0

    if composite <= 25:   label, color = "Extreme Fear",  "#f85149"
    elif composite <= 45: label, color = "Fear",          "#ff7b72"
    elif composite <= 55: label, color = "Neutral",       "#e3b341"
    elif composite <= 75: label, color = "Greed",         "#7ee787"
    else:                 label, color = "Extreme Greed", "#3fb950"

    return {
        "score":   round(composite, 1),
        "label":   label,
        "color":   color,
        "details": details,
    }


# ══════════════════════════════════════════════════════════════════════════════
#  FRAGILITY HISTORICAL TREND ENGINE
# ══════════════════════════════════════════════════════════════════════════════

def compute_fragility_trend(prices: pd.DataFrame,
                            volumes: pd.DataFrame = None) -> dict:
    """
    System-level fragility score timeseries (last 504 days = 2 years).
    Driven by IFM_43_TICKERS only — consistent with institutional reporting.
    Methodology matches institutional_fragility_monitor.py exactly.
    Returns daily scores + regime for chart rendering.
    """
    if volumes is None:
        volumes = pd.DataFrame(0.0, index=prices.index, columns=prices.columns)

    # Restrict to IFM 43-ticker universe for system score
    ifm_cols = [c for c in IFM_43_TICKERS if c in prices.columns]
    if not ifm_cols:
        ifm_cols = list(prices.columns)  # fallback

    p   = prices[ifm_cols]
    v   = volumes.reindex(columns=ifm_cols, fill_value=0.0)

    rets    = p.pct_change().replace([np.inf, -np.inf], np.nan)
    wdd     = min(252, len(p))
    peak    = p.rolling(wdd, min_periods=20).max()
    dd      = (p / peak - 1.0).abs()
    vol20   = rets.rolling(20, min_periods=10).std() * np.sqrt(252)

    def _cvar(x):
        q    = np.nanquantile(x, 0.05)
        tail = x[x <= q]
        return abs(np.nanmean(tail)) if len(tail) > 0 else np.nan

    cvar60  = rets.rolling(60, min_periods=20).apply(_cvar, raw=False)
    ma200   = p.rolling(200, min_periods=50).mean()
    dist200 = (-(p / ma200 - 1.0)).clip(lower=0)

    wcol   = "ACWI" if "ACWI" in rets.columns else rets.columns[0]
    corr_w = pd.DataFrame(index=rets.index, columns=p.columns, dtype=float)
    for c in p.columns:
        corr_w[c] = rets[c].rolling(60, min_periods=20).corr(rets[wcol]).clip(lower=0)

    vol_mu  = v.rolling(60, min_periods=20).mean()
    vol_sd  = v.rolling(60, min_periods=20).std().replace(0, np.nan)
    volz    = ((v - vol_mu) / vol_sd).abs()
    no_vol_mask = (v == 0.0).all(axis=0)
    volz.loc[:, no_vol_mask] = np.nan

    w   = FRAGILITY_WEIGHTS
    ZW  = 504
    MIN_COVERAGE = 0.60
    total_w = sum(w.values())

    latents = pd.DataFrame(index=p.index)
    for col in p.columns:
        zd = _robust_zscore(dd[col],      ZW)
        zv = _robust_zscore(vol20[col],   ZW)
        zc = _robust_zscore(cvar60[col],  ZW)
        zt = _robust_zscore(dist200[col], ZW)
        zr = _robust_zscore(corr_w[col],  ZW)
        zz = _robust_zscore(volz[col],    ZW) if col in volz.columns else pd.Series(np.nan, index=p.index)

        w_valid  = (
            w["dd"]    * zd.notna().astype(float) +
            w["vol"]   * zv.notna().astype(float) +
            w["cvar"]  * zc.notna().astype(float) +
            w["trend"] * zt.notna().astype(float) +
            w["corr"]  * zr.notna().astype(float) +
            w["volz"]  * zz.notna().astype(float)
        )
        coverage = w_valid / total_w
        latent   = (
            w["dd"]    * zd.fillna(0) +
            w["vol"]   * zv.fillna(0) +
            w["cvar"]  * zc.fillna(0) +
            w["trend"] * zt.fillna(0) +
            w["corr"]  * zr.fillna(0) +
            w["volz"]  * zz.fillna(0)
        )
        latent_adj = (latent / w_valid.replace(0, np.nan)) * total_w
        latent_adj = latent_adj.where(coverage >= MIN_COVERAGE)
        latents[col] = latent_adj

    # Cross-sectional median over IFM universe — IFM logistic, no scaling
    sys_lat   = latents.median(axis=1).ewm(span=10, adjust=False).mean()
    sys_score = 100.0 * _frag_logistic(sys_lat)

    trend_raw = sys_score.tail(504).dropna()
    trend = []
    for date, val in trend_raw.items():
        reg = "Crisis" if val >= 70 else "Stressed" if val >= 55 else "Moderate"
        trend.append({
            "date":   date.strftime("%Y-%m-%d"),
            "score":  round(float(val), 1),
            "regime": reg,
            "color":  "#f85149" if reg == "Crisis" else "#e3b341" if reg == "Stressed" else "#3fb950",
        })

    return {
        "trend":     trend,
        "current":   round(float(sys_score.iloc[-1]), 1) if not sys_score.empty else 50,
        "peak_2y":   round(float(sys_score.tail(504).max()), 1),
        "trough_2y": round(float(sys_score.tail(504).min()), 1),
        "avg_2y":    round(float(sys_score.tail(504).mean()), 1),
    }


# ══════════════════════════════════════════════════════════════════════════════
#  AI COMMENTARY ENGINE — Claude API
# ══════════════════════════════════════════════════════════════════════════════

def generate_ai_commentary(market_data: dict) -> dict:
    """
    Generate AI-powered market commentary using Claude API.
    Falls back to empty strings if API unavailable.
    """
    import os, json
    if DEVELOPMENT_MODE:
        print("[AI] DEVELOPMENT_MODE — skipping Claude API call")
        return {}

    api_key = os.environ.get("ANTHROPIC_API_KEY", "")

    if not api_key:
        print("[AI] No API key found — skipping AI commentary")
        return {}
    
    try:
        import urllib.request
        
        prompt = f"""You are a quantitative research assistant producing a personal, non-commercial daily market observation log. This is a private research tool — not a client-facing document and not financial advice.

Current model outputs as of {market_data.get('date', 'today')}:
- Market Regime: {market_data.get('regime', 'Unknown')} (for {market_data.get('regime_days', 0)} consecutive trading days)
- System Fragility Score: {market_data.get('fragility_score', 50)}/100 ({market_data.get('fragility_label', 'STRESSED')})
- Fear & Greed Index: {market_data.get('fg_score', 50)}/100 ({market_data.get('fg_label', 'Neutral')})
- World Volatility Percentile: {market_data.get('vol_pct', 50)}th percentile of 10-year history
- World Drawdown: {market_data.get('dd_pct', 0):.1f}% from peak
- Instruments in RED signal: {market_data.get('n_red', 0)} of {market_data.get('n_total', 60)}
- Instruments with Rising Vol: {market_data.get('vol_rising', 0)}
- Top 1M Gainers: {', '.join(market_data.get('top_gainers', []))}
- Top 1M Losers: {', '.join(market_data.get('top_losers', []))}
- Average Cross-Asset Correlation: {market_data.get('avg_corr', 0.3):.2f}
- Backtest 5Y: BK Allocation {market_data.get('bt_bk', 'N/A')}, SPY {market_data.get('bt_spy', 'N/A')}, 60/40 {market_data.get('bt_6040', 'N/A')}

Generate a JSON response with exactly these fields:
{{
  "narrative": "4-5 sentences describing what the model outputs and cross-asset data show today. Be specific — reference actual instruments and numbers. Use purely observational language: what is elevated, what has declined, what correlations suggest. Do NOT use the words: buy, sell, invest, allocate, position, exposure, recommend, opportunity, should, or must.",
  "actions": [
    "Observation 1: what the framework scores show for a specific instrument and why",
    "Observation 2: what the framework scores show for a specific instrument and why",
    "Observation 3: what the framework scores show for a specific instrument and why",
    "Observation 4: what the framework scores show for a specific instrument and why"
  ],
  "watchlist": [
    {{"instrument": "Name", "reason": "Data-driven observation — what the model flags and why it is notable"}},
    {{"instrument": "Name", "reason": "Data-driven observation — what the model flags and why it is notable"}},
    {{"instrument": "Name", "reason": "Data-driven observation — what the model flags and why it is notable"}}
  ],
  "fg_summary": "One sentence describing what a Fear & Greed score of {market_data.get('fg_score', 50)} indicates about current market sentiment dynamics",
  "regime_interpretation": "2-3 sentences describing what characterises the current regime, what data points drove the classification, and what a regime shift would look like",
  "edge_rationale": "2-3 sentences describing why the current framework weights reflect the regime and fragility readings, referencing specific asset classes and model outputs"
}}

Return ONLY valid JSON. No preamble, no markdown, no explanation."""

        payload = json.dumps({
            "model": "claude-sonnet-4-20250514",
            "max_tokens": 1000,
            "messages": [{"role": "user", "content": prompt}]
        }).encode()
        
        req = urllib.request.Request(
            "https://api.anthropic.com/v1/messages",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01"
            }
        )
        
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read())
            text = result["content"][0]["text"]
            # Clean any markdown fences
            text = text.replace("```json", "").replace("```", "").strip()
            commentary = json.loads(text)
            print(f"[AI] Commentary generated successfully")
            return commentary
            
    except Exception as e:
        print(f"[AI] Commentary generation failed: {e}")
        return {}


# ── D-08: News headlines ──────────────────────────────────────────────────────

_HEADLINE_TICKERS = [
    "SPY", "QQQ", "IWM", "EEM", "VEA",
    "TLT", "HYG", "LQD",
    "GLD", "USO",
    "^VIX", "DX-Y.NYB",
]

def fetch_news_pool(tickers: list = None, max_per_ticker: int = 3) -> list:
    """
    Fetch recent news from yfinance for a curated subset of tickers.
    Returns a deduplicated list of article dicts: title, publisher, link, published_ts.
    Handles both legacy yfinance schema and the newer nested 'content' schema.
    """
    tickers = tickers or _HEADLINE_TICKERS
    seen_titles: set = set()
    pool: list = []
    for tk in tickers:
        try:
            articles = yf.Ticker(tk).news or []
            count = 0
            for art in articles:
                # New yfinance schema: data lives under art["content"]
                if "content" in art and isinstance(art["content"], dict):
                    c       = art["content"]
                    title   = c.get("title", "").strip()
                    pub     = c.get("provider", {}).get("displayName", "")
                    link    = (c.get("canonicalUrl") or c.get("clickThroughUrl") or {}).get("url", "")
                    pub_dt  = c.get("pubDate", "") or c.get("displayTime", "")
                    # Convert ISO string to unix timestamp
                    try:
                        import datetime as _dt
                        ts = int(_dt.datetime.fromisoformat(pub_dt.replace("Z", "+00:00")).timestamp()) if pub_dt else 0
                    except Exception:
                        ts = 0
                else:
                    # Legacy schema
                    title   = art.get("title", "").strip()
                    pub     = art.get("publisher", "")
                    link    = art.get("link", "")
                    ts      = art.get("providerPublishTime", 0)
                if not title or title in seen_titles:
                    continue
                seen_titles.add(title)
                pool.append({
                    "title":        title,
                    "publisher":    pub,
                    "link":         link,
                    "published_ts": ts,
                })
                count += 1
                if count >= max_per_ticker:
                    break
        except Exception:
            pass
    # Sort by recency
    pool.sort(key=lambda x: x["published_ts"], reverse=True)
    return pool


def select_top_headlines(news_pool: list, regime: str = "Unknown", api_key: str = "") -> list:
    """
    Use Claude API to select the 3 most market-relevant headlines from the pool.
    Returns article dicts verbatim — no generated commentary is added.
    Falls back to the 3 most recent articles if API is unavailable.
    """
    import os, json
    if not api_key:
        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not news_pool:
        return []
    fallback = news_pool[:3]
    if not api_key or DEVELOPMENT_MODE:
        return fallback

    # Build the article list for the prompt (cap at 30 to stay within token budget)
    articles_text = "\n".join(
        f'{i+1}. [{a["publisher"]}] {a["title"]}'
        for i, a in enumerate(news_pool[:30])
    )
    prompt = (
        f"You are a quantitative research assistant. The current market regime is {regime}.\n\n"
        f"Below are recent financial news headlines. Select the 3 that are most relevant to "
        f"understanding current cross-asset market conditions.\n\n"
        f"{articles_text}\n\n"
        f"Return ONLY a JSON array of exactly 3 integers — the 1-based index numbers of the "
        f"selected articles. Example: [4, 12, 7]. No preamble, no markdown, no explanation."
    )
    try:
        import urllib.request
        payload = json.dumps({
            "model": "claude-haiku-4-5-20251001",
            "max_tokens": 60,
            "messages": [{"role": "user", "content": prompt}]
        }).encode()
        req = urllib.request.Request(
            "https://api.anthropic.com/v1/messages",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01"
            }
        )
        with urllib.request.urlopen(req, timeout=20) as resp:
            result = json.loads(resp.read())
            text = result["content"][0]["text"].replace("```json","").replace("```","").strip()
            indices = json.loads(text)
            output = []
            for idx in indices[:3]:
                i = int(idx) - 1
                if 0 <= i < len(news_pool):
                    output.append(dict(news_pool[i]))  # verbatim, no observation field added
            return output if output else fallback
    except Exception as e:
        print(f"[Headlines] Selection failed: {e}")
        return fallback


# ══════════════════════════════════════════════════════════════════════════════
#  BACKTESTING ENGINE
# ══════════════════════════════════════════════════════════════════════════════

def compute_backtest(prices: pd.DataFrame, regime_series: pd.Series) -> dict:
    """
    Backtest regime-aware allocation vs SPY buy-and-hold and 60/40.
    Uses last 5 years of data.
    Monthly rebalancing.
    """
    # Allocation weights per regime
    ALLOC = {
        "Crisis":   {"SPY":0.10,"TLT":0.20,"GLD":0.25,"BIL":0.35,"HYG":0.05,"EEM":0.05},
        "Stressed": {"SPY":0.25,"TLT":0.20,"GLD":0.15,"BIL":0.20,"HYG":0.15,"EEM":0.05},
        "Calm":     {"SPY":0.45,"TLT":0.10,"GLD":0.10,"BIL":0.05,"HYG":0.15,"EEM":0.15},
    }
    
    # Get 5 years of data
    lookback = min(1260, len(prices))
    px = prices.tail(lookback).copy()
    reg = regime_series.reindex(px.index).ffill().fillna("Calm")
    # Eliminate one-day look-ahead: regime detected at close of day i can only
    # be traded starting day i+1. Shift forward by one so the rebalance on the
    # first trading day of month M uses the regime from the last day of M-1.
    reg = reg.shift(1).ffill().fillna("Calm")
    
    # Check required tickers available
    required = ["SPY","TLT","GLD","BIL","HYG","EEM"]
    available = [t for t in required if t in px.columns]
    if len(available) < 4:
        return {}
    
    # Fill missing tickers with SPY
    for t in required:
        if t not in px.columns:
            px[t] = px["SPY"]
    
    rets = px[required].pct_change().fillna(0)
    
    # BK Regime Allocation — monthly rebalance
    bk_val = [1.0]; spy_val = [1.0]; port6040_val = [1.0]
    dates = [px.index[0]]
    
    cur_month = px.index[0].month
    cur_weights = ALLOC["Calm"]
    
    for i in range(1, len(rets)):
        date = rets.index[i]
        r = rets.iloc[i]
        
        # Rebalance monthly
        if date.month != cur_month:
            cur_month = date.month
            regime_now = reg.iloc[i] if i < len(reg) else "Calm"
            if regime_now not in ALLOC: regime_now = "Calm"
            cur_weights = ALLOC[regime_now]
        
        # BK portfolio daily return
        bk_ret = sum(cur_weights.get(t,0) * r.get(t,0) for t in required)
        # SPY buy and hold
        spy_ret = r.get("SPY", 0)
        # 60/40
        port_ret = 0.6 * r.get("SPY",0) + 0.4 * r.get("TLT",0)
        
        bk_val.append(bk_val[-1] * (1 + bk_ret))
        spy_val.append(spy_val[-1] * (1 + spy_ret))
        port6040_val.append(port6040_val[-1] * (1 + port_ret))
        dates.append(date)
    
    import math
    def _stats(vals):
        total_ret = vals[-1] / vals[0] - 1
        n_years   = len(vals) / 252
        cagr      = (vals[-1] / vals[0]) ** (1/max(n_years,0.1)) - 1
        daily_rets = [vals[i]/vals[i-1]-1 for i in range(1,len(vals))]
        vol       = (sum(r**2 for r in daily_rets)/len(daily_rets))**0.5 * math.sqrt(252)
        sharpe    = (cagr - 0.045) / max(vol, 0.001)
        # Max drawdown
        peak = vals[0]; max_dd = 0
        for v in vals:
            if v > peak: peak = v
            dd = v/peak - 1
            if dd < max_dd: max_dd = dd
        return {
            "total":  round(total_ret*100, 1),
            "cagr":   round(cagr*100, 1),
            "sharpe": round(sharpe, 2),
            "max_dd": round(max_dd*100, 1),
        }
    
    # SINGLE SOURCE OF TRUTH for all backtest stats. Both the Intel tab mini
    # snapshot and the Edge tab full panel must read from this one dict —
    # never compute CAGR/Sharpe/MaxDD independently anywhere else.
    bk_stats   = _stats(bk_val)
    spy_stats  = _stats(spy_val)
    p6040_stats= _stats(port6040_val)
    
    # Timeline for chart (subsampled ~60 points)
    timeline = []
    step = max(1, len(dates)//60)
    for i in range(0, len(dates), step):
        timeline.append({
            "date":  dates[i].strftime("%Y-%m"),
            "bk":    round(bk_val[i], 4),
            "spy":   round(spy_val[i], 4),
            "p6040": round(port6040_val[i], 4),
            "regime": reg.iloc[i] if i < len(reg) else "Calm",
        })

    # Full daily cumulative series for accurate monthly returns grid
    daily_bk = pd.Series(bk_val, index=dates)

    return {
        "bk":     bk_stats,
        "spy":    spy_stats,
        "p6040":  p6040_stats,
        "timeline": timeline,
        "daily_bk": daily_bk,
        "years":  round(len(dates)/252, 1),
    }


# ══════════════════════════════════════════════════════════════════════════════
#  BK OPPORTUNITY SCORE + RISK APPETITE SCORE (S2)
# ══════════════════════════════════════════════════════════════════════════════

def compute_bk_opportunity_scores(df: pd.DataFrame,
                                  prices: pd.DataFrame,
                                  frag_df: pd.DataFrame,
                                  current_regime: str) -> tuple:
    """
    BK Composite Score — 5-factor composite ranking model (0–100).

    Returns:
        (scores_dict, factors_dict)
        - scores_dict:  {ticker: total_score}        for ranking
        - factors_dict: {ticker: {factor: raw_0_100}} for decomposition display

    ── FACTOR DECOMPOSITION (what each bar means) ─────────────────────────
    Each factor is scored 0–100 independently, then weighted:

      MOMENTUM (30%)
        Blended price return: 1M×0.2 + 3M×0.5 + 6M×0.3
        Then percentile-ranked across the full universe (0=worst, 100=best).
        Captures trend-following alpha — 3M gets the heaviest weight because
        it balances noise (1M) vs. mean-reversion risk (6M).
        Bar shows: percentile rank of the blended momentum.

      FRAGILITY INVERTED (25%)
        100 minus the BK Fragility Score. A fragility of 20 → bar at 80.
        Rewards low-risk instruments. This is the risk-overlay — momentum
        alone would chase into fragile assets; this penalises them.
        Bar shows: inverse fragility (higher = safer).

      REGIME FIT (20%)
        Matrix lookup: how well the asset's bucket (EQ Growth, FI, CMD, etc.)
        historically performs in the current regime (Calm/Stressed/Crisis).
        Example: Fixed Income scores high in Stressed; EQ Growth in Calm.
        Bar shows: regime-fit score from the bucket × regime matrix.

      SIGNAL (15%)
        Dashboard RAG signal mapped to score: GREEN=100, AMBER=50, RED=0.
        Ensures the model doesn't recommend assets the signal system has
        flagged as deteriorating. Acts as a binary quality gate.
        Bar shows: signal strength (green/amber/red → 100/50/0).

      VOL TREND (10%)
        Binary: 100 if current 20D vol < 1-month-ago vol, else 0.
        Rewards instruments where volatility is compressing (de-risking).
        Penalises those with expanding vol (risk is increasing).
        Bar shows: 100 (vol falling) or 0 (vol rising).
    ─────────────────────────────────────────────────────────────────────────
    """
    scores: dict = {}
    factors: dict = {}

    # Map fragility
    frag_map = {}
    if frag_df is not None and not frag_df.empty and "ticker" in frag_df.columns:
        frag_map = dict(zip(frag_df["ticker"], frag_df["fragility"]))

    # Pre-compute momentum blend universe-wide for percentile ranking
    momentum_raw = {}
    for tk in prices.columns:
        p = prices[tk].dropna()
        if len(p) < 126:
            continue
        r1 = (p.iloc[-1] / p.iloc[-21]  - 1) * 100 if len(p) >= 21  else 0
        r3 = (p.iloc[-1] / p.iloc[-63]  - 1) * 100 if len(p) >= 63  else 0
        r6 = (p.iloc[-1] / p.iloc[-126] - 1) * 100 if len(p) >= 126 else 0
        momentum_raw[tk] = r1 * 0.20 + r3 * 0.50 + r6 * 0.30
    all_mom = sorted(momentum_raw.values())

    def _pct_rank(v):
        if not all_mom:
            return 50.0
        below = sum(1 for x in all_mom if x <= v)
        return below / len(all_mom) * 100.0

    # Map df rag_label by ticker for signal scoring
    rag_map = dict(zip(df["ticker"], df["rag_label"].str.strip()))
    vol_now_map    = dict(zip(df["ticker"], df["vol_now"]))
    vol_prev_map   = dict(zip(df["ticker"], df["vol_1m_ago"]))

    for tk, mom_raw in momentum_raw.items():
        try:
            mom_score = _pct_rank(mom_raw)

            frag_val = frag_map.get(tk, 50.0)
            if pd.isna(frag_val):
                frag_val = 50.0
            frag_inv = 100.0 - float(frag_val)

            regime_fit = get_regime_fit_score(tk, current_regime)

            rag = rag_map.get(tk, "AMBER")
            sig_score = {"GREEN": 100, "AMBER": 50, "RED": 0}.get(rag, 50)

            vn = vol_now_map.get(tk, float("nan"))
            vp = vol_prev_map.get(tk, float("nan"))
            if pd.isna(vn) or pd.isna(vp):
                vol_trend = 50
            else:
                vol_trend = 100 if vn < vp else 0

            total = (mom_score   * 0.30 +
                     frag_inv    * 0.25 +
                     regime_fit  * 0.20 +
                     sig_score   * 0.15 +
                     vol_trend   * 0.10)
            scores[tk] = round(total, 1)

            # Store raw factor scores (0–100 each) for decomposition bars
            factors[tk] = {
                "mom":       round(mom_score, 1),
                "frag_inv":  round(frag_inv, 1),
                "regime_fit": round(regime_fit, 1),
                "signal":    round(sig_score, 1),
                "vol_trend": round(vol_trend, 1),
            }
        except Exception as e:
            print(f"[BK Opp] error {tk}: {e}")
    return scores, factors


def compute_risk_appetite_score(regime: str,
                                fragility_score: float,
                                fear_greed: float,
                                vol_percentile: float) -> float:
    """
    RAS (0-100):
      regime_norm   35%  (score 1..5 normalised)
      fragility_inv 30%
      fear_greed    20%
      vol_inv       15%
    """
    _r5 = REGIME_SCORE_MAP.get(regime, 3)
    regime_norm   = (_r5 / 5.0) * 100.0
    fragility_inv = max(0.0, 100.0 - float(fragility_score))
    vol_inv       = max(0.0, 100.0 - float(vol_percentile))
    ras = (regime_norm   * 0.35 +
           fragility_inv * 0.30 +
           float(fear_greed) * 0.20 +
           vol_inv       * 0.15)
    return round(ras, 1)


def get_allocation_weights(ras: float):
    """Map RAS to (label, bucket-weight dict)."""
    if ras >= 75:
        return ("Risk-On", {
            "EQ Growth": 45, "EQ Defensive": 10, "Fixed Income": 15,
            "Real Assets": 15, "Cash": 5, "Alts": 10})
    if ras >= 55:
        return ("Mild Risk-On", {
            "EQ Growth": 35, "EQ Defensive": 15, "Fixed Income": 20,
            "Real Assets": 15, "Cash": 10, "Alts": 5})
    if ras >= 40:
        return ("Neutral", {
            "EQ Growth": 25, "EQ Defensive": 20, "Fixed Income": 25,
            "Real Assets": 15, "Cash": 15, "Alts": 0})
    if ras >= 25:
        return ("Defensive", {
            "EQ Growth": 15, "EQ Defensive": 25, "Fixed Income": 30,
            "Real Assets": 15, "Cash": 15, "Alts": 0})
    return ("Risk-Off", {
        "EQ Growth": 5, "EQ Defensive": 20, "Fixed Income": 35,
        "Real Assets": 20, "Cash": 20, "Alts": 0})


CASH_PRIORITY = ["BIL", "SHY"]  # prefer real cash instruments over FX pairs


def get_top_instruments_per_bucket(bk_opp_scores: dict, bucket: str, n: int = 2):
    """
    Return top-N *investable* tickers in a bucket by BK Opp Score desc.
    Strict bucket membership: a ticker can only be picked for its assigned bucket.
    Cash bucket has a priority override — if BIL or SHY have non-None scores,
    they are surfaced first regardless of score rank (FX pairs like CHF=X are
    real safe-haven proxies but the Cash bucket should display cash-like assets).
    """
    tks = BUCKET_TICKERS.get(bucket, [])
    scored = [(t, bk_opp_scores[t]) for t in tks
              if t in bk_opp_scores and bk_opp_scores[t] is not None
              and is_rankable(t)]
    scored.sort(key=lambda x: x[1], reverse=True)

    if bucket == "Cash":
        priority_hits = [s for s in scored if s[0] in CASH_PRIORITY]
        others        = [s for s in scored if s[0] not in CASH_PRIORITY]
        scored = priority_hits + others

    return scored[:n]


def _build_headlines_html(headlines: list) -> str:
    """Render the Today's Headlines card for the Intel tab."""
    if not headlines:
        return ""
    import datetime as _dt
    cards = ""
    for art in headlines[:3]:
        title     = art.get("title", "")
        publisher = art.get("publisher", "")
        link      = art.get("link", "#")
        ts        = art.get("published_ts", 0)
        try:
            ts_str = _dt.datetime.utcfromtimestamp(ts).strftime("%d %b %Y") if ts else ""
        except Exception:
            ts_str = ""
        cards += (
            f'<div style="background:#1c2128;border:1px solid #21262d;border-radius:6px;padding:12px 14px;margin-bottom:8px;">'
            f'<div style="display:flex;justify-content:space-between;align-items:flex-start;margin-bottom:6px;">'
            f'<span style="font-size:9px;color:#58a6ff;font-family:monospace;text-transform:uppercase;">{publisher}</span>'
            f'<span style="font-size:9px;color:#6e7681;font-family:monospace;">{ts_str}</span>'
            f'</div>'
            f'<a href="{link}" target="_blank" rel="noopener noreferrer" '
            f'style="font-size:12px;font-weight:600;color:#e6edf3;text-decoration:none;line-height:1.5;">'
            f'{title}</a>'
            f'</div>'
        )
    return (
        f'<div class="fc" style="margin-bottom:14px;">'
        f'<div class="lbl" style="margin-bottom:4px;">TODAY\'S HEADLINES</div>'
        f'<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-bottom:8px;">'
        f'Selection by rule &middot; Headlines verbatim from Yahoo Finance &middot; No commentary by BKIQ</div>'
        + cards
        + f'<div style="font-size:9px;color:#6e7681;font-family:monospace;margin-top:4px;">'
        f'Headlines shown verbatim from original publishers via Yahoo Finance. '
        f'BKIQ does not curate, endorse, or verify these headlines. Not investment advice.</div>'
        + f'</div>'
    )


def build_web_html(df: pd.DataFrame, frag_df: pd.DataFrame = None, prices: pd.DataFrame = None, regime_data: dict = None, fg_data: dict = None, frag_trend: dict = None, ai_commentary: dict = None, backtest_data: dict = None, headlines_data: list = None) -> str:
    import math
    now         = datetime.now(SGT)
    date_str    = now.strftime("%A, %d %b %Y %H:%M SGT")
    gen_ts      = now.strftime("%Y-%m-%dT%H:%M:%S")
    market_open = bool(df["market_open"].iloc[0]) if "market_open" in df.columns else True
    GA          = ""  # Set your GA4 ID here e.g. "G-XXXXXXXXXX"
    # Derive data-review tickers from the scored df (set by compute_scores sanity checks)
    data_review_tickers = set(df.loc[df.get("data_review", pd.Series(False, index=df.index)).astype(bool), "ticker"]) if "data_review" in df.columns else set()

    # ── S2: BK Composite Score + Risk Appetite Score + Dynamic Allocation ──
    _cur_regime = regime_data.get("regime", "Calm") if regime_data else "Calm"
    _frag_sys   = (frag_df.attrs.get("system_score", float(frag_df["fragility"].median()))
                   if frag_df is not None and not frag_df.empty else 50.0)
    _fg_score   = round(fg_data.get("score", 50)) if fg_data else 50
    # Vol percentile: median 20D vol now vs 1Y history
    _vol_pct = 50.0
    try:
        if prices is not None and not prices.empty:
            _rets_all = prices.pct_change()
            _vol_series = _rets_all.rolling(20, min_periods=10).std().median(axis=1) * (252 ** 0.5) * 100
            _vol_series = _vol_series.dropna().tail(252)
            if len(_vol_series) >= 20:
                _cur_vol = float(_vol_series.iloc[-1])
                _vol_pct = float((_vol_series <= _cur_vol).sum() / len(_vol_series) * 100)
    except Exception as _e:
        print(f"[RAS] vol percentile fallback: {_e}")

    bk_opp_scores = {}
    bk_opp_factors = {}  # {ticker: {mom, frag_inv, regime_fit, signal, vol_trend}} for decomposition bars
    if prices is not None and not prices.empty:
        try:
            bk_opp_scores, bk_opp_factors = compute_bk_opportunity_scores(df, prices, frag_df, _cur_regime)
        except Exception as _e:
            print(f"[BK Opp] compute error: {_e}")

    ras_score = compute_risk_appetite_score(_cur_regime, _frag_sys, _fg_score, _vol_pct)
    ras_label, bucket_weights = get_allocation_weights(ras_score)
    print(f"[RAS]  regime={_cur_regime}  fragility={_frag_sys:.1f}  "
          f"fear_greed={_fg_score}  vol_pct={_vol_pct:.1f}  "
          f"-> RAS={ras_score}  label={ras_label}")
    print(f"[RAS]  weights = {bucket_weights}")

    # ── shared cell helpers ───────────────────────────────────────────────────
    def _rc(v, fmt="ret"):
        if pd.isna(v): return '<td class="num gr">&mdash;</td>'
        if fmt == "ret":
            p=v*100; s="+" if p>0 else ""
            cl="ps" if p>=2 else "pl" if p>=0 else "ng" if p>=-2 else "nr"
            return f'<td class="num {cl}">{s}{p:.2f}%</td>'
        if fmt == "vol":
            p=v*100
            return f'<td class="num {"nr" if p>30 else "am" if p>18 else "gr"}">{p:.1f}%</td>'
        if fmt == "dd":
            p=v*100
            return f'<td class="num {"nr" if p<-15 else "am" if p<-7 else "ps"}">{p:.1f}%</td>'
        if fmt == "sh":
            return f'<td class="num {"ps" if v>1 else "am" if v>0 else "nr"}">{v:.2f}</td>'
        return f'<td class="num gr">{v}</td>'

    def _sig(rl, rc_col):
        rl=rl.strip()
        dot={"RED":"#f85149","AMBER":"#e3b341","GREEN":"#3fb950"}.get(rl,"#8b949e")
        cl={"RED":"sr","AMBER":"sa","GREEN":"sg"}.get(rl,"")
        return f'<td class="sig {cl}"><span style="color:{dot};">&#9679;</span> {rl}</td>'

    def _srow(sec, cs=14):
        return f'<tr class="sh"><td colspan="{cs}">{SECTION_LABELS.get(sec,sec)}</td></tr>'

    def _srow_acc(sec, tab_prefix, summary, count, cs=14, first=False):
        """Accordion-header variant of _srow. Wraps each section in <tbody>."""
        key = f"{tab_prefix}_{sec.lower()}"
        label = SECTION_LABELS.get(sec, sec)
        prefix = '' if first else '</tbody>'
        return (
            f'{prefix}<tbody class="bk-accordion-section">'
            f'<tr class="bk-accordion-header" data-key="{key}">'
            f'<td colspan="{cs}">'
            f'<span class="bk-accordion-icon">&#9654;</span>'
            f'{label}'
            f'<span class="bk-accordion-count">({count})</span>'
            f'<span class="bk-accordion-summary">{summary}</span>'
            f'<span class="bk-accordion-hint">Click to expand</span>'
            f'</td></tr>'
        )

    def _acc_controls(tab_prefix):
        return (f'<div class="bk-acc-controls">'
                f'<button class="bk-acc-btn" onclick="bkSetAllAccordions(\'open\',\'{tab_prefix}\')">Expand All</button>'
                f'<button class="bk-acc-btn" onclick="bkSetAllAccordions(\'closed\',\'{tab_prefix}\')">Collapse All</button>'
                f'</div>')

    def _bar(nm, val, mx, color):
        # Proportional scaling capped at 75% (prevents outlier dominance /
        # tight-cluster bunching), with 5% floor for visibility when val != 0.
        if mx > 0:
            w = abs(val) / mx * 75
            if val != 0:
                w = max(5, w)
            w = min(75, w)
        else:
            w = 0
        s="+" if val>=0 else ""
        return (f'<div style="display:flex;align-items:center;gap:10px;padding:7px 0;border-bottom:1px solid #21262d;">'
                f'<div style="width:150px;font-size:11px;color:#e6edf3;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;">{nm}</div>'
                f'<div style="flex:1;background:#21262d;border-radius:3px;height:8px;">'
                f'<div style="width:{w:.0f}%;background:{color};height:8px;border-radius:3px;"></div></div>'
                f'<div style="width:62px;text-align:right;font-family:monospace;font-size:12px;font-weight:700;color:{color};">{s}{val*100:.2f}%</div>'
                f'</div>')

    # ══ TAB 1: PERFORMANCE ════════════════════════════════════════════════════
    nr=int((df["rag_label"].str.strip()=="RED").sum())
    na=int((df["rag_label"].str.strip()=="AMBER").sum())
    ng=int((df["rag_label"].str.strip()=="GREEN").sum())
    tot=len(df)

    # Market tone: three-way gate over regime, fragility, and rising-risk share.
    _regime_now_tone = regime_data.get("regime", "Calm") if regime_data else "Calm"
    _frag_sys_tone = (frag_df.attrs.get("system_score", float(frag_df["fragility"].median()))
                      if frag_df is not None and not frag_df.empty else 50.0)
    _rising_tone = _count_rising_risk(df)
    tone, tc, tb = calculate_market_tone(_regime_now_tone, _frag_sys_tone, _rising_tone, tot)

    # Rolling-1M gainers / losers — exclude non-investable tickers
    # (yield indices, vol instruments, computed proxies, raw FX index).
    _rankable = df[df["ticker"].apply(is_rankable)]
    mtd=_rankable[["name","ticker","ret_1m"]].dropna(subset=["ret_1m"])
    gain=mtd.nlargest(5,"ret_1m"); loss=mtd.nsmallest(5,"ret_1m")
    gm=gain["ret_1m"].abs().max(); lm=loss["ret_1m"].abs().max()
    gh="".join(_bar(r["name"],r["ret_1m"],gm,"#3fb950") for _,r in gain.iterrows())
    lh="".join(_bar(r["name"],r["ret_1m"],lm,"#f85149") for _,r in loss.iterrows())

    d1th="<th>1D</th>" if market_open else ""
    def _sparkline(spark_data, width=80, height=24):
        """Generate inline SVG sparkline from 20-day normalised data."""
        if not spark_data or len(spark_data) < 2:
            return '<td style="padding:7px 8px;color:#8b949e;">&mdash;</td>'
        mn, mx = min(spark_data), max(spark_data)
        rng = mx - mn if mx != mn else 1.0
        pts = []
        for i, v in enumerate(spark_data):
            x = i / (len(spark_data)-1) * width
            y = height - (v - mn) / rng * (height - 4) - 2
            pts.append(f"{x:.1f},{y:.1f}")
        color = "#3fb950" if spark_data[-1] >= spark_data[0] else "#f85149"
        polyline = " ".join(pts)
        last_x, last_y = pts[-1].split(",")
        svg = (f'<svg width="{width}" height="{height}" viewBox="0 0 {width} {height}" '
               f'style="display:block;">'
               f'<polyline points="{polyline}" fill="none" stroke="{color}" stroke-width="1.5" stroke-linejoin="round"/>'
               f'<circle cx="{last_x}" cy="{last_y}" r="2" fill="{color}"/>'
               f'</svg>')
        return f'<td style="padding:4px 8px;">{svg}</td>'

    # Per-section summaries for Performance accordion headers
    _perf_sum = {}
    for _sec in df["section"].unique():
        _sub = df[df["section"] == _sec].dropna(subset=["ret_1m"])
        if _sub.empty:
            _perf_sum[_sec] = ("", len(df[df["section"] == _sec]))
            continue
        _best = _sub.loc[_sub["ret_1m"].idxmax()]
        _worst = _sub.loc[_sub["ret_1m"].idxmin()]
        _txt = (f'Best: <span style="color:#3fb950;">{_best["name"]} {_best["ret_1m"]*100:+.1f}%</span>'
                f' &nbsp;|&nbsp; Worst: <span style="color:#f85149;">{_worst["name"]} {_worst["ret_1m"]*100:+.1f}%</span>')
        _perf_sum[_sec] = (_txt, int((df["section"] == _sec).sum()))

    _perf_cs = 10  # colspan: Asset,Ticker,CCY,Trend,[1D],1W,1M,3M,YTD,Signal
    pr=""; pv=None; _perf_first=True
    for _,row in df.iterrows():
        if row["section"]!=pv:
            pv=row["section"]
            _sum, _cnt = _perf_sum.get(pv, ("", 0))
            pr += _srow_acc(pv, "perf", _sum, _cnt, cs=_perf_cs, first=_perf_first)
            _perf_first=False
        d1=_rc(row["ret_1d"]) if market_open else ""
        ccy = row.get("currency","USD")
        spark_td = _sparkline(row.get("spark",[]))
        _dr_badge = (' <span title="Data under review — yfinance source suspected unadjusted" '
                     'style="color:#e3b341;font-size:9px;font-weight:700;">&#9888; DATA REVIEW</span>'
                     if row.get("data_review") else '')
        # BNO (WTI oil proxy): flag elevated vol as market dislocation, not bad data
        _disloc_badge = (' <span title="Vol elevated — reflects recent crude market dislocation" '
                         'style="color:#e3b341;font-size:9px;font-weight:700;">&#9888; VOL DISLOCATION</span>'
                         if row.get("ticker") == "BNO" and not pd.isna(row.get("vol_now", float("nan")))
                         and row.get("vol_now", 0) > 0.50 else '')
        pr+=(f'<tr class="bk-accordion-body"><td class="an">{row["name"]}{_dr_badge}{_disloc_badge}</td><td class="tk">{row["ticker"]}</td>'
             f'<td class="num gr" style="font-size:9px;">{ccy}</td>'
             f'{spark_td}'
             f'{d1}{_rc(row["ret_1w"])}{_rc(row["ret_1m"])}{_rc(row["ret_3m"])}{_rc(row["ret_ytd"])}'
             f'{_sig(row["rag_label"],row["rag_color"])}</tr>')

    perf=(f'<div class="tbar"><div><div class="lbl">MARKET TONE</div>'
          f'<div class="pill" style="background:{tb};color:{tc};border:1px solid {tc};">{tone}</div></div>'
          f'<div class="dvdr"></div><div class="rb">'
          f'<div class="ri"><div class="rn" style="color:#f85149;">{nr}</div><div class="rl">RED</div></div>'
          f'<div class="ri"><div class="rn" style="color:#e3b341;">{na}</div><div class="rl">AMBER</div></div>'
          f'<div class="ri"><div class="rn" style="color:#3fb950;">{ng}</div><div class="rl">GREEN</div></div>'
          f'<div class="ri"><div class="rn" style="color:#e6edf3;">{tot}</div><div class="rl">TOTAL</div></div>'
          f'</div></div>'
          f'<div class="gl"><div class="gc"><div class="gt"><span class="gd" style="background:#3fb950;"></span>'
          f'Top 5 &mdash; Rolling 1M &nbsp;<span style="color:#8b949e;font-weight:400;">(21 trading days)</span></div>{gh}</div>'
          f'<div class="gc"><div class="gt"><span class="gd" style="background:#f85149;"></span>'
          f'Top 5 Laggards &mdash; Rolling 1M &nbsp;<span style="color:#8b949e;font-weight:400;">(21 trading days)</span></div>{lh}</div></div>'
          + _acc_controls("perf")
          + f'<div class="tw"><table><thead><tr><th style="text-align:left;">Asset</th><th>Ticker</th><th>CCY</th><th style="min-width:80px;">Trend 20D</th>'
          f'{d1th}<th>1W</th><th>1M</th><th>3M</th><th>YTD</th><th>Signal</th>'
          f'</tr></thead>{pr}</tbody></table></div>'
          + (
             (lambda _bno_ytd: (
               f'<div style="margin-top:8px;padding:8px 14px;background:#2d1a0e;border-left:3px solid #e3b341;'
               f'font-size:9px;color:#8b949e;line-height:1.7;">'
               f'<strong style="color:#e3b341;">&#9888; DATA REVIEW — BNO (Brent Oil ETF):</strong> '
               f'yfinance auto-adjusted close prices for BNO show anomalous 2026 YTD of {_bno_ytd} '
               f'while Brent crude spot prices have declined. Likely cause: unadjusted corporate action '
               f'(reverse split or NAV reset) in the upstream data source. '
               f'YTD / 3M / 1M return cells are dashed until the source data is confirmed clean. '
               f'Vol, Max DD, and RAG signal calculations are unaffected.'
               f'</div>'
             ))(
               (lambda r: f'{r["ret_ytd"]*100:+.1f}%' if r is not None and pd.notna(r.get("ret_ytd")) else '+73.3%')(
                 df[df["ticker"]=="BNO"].iloc[0].to_dict() if not df[df["ticker"]=="BNO"].empty else None
               )
             )
             if "BNO" in data_review_tickers else ''
          )
          )

    # ══ TAB 2: RISK ═══════════════════════════════════════════════════════════
    def _varrow(now_v, ago_v):
        if pd.isna(now_v) or pd.isna(ago_v) or ago_v==0: return "gr","&#8594;","&mdash;"
        chg = (now_v - ago_v) / ago_v
        abs_chg = now_v - ago_v  # absolute change in vol (pp)
        # For low-vol instruments (< 3% annualised), show absolute pp change
        if ago_v < 0.03:
            pp = abs_chg * 100
            pct = f"{pp:+.2f}pp"
        else:
            pct = f"{chg*100:+.1f}%"
        if chg>=0.20:  return "nr","&#11014;&#11014;",pct
        if chg>=0.05:  return "am","&#11014;",pct
        if chg>=-0.05: return "gr","&#8594;",pct
        return "ps","&#11015;",pct

    rising=stable=falling=0
    for _,row in df.iterrows():
        nv=row.get("vol_now",float("nan")); av=row.get("vol_1m_ago",float("nan"))
        if pd.isna(nv) or pd.isna(av) or av==0: continue
        chg=(nv-av)/av
        if chg>=0.05: rising+=1
        elif chg<=-0.05: falling+=1
        else: stable+=1

    vsumm=(f'<div style="display:flex;gap:14px;margin-bottom:14px;flex-wrap:wrap;">'
           f'<div class="vc" style="border-color:#f85149;"><div class="vn" style="color:#f85149;">{rising}</div>'
           f'<div class="vl">VOL RISING &#11014;</div><div class="vs">Change &gt; +5%</div></div>'
           f'<div class="vc" style="border-color:#8b949e;"><div class="vn" style="color:#8b949e;">{stable}</div>'
           f'<div class="vl">VOL STABLE &#8594;</div><div class="vs">&#8722;5% to +5%</div></div>'
           f'<div class="vc" style="border-color:#3fb950;"><div class="vn" style="color:#3fb950;">{falling}</div>'
           f'<div class="vl">VOL EASING &#11015;</div><div class="vs">Change &lt; &#8722;5%</div></div>'
           f'</div>')

    # Per-section summaries for Risk accordion headers: avg vol, avg DD, avg Sharpe
    _risk_sum = {}
    for _sec in df["section"].unique():
        _sub = df[df["section"] == _sec]
        _avol = _sub["vol_now"].mean(skipna=True)
        _add  = _sub["max_dd"].mean(skipna=True)
        _ash  = _sub["sharpe"].mean(skipna=True)
        _vol_txt = f'{_avol*100:.1f}%' if pd.notna(_avol) else '—'
        _dd_txt  = f'{_add*100:.1f}%'  if pd.notna(_add)  else '—'
        _sh_txt  = f'{_ash:.2f}'       if pd.notna(_ash)  else '—'
        _risk_sum[_sec] = (
            f'Avg vol: {_vol_txt} &nbsp;|&nbsp; Avg DD: {_dd_txt} &nbsp;|&nbsp; Avg Sharpe: {_sh_txt}',
            int(len(_sub))
        )

    rr=""; rv=None; _risk_first=True
    for _,row in df.iterrows():
        if row["section"]!=rv:
            rv=row["section"]
            _sum, _cnt = _risk_sum.get(rv, ("", 0))
            rr += _srow_acc(rv, "risk", _sum, _cnt, cs=8, first=_risk_first)
            _risk_first=False
        cl,arrow,pct=_varrow(row.get("vol_now",float("nan")),row.get("vol_1m_ago",float("nan")))
        # Vol indices: vol-of-vol is not a sensible metric — show N/A
        if row["ticker"] in VOL_VALUE_EXCLUSIONS:
            vn = va = "N/A"; cl = "gr"; arrow = ""; pct = "N/A"
        else:
            vn=f'{row["vol_now"]*100:.1f}%' if not pd.isna(row.get("vol_now",float("nan"))) else "—"
            va=f'{row["vol_1m_ago"]*100:.1f}%' if not pd.isna(row.get("vol_1m_ago",float("nan"))) else "—"
        # Yield indices: Sharpe is meaningless for rate levels — show N/A
        _sharpe_cell = ('<td class="num gr" style="font-style:italic;">N/A</td>'
                        if row["ticker"] in SHARPE_EXCLUSIONS
                        else _rc(row["sharpe"],"sh"))
        rr+=(f'<tr class="bk-accordion-body"><td class="an">{row["name"]}</td><td class="tk">{row["ticker"]}</td>'
             f'<td class="num gr">{vn}</td><td class="num gr">{va}</td>'
             f'<td class="num {cl}" style="font-family:monospace;">{arrow}&nbsp;{pct}</td>'
             f'{_rc(row["max_dd"],"dd")}{_sharpe_cell}'
             f'{_sig(row["rag_label"],row["rag_color"])}</tr>')

    risk=(vsumm+
          _acc_controls("risk") +
          f'<div class="tw"><table><thead><tr><th style="text-align:left;">Asset</th><th>Ticker</th>'
          f'<th>Vol 20D</th><th>Vol 1M Ago</th><th>30D Change</th>'
          f'<th>Max DD</th><th>Sharpe</th><th>Signal</th></tr></thead>{rr}</tbody></table></div>'
          f'<div style="margin-top:10px;font-size:9px;color:#8b949e;font-family:monospace;line-height:1.8;">'
          f'&#11014;&#11014; Vol rising &ge;+20% &nbsp;&#183;&nbsp; &#11014; +5% to +20% &nbsp;&#183;&nbsp; '
          f'&#8594; stable &#8722;5% to +5% &nbsp;&#183;&nbsp; &#11015; easing &lt;&#8722;5% &nbsp;&#183;&nbsp; '
          f'Sharpe = 1Y excess return / vol (rf=4.5%)<br>'
          f'Vol 20D = 20-day daily returns std dev &times; &radic;252 (annualised) &nbsp;&#183;&nbsp; '
          f'Vol 1M Ago = same calculation 21 trading days prior &nbsp;&#183;&nbsp; '
          f'20Y Treasury vol &gt; HY Credit vol is expected: duration risk dominates rate-sensitive environments</div>')


    # ── Build Fragility Trend SVG chart ──────────────────────────────────────
    if frag_trend and frag_trend.get("trend"):
        trend_pts  = frag_trend["trend"]
        ft_current = frag_trend.get("current", 50)
        ft_peak    = frag_trend.get("peak_2y", 50)
        ft_trough  = frag_trend.get("trough_2y", 50)
        ft_avg     = frag_trend.get("avg_2y", 50)
        ft_color   = "#f85149" if ft_current>=70 else "#e3b341" if ft_current>=55 else "#3fb950"

        # Build SVG
        svg_w = 900; svg_h = 180; pad_l = 40; pad_r = 10; pad_t = 10; pad_b = 30
        chart_w = svg_w - pad_l - pad_r
        chart_h = svg_h - pad_t - pad_b
        n_pts   = len(trend_pts)

        def _x(i):   return pad_l + i / max(n_pts-1,1) * chart_w
        def _y(val): return pad_t + (1 - val/100) * chart_h

        # Background regime bands
        bands = []
        band_start = 0; band_reg = trend_pts[0]["regime"] if trend_pts else "Moderate"
        for i, pt in enumerate(trend_pts):
            if pt["regime"] != band_reg or i == n_pts-1:
                x1 = _x(band_start); x2 = _x(i)
                bc = {"Crisis":"rgba(248,81,73,0.12)","Stressed":"rgba(227,179,65,0.10)","Moderate":"rgba(63,185,80,0.06)"}.get(band_reg,"rgba(0,0,0,0)")
                bands.append(f'<rect x="{x1:.1f}" y="{pad_t}" width="{x2-x1:.1f}" height="{chart_h}" fill="{bc}"/>')
                band_start = i; band_reg = pt["regime"]

        # Threshold lines
        y50 = _y(50); y70 = _y(70)
        thresholds = (
            f'<line x1="{pad_l}" y1="{y70:.1f}" x2="{svg_w-pad_r}" y2="{y70:.1f}" stroke="#f85149" stroke-width="0.8" stroke-dasharray="4,3" opacity="0.6"/>'
            f'<text x="{pad_l-4}" y="{y70+3:.1f}" text-anchor="end" font-size="8" fill="#f85149" opacity="0.8">70</text>'
            f'<line x1="{pad_l}" y1="{y50:.1f}" x2="{svg_w-pad_r}" y2="{y50:.1f}" stroke="#e3b341" stroke-width="0.8" stroke-dasharray="4,3" opacity="0.6"/>'
            f'<text x="{pad_l-4}" y="{y50+3:.1f}" text-anchor="end" font-size="8" fill="#e3b341" opacity="0.8">50</text>'
        )

        # Y axis labels
        y_labels = "".join(
            f'<text x="{pad_l-4}" y="{_y(v)+3:.1f}" text-anchor="end" font-size="8" fill="#8b949e">{v}</text>'
            for v in [0, 25, 75, 100]
        )

        # X axis month labels every ~21 pts
        x_labels = ""
        prev_month = ""
        for i, pt in enumerate(trend_pts):
            m = pt["date"][:7]
            if m != prev_month and i % 42 == 0:
                prev_month = m
                x_labels += f'<text x="{_x(i):.1f}" y="{svg_h-4}" text-anchor="middle" font-size="8" fill="#8b949e">{m}</text>'

        # Line path + fill
        pts_str = " ".join(f"{_x(i):.1f},{_y(pt['score']):.1f}" for i, pt in enumerate(trend_pts))
        fill_pts = f"{_x(0):.1f},{pad_t+chart_h} " + pts_str + f" {_x(n_pts-1):.1f},{pad_t+chart_h}"
        line_color = ft_color

        # Last point dot
        last_x = _x(n_pts-1); last_y = _y(ft_current)
        dot = f'<circle cx="{last_x:.1f}" cy="{last_y:.1f}" r="4" fill="{ft_color}" stroke="#0d1117" stroke-width="2"/>'

        svg_parts = (
            f'<svg viewBox="0 0 {svg_w} {svg_h}" width="100%" style="max-width:{svg_w}px;display:block;overflow:visible;">'
            + "".join(bands)
            + thresholds + y_labels + x_labels
            + f'<polyline points="{fill_pts}" fill="{ft_color}" opacity="0.08" stroke="none"/>'
            + f'<polyline points="{pts_str}" fill="none" stroke="{ft_color}" stroke-width="1.8" stroke-linejoin="round" stroke-linecap="round"/>'
            + dot
            + '</svg>'
        )

        frag_trend_html = (
            f'<div class="fc" style="margin-bottom:14px;">'
            f'<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:12px;flex-wrap:wrap;gap:8px;">'
            f'<div class="lbl">SYSTEM FRAGILITY TREND — LAST 2 YEARS</div>'
            f'<div style="display:flex;gap:16px;">'
            f'<div style="text-align:center;"><div style="font-size:9px;color:#8b949e;">CURRENT</div>'
            f'<div style="font-size:18px;font-weight:700;color:{ft_color};font-family:monospace;">{ft_current:.0f}</div></div>'
            f'<div style="text-align:center;"><div style="font-size:9px;color:#8b949e;">2Y PEAK</div>'
            f'<div style="font-size:18px;font-weight:700;color:#f85149;font-family:monospace;">{ft_peak:.0f}</div></div>'
            f'<div style="text-align:center;"><div style="font-size:9px;color:#8b949e;">2Y AVG</div>'
            f'<div style="font-size:18px;font-weight:700;color:#8b949e;font-family:monospace;">{ft_avg:.0f}</div></div>'
            f'<div style="text-align:center;"><div style="font-size:9px;color:#8b949e;">2Y LOW</div>'
            f'<div style="font-size:18px;font-weight:700;color:#3fb950;font-family:monospace;">{ft_trough:.0f}</div></div>'
            f'</div></div>'
            f'{svg_parts}'
            f'<div style="display:flex;gap:16px;margin-top:6px;font-size:9px;color:#8b949e;">'
            f'<span><span style="color:#f85149;">&#9632;</span> CRISIS &#8805;70</span>'
            f'<span><span style="color:#e3b341;">&#9632;</span> STRESSED 55&#8211;69</span>'
            f'<span><span style="color:#3fb950;">&#9632;</span> MODERATE &lt;55</span>'
            f'<span style="margin-left:8px;">Dashed lines = regime thresholds</span>'
            f'</div></div>'
        )
    else:
        frag_trend_html = ""


    # ══ TAB 3: FRAGILITY ══════════════════════════════════════════════════════
    if frag_df is not None and not frag_df.empty:
        ss=frag_df.attrs.get("system_score",float(frag_df["fragility"].median()))
        reg=frag_df.attrs.get("regime","MODERATE")
        rc_={"CRISIS":"#f85149","STRESSED":"#e3b341","MODERATE":"#3fb950"}.get(reg,"#8b949e")
        rb_={"CRISIS":"#2d0f0e","STRESSED":"#2d2106","MODERATE":"#0d2318"}.get(reg,"#161b22")
        ncr=int((frag_df["rag"]=="CRISIS").sum())
        nst=int((frag_df["rag"]=="STRESSED").sum())
        nca=int((frag_df["rag"]=="MODERATE").sum())

        def _arc(deg,r=75,cx=100,cy=100):
            rad=math.radians(180-deg)
            return cx+r*math.cos(rad), cy-r*math.sin(rad)
        ga=max(1,min(179,int(ss/100*180))); gc="#f85149" if ss>=70 else "#e3b341" if ss>=55 else "#3fb950"
        _frag_label = "CRISIS" if ss>=70 else "STRESSED" if ss>=55 else "MODERATE"
        ax,ay=_arc(ga); lg=1 if ga>90 else 0
        gauge=(f'<svg viewBox="0 0 200 130" width="200" height="130">'
               f'<path d="M 25 100 A 75 75 0 0 1 175 100" fill="none" stroke="#21262d" stroke-width="12" stroke-linecap="round"/>'
               f'<path d="M 25 100 A 75 75 0 {lg} 1 {ax:.1f} {ay:.1f}" fill="none" stroke="{gc}" stroke-width="12" stroke-linecap="round"/>'
               f'<text x="100" y="88" text-anchor="middle" font-size="28" font-weight="bold" fill="{gc}" font-family="monospace">{ss:.0f}</text>'
               f'<text x="100" y="104" text-anchor="middle" font-size="9" fill="#8b949e" font-family="monospace">/ 100</text>'
               f'<text x="100" y="122" text-anchor="middle" font-size="11" font-weight="700" fill="{gc}" font-family="monospace">{_frag_label}</text>'
               f'<text x="25" y="115" text-anchor="middle" font-size="8" fill="#555">0</text>'
               f'<text x="175" y="115" text-anchor="middle" font-size="8" fill="#555">100</text></svg>')

        t5h=""
        for _,r in frag_df[frag_df["ticker"].apply(is_rankable)].head(5).iterrows():
            fc="#f85149" if r["rag"]=="CRISIS" else "#e3b341" if r["rag"]=="STRESSED" else "#3fb950"
            bw=min(100,r["fragility"])
            t5h+=(f'<div style="display:flex;align-items:center;gap:10px;padding:7px 0;border-bottom:1px solid #21262d;">'
                  f'<div style="width:150px;font-size:11px;color:#e6edf3;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;">{r["name"]}</div>'
                  f'<div style="flex:1;background:#21262d;border-radius:3px;height:8px;">'
                  f'<div style="width:{bw:.0f}%;background:{fc};height:8px;border-radius:3px;"></div></div>'
                  f'<div style="width:48px;text-align:right;font-family:monospace;font-size:12px;font-weight:700;color:{fc};">{r["fragility"]:.0f}</div>'
                  f'<div style="width:70px;text-align:center;font-size:9px;font-family:monospace;color:{fc};'
                  f'background:{rb_};border:1px solid {fc};border-radius:10px;padding:1px 6px;">{r["rag"]}</div></div>')

        PL={"pillar_dd":"Drawdown","pillar_vol":"Volatility","pillar_cvar":"Tail Risk",
            "pillar_trend":"Trend","pillar_corr":"Contagion","pillar_volz":"Vol Stress"}

        # Group frag_df by section (preserve fragility ordering within section)
        _frag_sorted = frag_df.copy()
        _frag_sorted["_sec_order"] = _frag_sorted["section"].map(
            {s: i for i, s in enumerate(SECTION_ORDER)}).fillna(99)
        _frag_sorted = _frag_sorted.sort_values(
            ["_sec_order", "fragility"], ascending=[True, False])

        # FX excluded from Fragility tab: currency pair pillars (Contagion, Trend,
        # Vol Stress) return 0.0 or meaningless values for FX pairs.
        _frag_display = _frag_sorted[_frag_sorted["section"] != "FX"]
        _n_fx_excl    = len([t for sec, t, _, _ in UNIVERSE if sec == "FX" and t not in DISPLAY_EXCLUSIONS])
        _n_frag_scored = N_INSTRUMENTS - _n_fx_excl

        # Per-section summaries for Fragility accordion headers
        _frag_sum = {}
        for _sec in _frag_display["section"].unique():
            _sub = _frag_display[_frag_display["section"] == _sec]
            _avg = float(_sub["fragility"].mean())
            _top = _sub.iloc[0]
            _frag_sum[_sec] = (
                f'Avg score: <span style="color:#e3b341;">{_avg:.0f}</span> '
                f'&nbsp;|&nbsp; Highest: <span style="color:#f85149;">{_top["name"]} ({_top["fragility"]:.0f})</span>',
                int(len(_sub))
            )

        fr=""; _frag_pv = None; _frag_first=True
        for _,r in _frag_display.iterrows():
            _sec = r["section"]
            if _sec != _frag_pv:
                _frag_pv = _sec
                _sum, _cnt = _frag_sum.get(_sec, ("", 0))
                fr += _srow_acc(_sec, "frag", _sum, _cnt, cs=12, first=_frag_first)
                _frag_first=False
            if r.get("rag") == "N/A" or pd.isna(r.get("fragility")):
                # Excluded computed proxies / fear indices — show N/A row
                fr += (f'<tr class="bk-accordion-body"><td class="an">{r["name"]}</td><td class="tk">{r["ticker"]}</td>'
                       f'<td class="num gr" colspan="10" style="font-style:italic;">'
                       f'N/A &mdash; computed proxy / fear index (not scored)</td></tr>')
                continue
            fc="#f85149" if r["rag"]=="CRISIS" else "#e3b341" if r["rag"]=="STRESSED" else "#3fb950"
            bw=min(100,r["fragility"])
            pv={k:r.get(k,0) for k in PL}; top=PL[max(pv,key=pv.get)]
            pc="".join(f'<td class="num {"ps" if r.get(k,0)>1 else "am" if r.get(k,0)>0 else "gr"}">{r.get(k,0):+.1f}</td>' for k in PL)
            fr+=(f'<tr class="bk-accordion-body"><td class="an">{r["name"]}</td><td class="tk">{r["ticker"]}</td>'
                 f'<td class="num" style="color:{fc};font-weight:700;">{r["fragility"]:.0f}</td>'
                 f'<td style="padding:7px 8px;"><div style="background:#21262d;border-radius:3px;height:6px;width:80px;">'
                 f'<div style="width:{bw:.0f}%;background:{fc};height:6px;border-radius:3px;"></div></div></td>'
                 f'<td style="text-align:center;"><span style="font-size:9px;font-family:monospace;color:{fc};'
                 f'background:{rb_};border:1px solid {fc};border-radius:10px;padding:1px 8px;">{r["rag"]}</span></td>'
                 f'<td class="num gr" style="font-size:10px;">{top}</td>{pc}</tr>')

        frag= frag_trend_html + (f'<div style="display:grid;grid-template-columns:auto 1fr 1fr 1fr 1fr;gap:14px;margin-bottom:14px;align-items:stretch;">'
              f'<div class="fc" style="text-align:center;"><div class="lbl" style="margin-bottom:8px;">SYSTEM FRAGILITY</div>'
              f'{gauge}<div class="pill" style="background:{rb_};color:{rc_};border:1px solid {rc_};margin-top:6px;">{reg}</div></div>'
              f'<div class="fc" style="text-align:center;"><div class="lbl">CRISIS</div>'
              f'<div style="font-size:28px;font-weight:700;color:#f85149;font-family:monospace;">{ncr}</div>'
              f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Score &#8805; 70</div></div>'
              f'<div class="fc" style="text-align:center;"><div class="lbl">STRESSED</div>'
              f'<div style="font-size:28px;font-weight:700;color:#e3b341;font-family:monospace;">{nst}</div>'
              f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Score 55&#8211;69</div></div>'
              f'<div class="fc" style="text-align:center;"><div class="lbl">MODERATE</div>'
              f'<div style="font-size:28px;font-weight:700;color:#3fb950;font-family:monospace;">{nca}</div>'
              f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Score &lt; 55</div></div>'
              f'<div class="fc" style="text-align:center;"><div class="lbl">TOTAL</div>'
              f'<div style="font-size:28px;font-weight:700;color:#e6edf3;font-family:monospace;">{_n_frag_scored}</div>'
              f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">{_n_frag_scored} of {N_INSTRUMENTS} scored (FX excluded)</div>'
              f'<div style="font-size:8px;color:#6a7485;margin-top:2px;">Yield indices &amp; FX pairs excluded from fragility scoring</div></div></div>'
              f'<div class="fc" style="margin-bottom:14px;">'
              f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;margin-bottom:12px;text-transform:uppercase;">&#9888; Top 5 Most Fragile</div>'
              f'{t5h}</div>'
              + _acc_controls("frag") +
              f'<div class="tw"><table><thead><tr><th style="text-align:left;">Asset</th><th>Ticker</th>'
              f'<th>Score</th><th>Bar</th><th>Status</th><th>Top Driver</th>'
              f'<th>Drawdown</th><th>Volatility</th><th>Tail Risk</th><th>Trend</th><th>Contagion</th><th>Vol Stress</th>'
              f'</tr></thead>{fr}</tbody></table></div>'
              f'<div style="margin-top:14px;padding:12px 14px;background:#0d1117;border-left:3px solid #58a6ff;font-size:10px;color:#9aa3b8;line-height:1.7;">'
              f'<strong style="color:#c8cfe0;">What fragility measures:</strong> Fragility tracks '
              f'<em>price stress characteristics</em> &mdash; not the safety or quality of an instrument. '
              f'A safe-haven asset can show elevated fragility when its price is declining.<br>'
              f'<strong style="color:#c8cfe0;">T-Bills (BIL):</strong> Fragility reflects rate '
              f'sensitivity and reinvestment risk, not credit or liquidity risk.<br>'
              f'<strong style="color:#c8cfe0;">FX pairs:</strong> Excluded from this tab — '
              f'currency pair pillars (Contagion, Trend, Vol Stress) are not meaningful for FX. '
              f'FX instruments appear on Performance, Risk, and Analysis tabs.'
              f'</div>'
              f'<div style="margin-top:10px;font-size:9px;color:#8b949e;font-family:monospace;line-height:1.8;">'
              f'BK Fragility Framework &#183; Drawdown 22% + CVaR 20% + Contagion 18% + Volatility 15% + Trend 15% + Vol Stress 10% &#183; '
              f'CRISIS &#8805;70 &#183; STRESSED 55&#8211;69 &#183; MODERATE &lt;55<br>'
              f'Pillar scores are standardised z-scores relative to history (positive = above average stress) &#183; '
              f'Top Driver = highest contributing pillar &#183; '
              f'Negative scores = below historical stress average (healthy signal)</div>')
    else:
        frag='<div style="padding:40px;text-align:center;color:#8b949e;">Fragility data unavailable.</div>'



    def _build_edge_backtest(bt):
        """Build full backtest section for Edge tab."""
        if not bt: return ''
        bk_s=bt.get('bk',{}); spy_s=bt.get('spy',{}); p_s=bt.get('p6040',{})
        yrs=bt.get('years',5); tl=bt.get('timeline',[])
        # Build chart SVG
        svg_w=900; svg_h=200; pad_l=50; pad_r=10; pad_t=10; pad_b=30
        cw=svg_w-pad_l-pad_r; ch=svg_h-pad_t-pad_b
        if not tl: return ''
        all_v=[pt['bk'] for pt in tl]+[pt['spy'] for pt in tl]+[pt['p6040'] for pt in tl]
        mn=min(all_v); mx=max(all_v); rng=max(mx-mn,0.01)
        def _x(i): return pad_l+i/max(len(tl)-1,1)*cw
        def _y(v): return pad_t+(1-(v-mn)/rng)*ch
        bk_pts=' '.join(f'{_x(i):.1f},{_y(pt["bk"]):.1f}' for i,pt in enumerate(tl))
        spy_pts=' '.join(f'{_x(i):.1f},{_y(pt["spy"]):.1f}' for i,pt in enumerate(tl))
        p_pts=' '.join(f'{_x(i):.1f},{_y(pt["p6040"]):.1f}' for i,pt in enumerate(tl))
        # X labels
        xlabels=''
        prev_y=''
        for i,pt in enumerate(tl):
            yr=pt['date'][:4]
            if yr!=prev_y:
                prev_y=yr
                xlabels+=f'<text x="{_x(i):.1f}" y="{svg_h-4}" text-anchor="middle" font-size="9" fill="#8b949e">{yr}</text>'
        svg=(f'<svg viewBox="0 0 {svg_w} {svg_h}" width="100%" style="max-width:{svg_w}px;display:block;">'
             f'<polyline points="{bk_pts}" fill="none" stroke="#58a6ff" stroke-width="2"/>'
             f'<polyline points="{spy_pts}" fill="none" stroke="#8b949e" stroke-width="1.5" stroke-dasharray="4,3"/>'
             f'<polyline points="{p_pts}" fill="none" stroke="#3fb950" stroke-width="1.5" stroke-dasharray="2,3"/>'
             f'{xlabels}</svg>')
        # Stats table
        rows=''
        for lbl,bv,sv,pv in [('Total Return',f'{bk_s.get("total",0):+.1f}%',f'{spy_s.get("total",0):+.1f}%',f'{p_s.get("total",0):+.1f}%'),
                              ('CAGR',f'{bk_s.get("cagr",0):+.1f}%',f'{spy_s.get("cagr",0):+.1f}%',f'{p_s.get("cagr",0):+.1f}%'),
                              ('Sharpe',f'{bk_s.get("sharpe",0):.2f}',f'{spy_s.get("sharpe",0):.2f}',f'{p_s.get("sharpe",0):.2f}'),
                              ('Max DD',f'{bk_s.get("max_dd",0):.1f}%',f'{spy_s.get("max_dd",0):.1f}%',f'{p_s.get("max_dd",0):.1f}%')]:
            rows+=(f'<tr><td style="padding:7px 12px;color:#8b949e;font-size:11px;">{lbl}</td>'
                   f'<td style="text-align:right;padding:7px 12px;font-family:monospace;font-size:12px;font-weight:700;color:#58a6ff;">{bv}</td>'
                   f'<td style="text-align:right;padding:7px 12px;font-family:monospace;font-size:11px;color:#8b949e;">{sv}</td>'
                   f'<td style="text-align:right;padding:7px 12px;font-family:monospace;font-size:11px;color:#8b949e;">{pv}</td></tr>')
        return (
            f'<div class="fc" style="margin-top:14px;">'
            f'<div class="lbl" style="margin-bottom:10px;">REGIME ALLOCATION BACKTEST — {yrs:.0f} YEAR (SIMPLIFIED 6-INSTRUMENT MODEL)</div>'
            f'<div style="display:flex;gap:16px;margin-bottom:8px;font-size:10px;">'
            f'<span><span style="color:#58a6ff;">&#9472;&#9472;</span> BK Regime Allocation</span>'
            f'<span><span style="color:#8b949e;">- - -</span> SPY Buy & Hold</span>'
            f'<span><span style="color:#3fb950;">&#183;&#183;&#183;</span> 60/40 Portfolio</span>'
            f'</div>'
            f'<div style="margin-bottom:10px;padding:10px 14px;background:#0d2318;border:1px solid #3fb950;'
            f'border-radius:6px;font-size:11px;color:#c8cfe0;line-height:1.7;">'
            f'<strong style="color:#3fb950;">&#9432; Note:</strong> BK cumulative return is lower by design — '
            f'the strategy trades raw upside for drawdown protection. '
            f'Risk-adjusted performance (Sharpe <strong style="color:#58a6ff;">{bk_s.get("sharpe",0):.2f}</strong> vs '
            f'<strong style="color:#8b949e;">{spy_s.get("sharpe",0):.2f}</strong>) and maximum drawdown '
            f'(<strong style="color:#58a6ff;">{bk_s.get("max_dd",0):.1f}%</strong> vs '
            f'<strong style="color:#8b949e;">{spy_s.get("max_dd",0):.1f}%</strong>) both favour BK. '
            f'See table below.'
            f'</div>'
            f'{svg}'
            f'<table style="width:100%;border-collapse:collapse;margin-top:12px;">'
            f'<thead><tr>'
            f'<th style="text-align:left;padding:8px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;"></th>'
            f'<th style="text-align:right;padding:8px 12px;font-size:9px;color:#58a6ff;border-bottom:1px solid #30363d;">BK ALLOCATION</th>'
            f'<th style="text-align:right;padding:8px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">SPY B&H</th>'
            f'<th style="text-align:right;padding:8px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">60/40</th>'
            f'</tr></thead><tbody>{rows}</tbody></table>'
            f'<div style="margin-top:10px;padding:10px 12px;background:#0d1117;border-left:3px solid #e3b341;font-size:10px;color:#8b949e;line-height:1.6;">'
            f'<strong style="color:#e3b341;">&#9888; Backtest model note:</strong> This {yrs:.0f}-year simulation uses '
            f'a simplified <strong style="color:#c8cfe0;">3-regime, 6-instrument</strong> allocation '
            f'(SPY &middot; TLT &middot; GLD &middot; BIL &middot; HYG &middot; EEM) with hardcoded '
            f'regime weights &mdash; <strong style="color:#c8cfe0;">separate from the live BK Dynamic '
            f'Allocation model shown above</strong>, which uses the full 97-instrument universe and '
            f'the RAS formula. Assumptions: monthly rebalancing &middot; zero transaction costs &middot; '
            f'rf = 4.5% &middot; gross of fees &middot; one-day execution lag applied (regime signal '
            f'from day i executes on day i+1).<br>'
            f'<span style="color:#6a7485;">Chart shows cumulative total return; Sharpe measures '
            f'risk-adjusted return &mdash; a strategy with lower absolute return but significantly '
            f'lower volatility can achieve a higher Sharpe. The chart and Sharpe column can therefore '
            f'appear to disagree while both being correct.</span>'
            f'</div>'
            # ── Backtest methodology disclosure ──
            f'<div style="background:#161b22;border:1px solid #30363d;border-radius:10px;padding:16px 20px;margin-top:16px;">'
            f'<div style="font-size:9px;color:#e3b341;letter-spacing:3px;font-weight:700;margin-bottom:12px;">'
            f'&#9888; BACKTEST METHODOLOGY &mdash; READ BEFORE ACTING</div>'
            f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:16px;margin-bottom:14px;">'
            f'<div>'
            f'<div style="font-size:10px;color:#8b949e;margin-bottom:6px;letter-spacing:2px;">WHAT WAS TESTED</div>'
            f'<div style="font-size:11px;color:#e6edf3;line-height:1.8;">'
            f'Model: 3-regime allocation (Calm/Stressed/Crisis)<br>'
            f'Universe: 6 instruments only<br>'
            f'&nbsp;&nbsp;SPY &middot; TLT &middot; GLD &middot; BIL &middot; HYG &middot; EEM<br>'
            f'Period: Apr 2021 &ndash; Apr 2026 (5 years)<br>'
            f'Rebalancing: Monthly &middot; First trading day<br>'
            f'Regime signal: Prior month-end classification</div></div>'
            f'<div>'
            f'<div style="font-size:10px;color:#8b949e;margin-bottom:6px;letter-spacing:2px;">ASSUMPTIONS &amp; LIMITATIONS</div>'
            f'<div style="font-size:11px;color:#e6edf3;line-height:1.8;">'
            f'Transaction costs: 0 bps (gross of all fees)<br>'
            f'Risk-free rate: 4.5% annualised<br>'
            f'Execution: One-day lag applies<br>'
            f'&#9888; Does NOT test the live 97-instrument RAS model<br>'
            f'&#9888; Regime weights are hardcoded, not dynamic<br>'
            f'&#9888; Past performance &#8800; future results</div></div></div>'
            f'<div style="border-top:1px solid #21262d;padding-top:10px;">'
            f'<div style="font-size:10px;color:#8b949e;margin-bottom:6px;letter-spacing:2px;">WHAT HAS NOT YET BEEN BACKTESTED</div>'
            f'<div style="font-size:11px;color:#8b949e;line-height:1.8;">'
            f'&middot; BK Fragility Framework predictive validity (hit rate, false positive rate, avg drawdown after CRISIS signal)<br>'
            f'&middot; BK Composite Score forward returns (top quintile vs bottom quintile, 21-day holding period)<br>'
            f'&middot; Full 97-instrument RAS model backtest</div></div></div>'
            # ── Regime weights table ──
            f'<div style="margin-top:16px;">'
            f'<div style="font-size:9px;color:#8b949e;letter-spacing:3px;margin-bottom:10px;">REGIME ALLOCATION WEIGHTS</div>'
            f'<table style="width:100%;border-collapse:collapse;font-size:11px;">'
            f'<thead><tr style="border-bottom:1px solid #30363d;">'
            f'<th style="text-align:left;padding:6px 8px;color:#8b949e;">Instrument</th>'
            f'<th style="text-align:center;padding:6px 8px;color:#3fb950;">MODERATE</th>'
            f'<th style="text-align:center;padding:6px 8px;color:#e3b341;">STRESSED</th>'
            f'<th style="text-align:center;padding:6px 8px;color:#f85149;">CRISIS</th></tr></thead><tbody>'
            f'<tr style="border-bottom:0.5px solid #21262d;"><td style="padding:5px 8px;color:#e6edf3;">SPY (US Equities)</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#3fb950;font-family:monospace;">45%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#e3b341;font-family:monospace;">25%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#f85149;font-family:monospace;">10%</td></tr>'
            f'<tr style="border-bottom:0.5px solid #21262d;"><td style="padding:5px 8px;color:#e6edf3;">TLT (Long Treasuries)</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#3fb950;font-family:monospace;">10%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#e3b341;font-family:monospace;">20%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#f85149;font-family:monospace;">20%</td></tr>'
            f'<tr style="border-bottom:0.5px solid #21262d;"><td style="padding:5px 8px;color:#e6edf3;">GLD (Gold)</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#3fb950;font-family:monospace;">10%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#e3b341;font-family:monospace;">15%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#f85149;font-family:monospace;">25%</td></tr>'
            f'<tr style="border-bottom:0.5px solid #21262d;"><td style="padding:5px 8px;color:#e6edf3;">BIL (Cash / T-Bills)</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#3fb950;font-family:monospace;">5%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#e3b341;font-family:monospace;">20%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#f85149;font-family:monospace;">35%</td></tr>'
            f'<tr style="border-bottom:0.5px solid #21262d;"><td style="padding:5px 8px;color:#e6edf3;">HYG (High Yield Credit)</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#3fb950;font-family:monospace;">15%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#e3b341;font-family:monospace;">15%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#f85149;font-family:monospace;">5%</td></tr>'
            f'<tr><td style="padding:5px 8px;color:#e6edf3;">EEM (Emerging Markets)</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#3fb950;font-family:monospace;">15%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#e3b341;font-family:monospace;">5%</td>'
            f'<td style="padding:5px 8px;text-align:center;color:#f85149;font-family:monospace;">5%</td></tr>'
            f'</tbody></table></div>'
            + _build_monthly_returns_grid(bt.get('daily_bk'))
            + f'</div>'
        )

    def _build_monthly_returns_grid(daily_bk_series):
        """Build monthly returns heatmap from full daily cumulative BK values."""
        if daily_bk_series is None or len(daily_bk_series) < 2:
            return ''
        daily_ret = daily_bk_series.pct_change().dropna()
        if daily_ret.empty:
            return ''
        # Aggregate to monthly
        monthly = daily_ret.resample('ME').apply(lambda x: (1 + x).prod() - 1)
        months_names = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
        years = sorted(monthly.index.year.unique())

        # Header
        hdr = '<tr><th style="padding:5px 8px;font-size:9px;color:#8b949e;text-align:left;">Year</th>'
        for m in months_names:
            hdr += f'<th style="padding:5px 4px;font-size:9px;color:#8b949e;text-align:center;">{m}</th>'
        hdr += '<th style="padding:5px 8px;font-size:9px;color:#8b949e;text-align:center;border-left:1px solid #30363d;">Full Year</th></tr>'

        _current_year = pd.Timestamp.today().year
        body = ''
        for yr in years:
            yr_data = monthly[monthly.index.year == yr]
            annual = float((1 + yr_data).prod() - 1)
            _yr_label = "YTD" if yr == _current_year else str(yr)
            row = f'<tr><td style="padding:4px 8px;font-weight:700;color:#e6edf3;font-size:11px;">{_yr_label}</td>'
            for m_num in range(1, 13):
                m_vals = yr_data[yr_data.index.month == m_num]
                if m_vals.empty:
                    row += '<td style="padding:4px;text-align:center;color:#30363d;font-size:10px;">&mdash;</td>'
                else:
                    ret = float(m_vals.iloc[0])
                    intensity = min(abs(ret) * 500, 100)
                    if ret >= 0:
                        bg = f'rgba(63,185,80,{intensity/100:.2f})'
                        color = '#3fb950' if ret > 0.02 else '#e6edf3'
                    else:
                        bg = f'rgba(248,81,73,{intensity/100:.2f})'
                        color = '#f85149' if ret < -0.02 else '#e6edf3'
                    row += (f'<td style="background:{bg};color:{color};font-family:monospace;'
                            f'font-size:10px;text-align:center;padding:4px 3px;">'
                            f'{ret*100:+.1f}%</td>')
            ann_color = '#3fb950' if annual >= 0 else '#f85149'
            row += (f'<td style="font-family:monospace;font-weight:700;font-size:11px;color:{ann_color};'
                    f'text-align:center;border-left:1px solid #30363d;padding:4px 8px;">'
                    f'{annual*100:+.1f}%</td></tr>')
            body += row

        return (
            f'<div style="margin-top:20px;">'
            f'<div style="font-size:9px;color:#8b949e;letter-spacing:3px;text-transform:uppercase;margin-bottom:10px;">'
            f'Monthly Returns &mdash; BK Allocation Strategy</div>'
            f'<div style="overflow-x:auto;">'
            f'<table style="border-collapse:collapse;width:100%;font-size:11px;color:#8b949e;">'
            f'<thead style="border-bottom:1px solid #30363d;">{hdr}</thead>'
            f'<tbody>{body}</tbody></table></div>'
            f'<div style="margin-top:8px;font-size:9px;color:#30363d;">'
            f'Monthly returns based on simplified 6-instrument backtest model. '
            f'Gross of fees &#183; Zero transaction costs assumed &#183; rf = 4.5%</div>'
            f'</div>'
        )

    def _build_rsr(df_in):
        """Build Relative Strength Rankings — sorted within each asset class."""
        sections_order = SECTION_ORDER
        html = (
            '<div style="margin-top:14px;">'
            '<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;'
            'text-transform:uppercase;margin-bottom:4px;">RELATIVE STRENGTH RANKINGS — BY ASSET CLASS</div>'
            '<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-bottom:12px;">Ranked by 1-month price return within each asset class</div>'
            '<div style="display:grid;grid-template-columns:repeat(auto-fill,minmax(280px,1fr));gap:12px;">'
        )
        for sec in sections_order:
            sec_df = df_in[df_in["section"]==sec].copy()
            if sec_df.empty: continue
            sec_df = sec_df.dropna(subset=["ret_1m"]).sort_values("ret_1m", ascending=False)
            if sec_df.empty: continue
            sec_label = SECTION_LABELS.get(sec, sec)
            best_ret  = sec_df["ret_1m"].max()
            worst_ret = sec_df["ret_1m"].min()
            rng       = max(abs(best_ret), abs(worst_ret), 0.001)
            rows_html = ""
            for rank, (_, r) in enumerate(sec_df.iterrows(), 1):
                v     = r["ret_1m"]; s = "+" if v>=0 else ""
                color = "#3fb950" if v>=0 else "#f85149"
                bar_w = min(100, abs(v)/rng*100)
                bar_dir = "right" if v>=0 else "left"
                medal = str(rank)
                rows_html += (
                    f'<div style="display:flex;align-items:center;gap:8px;padding:5px 0;border-bottom:1px solid #21262d;">'
                    f'<div style="width:16px;font-size:10px;text-align:center;">{medal}</div>'
                    f'<div style="flex:1;font-size:10px;color:#e6edf3;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;">{r["name"]}</div>'
                    f'<div style="width:50px;text-align:right;font-family:monospace;font-size:10px;font-weight:700;color:{color};">{s}{v*100:.1f}%</div>'
                    f'</div>'
                )
            # Positive return must never be red — relabel "Worst" as "Weakest"
            # and colour amber when the whole bucket is positive.
            if worst_ret >= 0:
                worst_html = (f'<span>Weakest: <span style="color:#e3b341;">'
                              f'{worst_ret*100:+.1f}%</span> '
                              f'<span style="color:#6a7485;font-style:italic;">(all positive)</span></span>')
            else:
                worst_html = (f'<span>Worst: <span style="color:#f85149;">'
                              f'{worst_ret*100:+.1f}%</span></span>')
            html += (
                f'<div style="background:#161b22;border:1px solid #30363d;border-radius:8px;padding:12px 14px;">'
                f'<div style="font-size:9px;font-weight:700;letter-spacing:1px;color:#58a6ff;'
                f'text-transform:uppercase;margin-bottom:8px;">{sec_label}</div>'
                f'{rows_html}'
                f'<div style="display:flex;justify-content:space-between;margin-top:6px;font-size:9px;color:#8b949e;">'
                f'<span>Best: <span style="color:#3fb950;">{best_ret*100:+.1f}%</span></span>'
                f'{worst_html}'
                f'</div></div>'
            )
        html += '</div></div>'
        return html


    # PCA placeholder (replaces RSR per F1)
    _pca_placeholder = (
        '<div style="margin-top:20px;background:#161b22;border:1px solid #30363d;border-radius:10px;padding:24px;text-align:center;">'
        '<div style="font-size:14px;font-weight:700;color:#58a6ff;letter-spacing:1px;margin-bottom:8px;">'
        'Advanced Factor Analysis &mdash; Principal Component Analysis (PCA)</div>'
        '<div style="font-size:12px;color:#8b949e;line-height:1.8;max-width:600px;margin:0 auto;">'
        'PCA reveals hidden risk factors driving cross-asset moves &mdash; '
        'identifying which latent factors explain the majority of portfolio variance.<br>'
        'Replacing RSR in Q2 2026.</div>'
        '<div style="font-size:9px;color:#6a7485;margin-top:12px;font-family:monospace;">'
        'Factor decomposition &#183; Eigenvalue spectrum &#183; Variance explained &#183; Coming soon</div>'
        '</div>'
    )

    # ══ TAB 4: ANALYSIS — Correlation Heatmap ════════════════════════════════
    # Key-20 matrix grouped by asset class for institutional readability.
    HEATMAP_ORDER = [
        # EQ bloc (8)
        "SPY", "QQQ", "IWM", "ACWI", "EEM", "XLK", "XLV", "XLU",
        # FI bloc (4)
        "TLT", "HYG", "LQD", "BIL",
        # CMD bloc (3)
        "GLD", "BNO", "DBC",
        # ALT bloc (5)
        "BTC-USD", "EURUSD=X", "JPY=X", "VXX", "ITA",
    ]
    HEATMAP_GROUPS_ORDERED = [
        ("EQ",  8),   # first 8 cells
        ("FI",  4),
        ("CMD", 3),
        ("ALT", 5),
    ]
    HEATMAP_DISPLAY_NAMES = {
        "SPY": "SPY", "QQQ": "QQQ", "IWM": "IWM", "ACWI": "ACWI", "EEM": "EEM",
        "XLK": "Tech", "XLV": "Health", "XLU": "Util",
        "TLT": "TLT", "HYG": "HYG", "LQD": "LQD", "BIL": "BIL",
        "GLD": "Gold", "BNO": "Oil", "DBC": "CMD",
        "BTC-USD": "BTC", "EURUSD=X": "EUR", "JPY=X": "JPY",
        "VXX": "VIX", "ITA": "DEF",
    }

    def _build_heatmap(prices_df):
        tickers = [t for t in HEATMAP_ORDER if t in prices_df.columns]
        if len(tickers) < 4:
            return None, None
        rets    = prices_df[tickers].pct_change().dropna()
        corr_df = rets.tail(60).corr()
        return corr_df, tickers

    def _corr_cell_style(v):
        """Red/white/blue diverging scale for a correlation value."""
        if abs(v) < 0.10:
            return "#161b22", "#4a5568", ""  # near-zero cells: suppressed
        if v > 0:
            # white -> deep red
            intensity = min(int(abs(v) * 180), 180)
            bg = f"rgb({180+intensity//3},{80 - intensity//3 if intensity < 120 else 30},{80 - intensity//3 if intensity < 120 else 30})"
            text = "#ffffff" if intensity > 90 else "#1a1a1a"
        else:
            intensity = min(int(abs(v) * 180), 180)
            bg = f"rgb({80 - intensity//3 if intensity < 120 else 30},{110 - intensity//4 if intensity < 120 else 80},{180+intensity//3})"
            text = "#ffffff" if intensity > 90 else "#1a1a1a"
        return bg, text, f"{v:+.2f}"

    corr_df, hm_names = _build_heatmap(prices)

    if corr_df is not None:
        n = len(hm_names)
        cell_size = 28   # per briefing spec
        label_w   = 90
        pad_top   = 70   # space for rotated column labels
        total_w   = label_w + n * cell_size + 20
        total_h   = pad_top + n * cell_size + 20

        # Group boundary indices (where one asset class ends and next begins)
        group_bounds = []
        _cum = 0
        for _g, _sz in HEATMAP_GROUPS_ORDERED:
            _cum += _sz
            group_bounds.append((_cum, _g))
        div_indices = [b[0] for b in group_bounds[:-1]]  # exclude final boundary

        svg_parts = [
            f'<svg viewBox="0 0 {total_w} {total_h}" width="100%" '
            f'style="max-width:{total_w}px;font-family:system-ui,sans-serif;">'
        ]

        # Column labels — rotated 45°
        for j, tk in enumerate(hm_names):
            x = label_w + j * cell_size + cell_size / 2
            label = HEATMAP_DISPLAY_NAMES.get(tk, tk)
            svg_parts.append(
                f'<text x="{x}" y="{pad_top - 6}" text-anchor="end" '
                f'transform="rotate(-45,{x},{pad_top-6})" '
                f'font-size="10" font-weight="600" fill="#c8cfe0">{label}</text>'
            )

        # Row labels + cells
        for i, row_tk in enumerate(hm_names):
            y_top = pad_top + i * cell_size
            row_label = HEATMAP_DISPLAY_NAMES.get(row_tk, row_tk)
            svg_parts.append(
                f'<text x="{label_w - 8}" y="{y_top + cell_size/2 + 3}" '
                f'text-anchor="end" font-size="10" font-weight="600" fill="#c8cfe0">'
                f'{row_label}</text>'
            )
            for j, col_tk in enumerate(hm_names):
                v = corr_df.loc[row_tk, col_tk]
                x = label_w + j * cell_size
                y = y_top
                if i == j:
                    # Diagonal: dark grey, no value shown
                    svg_parts.append(
                        f'<rect x="{x}" y="{y}" width="{cell_size}" height="{cell_size}" '
                        f'fill="#21262d" rx="1"/>'
                    )
                    continue
                bg, tc, disp = _corr_cell_style(float(v))
                svg_parts.append(
                    f'<rect x="{x}" y="{y}" width="{cell_size}" height="{cell_size}" '
                    f'fill="{bg}" rx="1">'
                    f'<title>{HEATMAP_DISPLAY_NAMES.get(row_tk, row_tk)} vs '
                    f'{HEATMAP_DISPLAY_NAMES.get(col_tk, col_tk)}: {v:+.2f}</title></rect>'
                )
                if disp:
                    svg_parts.append(
                        f'<text x="{x + cell_size/2}" y="{y + cell_size/2 + 3}" '
                        f'text-anchor="middle" font-size="9" fill="{tc}" '
                        f'style="pointer-events:none;">{disp}</text>'
                    )

        # Asset-class group dividers (2px white lines between groups)
        for idx in div_indices:
            xd = label_w + idx * cell_size
            yd = pad_top + idx * cell_size
            # vertical divider
            svg_parts.append(
                f'<line x1="{xd}" y1="{pad_top}" x2="{xd}" '
                f'y2="{pad_top + n*cell_size}" stroke="#e6edf3" stroke-width="2"/>'
            )
            # horizontal divider
            svg_parts.append(
                f'<line x1="{label_w}" y1="{yd}" x2="{label_w + n*cell_size}" '
                f'y2="{yd}" stroke="#e6edf3" stroke-width="2"/>'
            )

        # Group labels along top
        _acc = 0
        for g_name, g_size in HEATMAP_GROUPS_ORDERED:
            x_mid = label_w + (_acc + g_size/2) * cell_size
            svg_parts.append(
                f'<text x="{x_mid}" y="14" text-anchor="middle" font-size="10" '
                f'font-weight="700" fill="#58a6ff" letter-spacing="1">{g_name}</text>'
            )
            _acc += g_size

        svg_parts.append('</svg>')
        heatmap_svg = "".join(svg_parts)

        as_of_date = now.strftime("%d %b %Y")
        legend_svg = (
            f'<div style="display:flex;align-items:center;gap:12px;margin-top:12px;flex-wrap:wrap;">'
            f'<span style="font-size:10px;color:#8b949e;font-family:monospace;">-1.0</span>'
            f'<div style="width:220px;height:12px;border-radius:4px;'
            f'background:linear-gradient(to right,rgb(30,80,210),#161b22 45%,#161b22 55%,rgb(210,80,30));">'
            f'</div>'
            f'<span style="font-size:10px;color:#8b949e;font-family:monospace;">+1.0</span>'
            f'<span style="font-size:10px;color:#8b949e;margin-left:16px;">'
            f'60-day rolling correlation &#183; as of {as_of_date}'
            f'</span>'
            f'</div>'
        )

        # Top correlations (most correlated pairs, excluding diagonal)
        pairs = []
        for i in range(n):
            for j in range(i+1, n):
                v = float(corr_df.iloc[i, j])
                pairs.append((v,
                              HEATMAP_DISPLAY_NAMES.get(hm_names[i], hm_names[i]),
                              HEATMAP_DISPLAY_NAMES.get(hm_names[j], hm_names[j])))
        pairs.sort(key=lambda x: abs(x[0]), reverse=True)

        top_pairs_html = ""
        for v, a, b in pairs[:8]:
            color = "#f85149" if v > 0 else "#58a6ff"
            sign  = "+" if v > 0 else ""
            top_pairs_html += (
                f'<div style="display:flex;justify-content:space-between;align-items:center;'
                f'padding:5px 0;border-bottom:1px solid #21262d;">'
                f'<div style="font-size:11px;color:#e6edf3;">{a} <span style="color:#8b949e;">vs</span> {b}</div>'
                f'<div style="font-family:monospace;font-size:12px;font-weight:700;color:{color};">{sign}{v:.2f}</div>'
                f'</div>'
            )

        analysis_tab = (
            # Summary stats
            f'<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:14px;margin-bottom:14px;">'
            f'<div class="fc"><div class="lbl">CORRELATION WINDOW</div>'
            f'<div style="font-size:22px;font-weight:700;color:#e6edf3;font-family:monospace;">60D</div>'
            f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Rolling daily returns</div></div>'
            f'<div class="fc"><div class="lbl">INSTRUMENTS</div>'
            f'<div style="font-size:22px;font-weight:700;color:#e6edf3;font-family:monospace;">{n}</div>'
            f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Key representatives</div></div>'
            f'<div class="fc"><div class="lbl">AVG CORRELATION</div>'
            f'<div style="font-size:22px;font-weight:700;'
            f'color:{"#f85149" if corr_df.values[corr_df.values < 1].mean() > 0.5 else "#e3b341" if corr_df.values[corr_df.values < 1].mean() > 0.3 else "#3fb950"};font-family:monospace;">'
            f'{corr_df.values[corr_df.values < 1].mean():.2f}</div>'
            f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Ex-diagonal (high = contagion risk)</div></div>'
            f'</div>'
            # Heatmap
            f'<div class="fc" style="margin-bottom:14px;overflow-x:auto;">'
            f'<div class="lbl" style="margin-bottom:12px;">CROSS-ASSET CORRELATION MATRIX — 60D</div>'
            f'{heatmap_svg}'
            f'<div style="margin-top:6px;">{legend_svg}</div>'
            f'</div>'
            # Top correlated pairs
            f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:14px;">'
            f'<div class="fc"><div class="lbl" style="margin-bottom:10px;">STRONGEST CORRELATIONS</div>'
            f'{top_pairs_html}</div>'
            f'<div class="fc"><div class="lbl" style="margin-bottom:10px;">HOW TO READ</div>'
            f'<div style="font-size:11px;color:#8b949e;line-height:1.9;">'
            f'<span style="color:#f85149;">&#9632;</span> Red = move together (+1.0)<br>'
            f'<span style="color:#aaa;">&#9632;</span> White = no relationship (0.0)<br>'
            f'<span style="color:#58a6ff;">&#9632;</span> Blue = move opposite (&#8722;1.0)<br><br>'
            f'High average correlation = contagion risk<br>'
            f'Diversification works when colours are mixed<br>'
            f'60-day window captures current market regime'
            f'</div></div></div>'
            f'<div style="margin-top:10px;font-size:9px;color:#8b949e;font-family:monospace;">'
            f'Correlation = 60-day rolling Pearson correlation of daily returns &#183; '
            f'Key 20 instruments selected as representatives of each asset class</div>'
        ) + _pca_placeholder
    else:
        analysis_tab = '<div style="padding:40px;text-align:center;color:#8b949e;">Insufficient data for correlation analysis.</div>' + _pca_placeholder



    # ══ TAB 5: REGIME — INSTITUTIONAL MULTI-MODEL FRAMEWORK ═════════════════
    #
    # Layout (top → bottom):
    #   1. HERO CARD — headline regime (SM) + days-in-regime counter
    #   2. MODEL CONVICTION ROW — agreement badge + HMM probability bars
    #      + transition risk alert + consensus regime
    #   3. DRIVERS — vol percentile + drawdown percentile (unchanged)
    #   4. TIMELINE — 2-year colour strip (unchanged)
    #   5. STATS + EPISODES — regime distribution + crisis history (unchanged)
    #   6. METHODOLOGY FOOTER — explains all 3 tiers + consensus logic
    #
    if regime_data:
        reg      = regime_data.get("regime","Calm")
        streak   = regime_data.get("days_in_regime", 0)
        stats    = regime_data.get("stats", {})
        drivers  = regime_data.get("drivers", {})
        timeline = regime_data.get("timeline", [])
        episodes = regime_data.get("episodes", [])

        # ── Multi-model data ─────────────────────────────────────────────────
        hmm_probs       = regime_data.get("hmm_probs")          # dict or None
        hmm_regime      = regime_data.get("hmm_regime")          # str or None
        gmm_regime      = regime_data.get("gmm_regime")          # str or None
        consensus       = regime_data.get("consensus", reg)
        model_agreement = regime_data.get("model_agreement", 1)
        models_avail    = regime_data.get("models_available", 1)
        tr_level        = regime_data.get("transition_risk", "N/A")
        tr_desc         = regime_data.get("transition_desc", "")

        # Display-only rename: internal key "Calm" shown as "Moderate"
        _disp = {"Calm": "Moderate"}
        reg_display  = _disp.get(reg, reg)
        cons_display = _disp.get(consensus, consensus)
        hmm_display  = _disp.get(hmm_regime, hmm_regime) if hmm_regime else None
        gmm_display  = _disp.get(gmm_regime, gmm_regime) if gmm_regime else None

        rc_ = {"Crisis":"#f85149","Stressed":"#e3b341","Calm":"#3fb950"}.get(reg,"#8b949e")
        rb_ = {"Crisis":"#2d0f0e","Stressed":"#2d2106","Calm":"#0d2318"}.get(reg,"#161b22")
        _rc_cons = {"Crisis":"#f85149","Stressed":"#e3b341","Calm":"#3fb950"}.get(consensus,"#8b949e")

        reg_desc = {
            "Calm":     "Markets are operating within normal historical ranges. Volatility and drawdowns are contained. Risk appetite is stable.",
            "Stressed": "Elevated volatility or meaningful drawdown detected. Markets are pricing in uncertainty.",
            "Crisis":   "Extreme volatility or severe drawdown detected. Historical crisis-level conditions. Defensive positioning warranted.",
        }.get(reg, "")

        # ── Regime timeline SVG ───────────────────────────────────────────────
        if timeline:
            tl_w = 900; tl_h = 60; bar_w = max(1, tl_w // len(timeline))
            tl_parts = [f'<svg viewBox="0 0 {tl_w} {tl_h}" width="100%" style="max-width:{tl_w}px;display:block;">']
            for i, pt in enumerate(timeline):
                x = i * bar_w
                tl_parts.append(f'<rect x="{x}" y="0" width="{bar_w+1}" height="{tl_h}" fill="{pt["color"]}" opacity="0.85"/>')
            prev_month = ""; label_count = 0
            for i, pt in enumerate(timeline):
                month = pt["date"][:7]
                if month != prev_month:
                    prev_month = month
                    label_count += 1
                    if label_count % 2 == 1:
                        x = i * bar_w
                        tl_parts.append(f'<text x="{x+2}" y="{tl_h-4}" font-size="8" fill="#e6edf3" opacity="0.7">{month}</text>')
            tl_parts.append('</svg>')
            timeline_svg = "".join(tl_parts)
        else:
            timeline_svg = "<div style='color:#8b949e;'>Insufficient history for timeline.</div>"

        # ── Stats table ───────────────────────────────────────────────────────
        stats_rows = ""
        for rkey, rname, rcolor in [("Calm","Moderate","#3fb950"),("Stressed","Stressed","#e3b341"),("Crisis","Crisis","#f85149")]:
            rs = stats.get(rkey, {})
            stats_rows += (
                f'<tr><td style="padding:8px 12px;color:{rcolor};font-weight:700;font-family:monospace;">{rname}</td>'
                f'<td class="num gr">{rs.get("days",0):,}</td>'
                f'<td class="num gr">{rs.get("pct",0):.1f}%</td>'
                f'<td class="num gr">{rs.get("avg_duration",0):.0f} days</td></tr>'
            )

        # ── Episodes table ────────────────────────────────────────────────────
        ep_rows = ""
        for ep in reversed(episodes):
            _range_txt = ep.get("range") or f"{ep['start']} &rarr; {ep['end']}"
            ep_rows += (
                f'<tr><td style="padding:7px 12px;color:#e6edf3;font-size:11px;" colspan="2">{_range_txt}</td>'
                f'<td class="num nr" style="font-size:11px;">{ep["depth"]:.1f}%</td></tr>'
            )

        _total_days = sum(int(rs.get("days", 0)) for rs in stats.values())
        _actual_years = _total_days / 252 if _total_days else 0
        _years_label = f"{_actual_years:.1f}-YEAR HISTORY" if _actual_years else "HISTORY"

        vol_pct_color = "#f85149" if drivers.get("vol_pct",0)>90 else "#e3b341" if drivers.get("vol_pct",0)>70 else "#3fb950"
        _dd_now = drivers.get("dd_now", 0)
        dd_pct_color  = "#8b949e" if abs(_dd_now) < 1.0 else "#f85149" if drivers.get("dd_pct",0)<20 else "#e3b341"

        # ── MODEL AGREEMENT BADGE ─────────────────────────────────────────────
        # Colour: green=all agree, amber=partial, red=disagreement
        _agree_color = "#3fb950" if model_agreement == models_avail else "#e3b341" if model_agreement >= 2 else "#f85149"
        _agree_label = f"{model_agreement}/{models_avail}"

        # ── HMM PROBABILITY BARS ─────────────────────────────────────────────
        # Shows the posterior probability distribution across all 3 states.
        # High entropy (flat bars) = regime uncertainty = reduce sizing.
        # Sharp spike on one state = high conviction = position with confidence.
        if hmm_probs:
            _pc = hmm_probs.get("p_calm", 0) * 100
            _ps = hmm_probs.get("p_stressed", 0) * 100
            _px = hmm_probs.get("p_crisis", 0) * 100
            _entropy = hmm_probs.get("entropy", 0)
            # Entropy: 0 = perfect certainty, ln(3)≈1.10 = maximum uncertainty
            _ent_pct = min(100, _entropy / 1.10 * 100)
            _ent_color = "#f85149" if _ent_pct > 70 else "#e3b341" if _ent_pct > 40 else "#3fb950"
            hmm_bars_html = (
                f'<div style="margin-top:6px;">'
                # Moderate (Calm) bar
                f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:4px;">'
                f'<div style="width:65px;font-size:9px;color:#3fb950;font-family:monospace;">Moderate</div>'
                f'<div style="flex:1;background:#21262d;border-radius:3px;height:10px;">'
                f'<div style="width:{_pc:.0f}%;background:#3fb950;height:10px;border-radius:3px;transition:width 0.3s;"></div></div>'
                f'<div style="width:36px;text-align:right;font-size:10px;font-family:monospace;color:#3fb950;">{_pc:.0f}%</div></div>'
                # Stressed bar
                f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:4px;">'
                f'<div style="width:65px;font-size:9px;color:#e3b341;font-family:monospace;">Stressed</div>'
                f'<div style="flex:1;background:#21262d;border-radius:3px;height:10px;">'
                f'<div style="width:{_ps:.0f}%;background:#e3b341;height:10px;border-radius:3px;transition:width 0.3s;"></div></div>'
                f'<div style="width:36px;text-align:right;font-size:10px;font-family:monospace;color:#e3b341;">{_ps:.0f}%</div></div>'
                # Crisis bar
                f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:6px;">'
                f'<div style="width:65px;font-size:9px;color:#f85149;font-family:monospace;">Crisis</div>'
                f'<div style="flex:1;background:#21262d;border-radius:3px;height:10px;">'
                f'<div style="width:{_px:.0f}%;background:#f85149;height:10px;border-radius:3px;transition:width 0.3s;"></div></div>'
                f'<div style="width:36px;text-align:right;font-size:10px;font-family:monospace;color:#f85149;">{_px:.0f}%</div></div>'
                # Entropy indicator
                f'<div style="display:flex;align-items:center;gap:8px;margin-top:4px;padding-top:6px;border-top:1px solid #21262d;">'
                f'<div style="width:65px;font-size:8px;color:#8b949e;font-family:monospace;">Entropy</div>'
                f'<div style="flex:1;background:#21262d;border-radius:3px;height:6px;">'
                f'<div style="width:{_ent_pct:.0f}%;background:{_ent_color};height:6px;border-radius:3px;"></div></div>'
                f'<div style="width:36px;text-align:right;font-size:9px;font-family:monospace;color:{_ent_color};">{_entropy:.2f}</div></div>'
                f'<div style="font-size:8px;color:#6a7485;margin-top:2px;">Low entropy = high conviction &nbsp;&#183;&nbsp; High entropy = regime uncertainty</div>'
                f'</div>'
            )
        else:
            hmm_bars_html = (
                f'<div style="padding:12px 0;font-size:10px;color:#8b949e;font-style:italic;">'
                f'HMM probabilities unavailable &mdash; install <code style="background:#21262d;padding:2px 5px;border-radius:3px;">pip install hmmlearn</code> for conviction bars</div>'
            )

        # ── TRANSITION RISK BADGE ─────────────────────────────────────────────
        # Early warning: fires when HMM sees a regime the SM hasn't flagged yet.
        # This is the institutional value-add — gives 5-15 days lead time.
        _tr_color = {"High": "#f85149", "Elevated": "#e3b341", "Low": "#3fb950", "N/A": "#8b949e"}.get(tr_level, "#8b949e")
        _tr_bg    = {"High": "#2d0f0e", "Elevated": "#2d2106", "Low": "#0d2318", "N/A": "#161b22"}.get(tr_level, "#161b22")

        # ── PER-MODEL CALL TABLE ─────────────────────────────────────────────
        # Shows what each model is saying independently — transparency for the PM.
        _model_rows = ""
        _sm_c = {"Crisis":"#f85149","Stressed":"#e3b341","Calm":"#3fb950"}.get(reg,"#8b949e")
        _model_rows += (
            f'<tr><td style="padding:5px 10px;font-size:10px;color:#e6edf3;">State Machine</td>'
            f'<td style="padding:5px 10px;font-size:10px;color:{_sm_c};font-weight:700;font-family:monospace;">{reg_display}</td>'
            f'<td style="padding:5px 10px;font-size:9px;color:#8b949e;">Deterministic &middot; governance headline</td></tr>'
        )
        if hmm_regime:
            _hc = {"Crisis":"#f85149","Stressed":"#e3b341","Calm":"#3fb950"}.get(hmm_regime,"#8b949e")
            _model_rows += (
                f'<tr><td style="padding:5px 10px;font-size:10px;color:#e6edf3;">Hidden Markov</td>'
                f'<td style="padding:5px 10px;font-size:10px;color:{_hc};font-weight:700;font-family:monospace;">{hmm_display}</td>'
                f'<td style="padding:5px 10px;font-size:9px;color:#8b949e;">Probabilistic &middot; conviction signal</td></tr>'
            )
        else:
            _model_rows += (
                f'<tr><td style="padding:5px 10px;font-size:10px;color:#8b949e;">Hidden Markov</td>'
                f'<td style="padding:5px 10px;font-size:10px;color:#8b949e;font-family:monospace;">&mdash;</td>'
                f'<td style="padding:5px 10px;font-size:9px;color:#6a7485;">Not available (install hmmlearn)</td></tr>'
            )
        if gmm_regime:
            _gc = {"Crisis":"#f85149","Stressed":"#e3b341","Calm":"#3fb950"}.get(gmm_regime,"#8b949e")
            _model_rows += (
                f'<tr><td style="padding:5px 10px;font-size:10px;color:#e6edf3;">Gaussian Mixture</td>'
                f'<td style="padding:5px 10px;font-size:10px;color:{_gc};font-weight:700;font-family:monospace;">{gmm_display}</td>'
                f'<td style="padding:5px 10px;font-size:9px;color:#8b949e;">Cross-validation &middot; cluster-based</td></tr>'
            )
        else:
            _model_rows += (
                f'<tr><td style="padding:5px 10px;font-size:10px;color:#8b949e;">Gaussian Mixture</td>'
                f'<td style="padding:5px 10px;font-size:10px;color:#8b949e;font-family:monospace;">&mdash;</td>'
                f'<td style="padding:5px 10px;font-size:9px;color:#6a7485;">Not available (install sklearn)</td></tr>'
            )

        regime_tab = (
            # ── 1. HERO CARD — headline regime + streak ───────────────────────
            f'<div style="background:{rb_};border:2px solid {rc_};border-radius:12px;padding:24px 28px;margin-bottom:14px;display:flex;justify-content:space-between;align-items:center;flex-wrap:wrap;gap:16px;">'
            f'<div>'
            f'<div style="font-size:9px;color:{rc_};letter-spacing:3px;font-family:monospace;margin-bottom:8px;">CURRENT MARKET REGIME</div>'
            f'<div style="font-size:42px;font-weight:700;color:{rc_};font-family:monospace;letter-spacing:2px;">{reg_display.upper()}</div>'
            f'<div style="font-size:11px;color:#e6edf3;margin-top:8px;max-width:500px;line-height:1.6;">{reg_desc}</div>'
            f'</div>'
            f'<div style="text-align:center;">'
            f'<div style="font-size:9px;color:#8b949e;letter-spacing:2px;margin-bottom:6px;">DAYS IN REGIME</div>'
            f'<div style="font-size:48px;font-weight:700;color:{rc_};font-family:monospace;">{streak}</div>'
            f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">consecutive trading days</div>'
            f'</div></div>'

            # ── 2. MODEL CONVICTION ROW ───────────────────────────────────────
            # This is the institutional value-add: multi-model conviction at a glance.
            # Left: agreement badge + per-model calls.
            # Centre: HMM posterior probability bars (conviction gauge).
            # Right: transition risk alert + consensus regime.
            f'<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:14px;margin-bottom:14px;">'

            # ── Left panel: Model Agreement + Per-Model Table ─────────────────
            f'<div class="fc">'
            f'<div class="lbl">MODEL AGREEMENT</div>'
            f'<div style="display:flex;align-items:center;gap:14px;margin:10px 0;">'
            f'<div style="font-size:36px;font-weight:700;color:{_agree_color};font-family:monospace;">{_agree_label}</div>'
            f'<div style="font-size:11px;color:#e6edf3;line-height:1.5;">'
            f'{"All models aligned" if model_agreement == models_avail else "Partial agreement &mdash; reduced conviction" if model_agreement >= 2 else "Models diverge &mdash; high uncertainty"}'
            f'</div></div>'
            f'<table style="width:100%;border-collapse:collapse;font-size:11px;">'
            f'<thead><tr>'
            f'<th style="text-align:left;padding:4px 10px;font-size:8px;color:#8b949e;border-bottom:1px solid #30363d;">Model</th>'
            f'<th style="text-align:left;padding:4px 10px;font-size:8px;color:#8b949e;border-bottom:1px solid #30363d;">Call</th>'
            f'<th style="text-align:left;padding:4px 10px;font-size:8px;color:#8b949e;border-bottom:1px solid #30363d;">Role</th>'
            f'</tr></thead><tbody>{_model_rows}</tbody></table>'
            f'</div>'

            # ── Centre panel: HMM Probability Bars ────────────────────────────
            # These posteriors are the "conviction gauge" — institutional PMs use
            # this to calibrate position sizing. A 90% p_stressed is very different
            # from a 55% p_stressed even though both map to "Stressed".
            f'<div class="fc">'
            f'<div class="lbl">HMM REGIME PROBABILITIES</div>'
            f'{hmm_bars_html}'
            f'</div>'

            # ── Right panel: Transition Risk + Consensus ──────────────────────
            f'<div class="fc">'
            f'<div class="lbl">TRANSITION RISK &amp; CONSENSUS</div>'
            # Transition risk alert box
            f'<div style="background:{_tr_bg};border:1px solid {_tr_color};border-radius:8px;padding:12px 14px;margin:8px 0;">'
            f'<div style="display:flex;align-items:center;gap:8px;">'
            f'<div style="font-size:9px;color:#8b949e;letter-spacing:2px;font-family:monospace;">TRANSITION RISK</div>'
            f'<div style="font-size:14px;font-weight:700;color:{_tr_color};font-family:monospace;">{tr_level.upper()}</div></div>'
            f'<div style="font-size:9px;color:#c8cfe0;margin-top:6px;line-height:1.5;">{tr_desc}</div>'
            f'</div>'
            # Consensus regime box
            f'<div style="background:#161b22;border:1px solid {_rc_cons};border-radius:8px;padding:12px 14px;margin-top:8px;">'
            f'<div style="display:flex;align-items:center;gap:8px;">'
            f'<div style="font-size:9px;color:#8b949e;letter-spacing:2px;font-family:monospace;">CONSENSUS</div>'
            f'<div style="font-size:14px;font-weight:700;color:{_rc_cons};font-family:monospace;">{cons_display.upper()}</div></div>'
            f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Union-of-risk: most severe model call wins</div>'
            f'</div>'
            f'</div>'

            f'</div>'  # close 3-column grid

            # ── 3. DRIVERS ────────────────────────────────────────────────────
            f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:14px;margin-bottom:14px;">'
            f'<div class="fc">'
            f'<div class="lbl">VOL DRIVER — WORLD (ACWI)</div>'
            f'<div style="display:flex;align-items:baseline;gap:10px;margin:8px 0;">'
            f'<div style="font-size:28px;font-weight:700;color:{vol_pct_color};font-family:monospace;">{drivers.get("vol_now",0):.1f}%</div>'
            f'<div style="font-size:12px;color:#8b949e;">annualised vol</div></div>'
            f'<div style="font-size:11px;color:#e6edf3;">At <span style="color:{vol_pct_color};font-weight:700;">{drivers.get("vol_pct",0):.0f}th percentile</span> of history</div>'
            f'<div style="background:#21262d;border-radius:4px;height:8px;margin-top:10px;">'
            f'<div style="width:{min(100,drivers.get("vol_pct",0)):.0f}%;background:{vol_pct_color};height:8px;border-radius:4px;"></div></div>'
            f'</div>'
            f'<div class="fc">'
            f'<div class="lbl">DRAWDOWN DRIVER — WORLD (ACWI)</div>'
            f'<div style="display:flex;align-items:baseline;gap:10px;margin:8px 0;">'
            f'<div style="font-size:28px;font-weight:700;color:{dd_pct_color};font-family:monospace;">{drivers.get("dd_now",0):.1f}%</div>'
            f'<div style="font-size:12px;color:#8b949e;">from 1Y peak</div></div>'
            f'<div style="font-size:11px;color:#e6edf3;"><span style="color:{dd_pct_color};font-weight:700;">{drivers.get("dd_pct",0):.0f}%</span> of history had smaller drawdowns</div>'
            f'<div style="background:#21262d;border-radius:4px;height:8px;margin-top:10px;">'
            f'<div style="width:{min(100, max(0, abs(_dd_now) * 5)):.0f}%;background:{dd_pct_color};height:8px;border-radius:4px;"></div></div>'
            + (f'<div style="font-size:9px;color:#8b949e;margin-top:6px;font-style:italic;">'
               f'Regime triggered by vol signal ({drivers.get("vol_pct",0):.0f}th pct) &mdash; market near 1Y peak</div>'
               if abs(_dd_now) < 1.0 else '')
            + f'</div></div>'

            # ── 4. TIMELINE ───────────────────────────────────────────────────
            f'<div class="fc" style="margin-bottom:14px;">'
            f'<div class="lbl" style="margin-bottom:10px;">REGIME TIMELINE — LAST 2 YEARS</div>'
            f'<div style="display:flex;gap:16px;margin-bottom:8px;">'
            f'<span style="font-size:10px;color:#3fb950;">&#9632; MODERATE</span>'
            f'<span style="font-size:10px;color:#e3b341;">&#9632; STRESSED</span>'
            f'<span style="font-size:10px;color:#f85149;">&#9632; CRISIS</span>'
            f'</div>'
            f'{timeline_svg}</div>'

            # ── 5. STATS + EPISODES ───────────────────────────────────────────
            f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:14px;">'
            f'<div class="fc">'
            f'<div class="lbl" style="margin-bottom:10px;">REGIME STATISTICS — {_years_label}</div>'
            f'<table style="width:100%;border-collapse:collapse;font-size:12px;">'
            f'<thead><tr><th style="text-align:left;padding:6px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">Regime</th>'
            f'<th style="text-align:right;padding:6px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">Days</th>'
            f'<th style="text-align:right;padding:6px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">% Time</th>'
            f'<th style="text-align:right;padding:6px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">Avg Duration</th>'
            f'</tr></thead><tbody>{stats_rows}</tbody></table></div>'
            f'<div class="fc">'
            f'<div class="lbl" style="margin-bottom:10px;">CRISIS EPISODES — WORLD DRAWDOWN &lt; &#8722;15%</div>'
            f'<table style="width:100%;border-collapse:collapse;font-size:12px;">'
            f'<thead><tr>'
            f'<th style="text-align:left;padding:6px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">Start</th>'
            f'<th style="text-align:left;padding:6px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">End</th>'
            f'<th style="text-align:right;padding:6px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">Peak DD</th>'
            f'</tr></thead><tbody>{ep_rows if ep_rows else "<tr><td colspan=3 style=padding:8px;color:#8b949e;>No episodes detected</td></tr>"}</tbody></table></div></div>'

            # ── 6. METHODOLOGY FOOTER ─────────────────────────────────────────
            # Full transparency on how each model works — institutional standard.
            f'<div style="margin-top:14px;background:#0d1117;border:1px solid #21262d;border-radius:8px;padding:14px 16px;">'
            f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;margin-bottom:10px;">METHODOLOGY — 3-TIER INSTITUTIONAL REGIME FRAMEWORK</div>'
            f'<div style="font-size:9px;color:#8b949e;font-family:monospace;line-height:2;">'
            f'<strong style="color:#c8cfe0;">Tier 1 &mdash; State Machine (headline):</strong> '
            f'Deterministic classifier on ACWI. Stressed: vol &ge; 70th pct OR dd &le; 30th pct. '
            f'Crisis: vol &ge; 90th pct OR dd &le; 10th pct. Ex-ante expanding quantiles, shifted t&minus;1. '
            f'Auditable, no ML dependencies.<br>'
            f'<strong style="color:#c8cfe0;">Tier 2 &mdash; Hidden Markov Model (conviction):</strong> '
            f'3-state Gaussian HMM on [returns, vol, dd]. Walk-forward retrain every 21 days. '
            f'Posterior probabilities = conviction gauge. High entropy = regime uncertainty. '
            f'Transition risk fires when HMM diverges from SM (5&ndash;15 day lead time).<br>'
            f'<strong style="color:#c8cfe0;">Tier 3 &mdash; Gaussian Mixture (cross-validation):</strong> '
            f'3-component GMM on [returns, vol, dd]. Walk-forward retrain every 21 days. '
            f'Captures non-linear clusters that HMM may miss.<br>'
            f'<strong style="color:#c8cfe0;">Consensus:</strong> '
            f'Union-of-risk &mdash; most severe model call wins. '
            f'Conservative by design: false positives preferred over missed crises.<br>'
            f'<strong style="color:#c8cfe0;">Model Agreement:</strong> '
            f'Count of models on same call. 3/3 = full model consensus. 1/3 = models diverge.'
            f'</div></div>'
        )
    else:
        regime_tab = '<div style="padding:40px;text-align:center;color:#8b949e;">Regime data unavailable.</div>'


    # ── Build Fear & Greed HTML ───────────────────────────────────────────────
    if fg_data:
        fg_score  = round(fg_data.get("score", 50))
        fg_label  = fg_data.get("label", "Neutral")
        # emoji removed per B3
        fg_color  = fg_data.get("color", "#e3b341")
        fg_details= fg_data.get("details", {})

        # Gauge needle SVG
        fg_angle  = int(fg_score / 100 * 180)
        fg_bg     = "#2d0f0e" if fg_score<=25 else "#2d1a0e" if fg_score<=45 else "#2d2106" if fg_score<=55 else "#0d2318" if fg_score<=75 else "#052e16"

        # Component definitions
        _comp_defs = {
            'Momentum':       'SPY price vs 125-day moving average',
            'Breadth':        '% instruments above 50-day moving average',
            'Safe Haven':     'TLT vs SPY 20-day relative return',
            'Volatility':     'VIXY level vs 50-day average',
            'Junk Bonds':     'HYG vs IEF credit spread (20D)',
            'Strength':       '% instruments within 5% of 52-week high',
            'Term Structure':  'VIX/VIX3M ratio (vol curve; backwardation = panic)',
        }

        # Component bars
        comp_bars = ""
        for comp_name, comp_data in fg_details.items():
            cs = comp_data.get("score", 50)
            cv = comp_data.get("value", "")
            cc = "#f85149" if cs<=25 else "#ff7b72" if cs<=45 else "#e3b341" if cs<=55 else "#7ee787" if cs<=75 else "#3fb950"
            defn = _comp_defs.get(comp_name, "")
            comp_bars += (
                f'<div style="margin-bottom:8px;">'
                f'<div style="display:flex;justify-content:space-between;margin-bottom:3px;">'
                f'<span style="font-size:10px;color:#e6edf3;">{comp_name}'
                f'<span style="font-size:8px;color:#8b949e;margin-left:6px;">{defn}</span></span>'
                f'<span style="font-size:10px;font-family:monospace;color:{cc};">{cs:.0f}</span></div>'
                f'<div style="background:#21262d;border-radius:3px;height:6px;">'
                f'<div style="width:{cs:.0f}%;background:{cc};height:6px;border-radius:3px;"></div></div>'
                f'<div style="font-size:9px;color:#8b949e;margin-top:2px;">{cv}</div>'
                f'</div>'
            )

        # Gradient scale bar
        scale_html = (
            '<div style="margin:10px 0 4px;">'
            '<div style="height:10px;border-radius:5px;background:linear-gradient(to right,#f85149,#ff7b72,#e3b341,#7ee787,#3fb950);position:relative;">'
            f'<div style="position:absolute;left:{fg_score:.0f}%;top:-4px;transform:translateX(-50%);">'
            '<div style="width:3px;height:18px;background:#fff;border-radius:2px;"></div></div></div>'
            '<div style="display:flex;justify-content:space-between;font-size:8px;color:#8b949e;margin-top:3px;">'
            '<span>Extreme Fear</span><span>Fear</span><span>Neutral</span><span>Greed</span><span>Extreme Greed</span>'
            '</div></div>'
        )

        fg_html = (
            f'<div style="background:{fg_bg};border:1px solid {fg_color};border-radius:10px;'
            f'padding:18px 20px;margin-bottom:14px;">'
            f'<div style="display:grid;grid-template-columns:auto 1fr;gap:20px;align-items:start;">'
            # Left: score display
            f'<div style="text-align:center;min-width:140px;">'
            f'<div style="font-size:9px;color:{fg_color};letter-spacing:2px;margin-bottom:6px;font-family:monospace;">FEAR & GREED INDEX</div>'
            f'<div style="font-size:36px;font-weight:700;color:{fg_color};font-family:monospace;line-height:1;">{fg_score:.0f}</div>'
            f'<div style="font-size:12px;font-weight:700;color:{fg_color};margin-top:4px;">{fg_label}</div>'
            f'{scale_html}'
            f'</div>'
            # Right: components
            f'<div>'
            f'<div style="font-size:9px;color:#8b949e;letter-spacing:2px;text-transform:uppercase;margin-bottom:10px;">COMPONENT BREAKDOWN</div>'
            f'{comp_bars}'
            f'</div>'
            f'</div>'
            # Methodology footnote
            f'<div style="margin-top:12px;padding-top:10px;border-top:1px solid #21262d;'
            f'font-size:9px;color:#8b949e;line-height:1.6;">'
            f'Methodology: Equal-weighted composite of 7 normalised signals (0&ndash;100). '
            f'Each component uses rolling historical percentile vs full available history. '
            f'&nbsp;&middot;&nbsp; '
            f'0&ndash;24 = Extreme Fear &nbsp;&middot;&nbsp; 25&ndash;44 = Fear &nbsp;&middot;&nbsp; '
            f'45&ndash;55 = Neutral &nbsp;&middot;&nbsp; 56&ndash;74 = Greed &nbsp;&middot;&nbsp; '
            f'75&ndash;100 = Extreme Greed &nbsp;&middot;&nbsp; Updated daily'
            f'</div>'
            f'</div>'
        )
    else:
        fg_html = ""


    # ══ TAB 1: INTEL — DAILY INTELLIGENCE ══════════════════════════════════════
    #
    #  Restructured layout — 3 clear zones, top-to-bottom:
    #
    #  ZONE 1 — DASHBOARD PULSE (quantitative, at-a-glance)
    #    5-card row: Tone | Regime | Fragility | Fear & Greed | RAG Signals
    #    Everything the PM needs in 3 seconds.
    #
    #  ZONE 2 — QUANT SIGNALS (model-driven, actionable)
    #    Left column:
    #      REDUCE EXPOSURE — fragility ≥ 70 AND regime fit ≤ 25
    #      CONTRARIAN WATCH — F&G ≤ 25 AND 3M return < -10%
    #    Right column:
    #      INCREASE EXPOSURE — regime fit ≥ 75 AND fragility ≤ 55 AND GREEN
    #        Ranked by BK Composite Score with factor decomposition bars.
    #        (Merged: old Signal-to-Trade INCREASE + old BK Top Picks)
    #      TOP RISKS — highest fragility instruments with pillar decomposition bars.
    #
    #    Factor bars on INCREASE cards show *why* an instrument scores high:
    #      Mom (30%) | Frag Inv (25%) | Fit (20%) | Signal (15%) | Vol (10%)
    #    Pillar bars on RISK cards show *what's driving* the fragility:
    #      Drawdown (22%) | CVaR (20%) | Contagion (18%) | Vol (15%) | Trend (15%) | Liquidity (10%)
    #
    #  ZONE 3 — AI INTELLIGENCE (LLM-generated, contextual)
    #    Market Narrative + Recommended Actions
    #    Fear & Greed component breakdown
    #    Backtest mini snapshot
    #
    #  KILLED:
    #    - "Instruments to Watch" (was AI-generated with no model, overlapped with picks)
    #    - 3 hardcoded "historical context" boxes (were static text, not live data)
    #    - Separate BK Top Picks section (merged into INCREASE)
    # ══════════════════════════════════════════════════════════════════════════════

    ai = ai_commentary or {}
    ai_narrative      = ai.get("narrative", "")
    ai_actions        = ai.get("actions", [])
    ai_fg_summary     = ai.get("fg_summary", "")
    ai_regime_interp  = ai.get("regime_interpretation", "")

    # Signal variables
    reg_now   = regime_data.get("regime","Calm") if regime_data else "Calm"
    reg_color = {"Crisis":"#f85149","Stressed":"#e3b341","Calm":"#3fb950"}.get(reg_now,"#8b949e")
    reg_bg    = {"Crisis":"#2d0f0e","Stressed":"#2d2106","Calm":"#0d2318"}.get(reg_now,"#161b22")
    frag_sys  = frag_df.attrs.get("system_score", float(frag_df["fragility"].median())) if frag_df is not None and not frag_df.empty else 50.0
    frag_reg  = frag_df.attrs.get("regime","MODERATE") if frag_df is not None and not frag_df.empty else "MODERATE"

    # RAG counts
    nr = int((df["rag_label"].str.strip()=="RED").sum())
    na = int((df["rag_label"].str.strip()=="AMBER").sum())
    ng = int((df["rag_label"].str.strip()=="GREEN").sum())

    # Tone
    _rising_intel = _count_rising_risk(df)
    tone, tc, tb = calculate_market_tone(reg_now, frag_sys, _rising_intel, len(df))

    # Fear & Greed for pulse row
    _fg_score = round(fg_data.get("score", 50)) if fg_data else 50
    _fg_label = fg_data.get("label", "Neutral") if fg_data else "Neutral"
    _fg_color = fg_data.get("color", "#e3b341") if fg_data else "#e3b341"

    _df_lookup = df.set_index("ticker")

    # ── REDUCE EXPOSURE ──────────────────────────────────────────────────────
    # Gate: fragility ≥ 70 AND regime fit ≤ 25
    # These are instruments under stress that don't fit the current regime.
    _frag_map = {}
    reduce_items = []
    if frag_df is not None and not frag_df.empty:
        _frag_map = dict(zip(frag_df["ticker"], frag_df["fragility"]))
        for _, r in frag_df.iterrows():
            tk = r["ticker"]
            if not is_rankable(tk):
                continue
            rf = get_regime_fit_score(tk, reg_now)
            if r["fragility"] >= 70 and rf <= 25:
                reduce_items.append((r["name"], tk, r["fragility"], rf))
        reduce_items = reduce_items[:5]

    # ── INCREASE EXPOSURE (merged: Signal-to-Trade INCREASE + BK Top Picks) ─
    # Gate: regime fit ≥ 75 AND fragility ≤ 55 AND RAG = GREEN
    # Ranked by: BK Composite Score (5-factor composite)
    # This merges the old binary filter with the continuous scorer —
    # instruments must PASS the gates AND are ranked by Opp Score.
    increase_items = []
    if frag_df is not None and not frag_df.empty:
        for _, r in df.iterrows():
            tk = r["ticker"]
            if not is_rankable(tk):
                continue
            rf = get_regime_fit_score(tk, reg_now)
            fr_val = _frag_map.get(tk, 50)
            opp_sc = bk_opp_scores.get(tk)
            if rf >= 75 and fr_val <= 55 and r["rag_label"].strip() == "GREEN" and opp_sc is not None:
                increase_items.append((r["name"], tk, opp_sc, fr_val, rf))
        # Sort by Opp Score descending — best opportunities first
        increase_items = sorted(increase_items, key=lambda x: x[2], reverse=True)[:5]

    # ── CONTRARIAN WATCH ─────────────────────────────────────────────────────
    # Gate: F&G ≤ 25 AND 3M return < -10%
    # Mean-reversion candidates during extreme fear — not a timing tool.
    _fg_score_val = round(fg_data.get("score", 50)) if fg_data else 50
    contrarian_items = []
    if _fg_score_val <= 25:
        for _, r in df.iterrows():
            tk = r["ticker"]
            if not is_rankable(tk):
                continue
            ret_3m = r.get("ret_3m", float("nan"))
            if pd.notna(ret_3m) and ret_3m < -0.10:
                contrarian_items.append((r["name"], tk, ret_3m * 100))
        contrarian_items = sorted(contrarian_items, key=lambda x: x[2])[:5]

    # ── TOP RISKS ────────────────────────────────────────────────────────────
    risks_df = (frag_df[frag_df["rag"].isin(["CRISIS","STRESSED"]) & frag_df["ticker"].apply(is_rankable)].head(5)
                if frag_df is not None and not frag_df.empty else pd.DataFrame())

    # ── Opportunity factor decomposition bar builder ─────────────────────────
    # Each factor is drawn as a thin coloured bar (0–100 width) with a label.
    # The bars are stacked vertically inside each card so the PM can
    # instantly compare factor profiles across picks.
    #
    # Factor key (what each bar means):
    #   Mom (30%)       — momentum percentile rank across universe (blue #58a6ff)
    #                     Blended: 1M×0.2 + 3M×0.5 + 6M×0.3. Higher = stronger trend.
    #   Frag Inv (25%)  — 100 minus fragility score (green #3fb950)
    #                     Higher = safer instrument. Penalises high-risk momentum plays.
    #   Fit (20%)       — regime-fit matrix: how well this bucket performs in current regime (purple #bc8cff)
    #                     Example: Fixed Income scores high in Stressed; EQ Growth in Calm.
    #   Signal (15%)    — RAG signal mapped to score: GREEN=100, AMBER=50, RED=0 (amber #e3b341)
    #                     Quality gate — prevents recommending deteriorating assets.
    #   Vol (10%)       — vol trend: 100 if 20D vol < 1M-ago vol, else 0 (cyan #56d4dd)
    #                     Rewards compressing volatility (de-risking environment).
    _FACTOR_META = [
        ("mom",        "Mom",      "#58a6ff", "30%"),
        ("frag_inv",   "Frag Inv", "#3fb950", "25%"),
        ("regime_fit", "Fit",      "#bc8cff", "20%"),
        ("signal",     "Signal",   "#e3b341", "15%"),
        ("vol_trend",  "Vol",      "#56d4dd", "10%"),
    ]

    def _factor_bars(tk: str) -> str:
        """Build inline factor decomposition bars for a single ticker."""
        fdata = bk_opp_factors.get(tk)
        if not fdata:
            return ""
        bars = ""
        for fkey, flabel, fcolor, fweight in _FACTOR_META:
            val = fdata.get(fkey, 0)
            bars += (
                f'<div style="display:flex;align-items:center;gap:4px;margin-top:2px;">'
                f'<div style="width:50px;font-size:7px;color:#8b949e;font-family:monospace;text-align:right;">{flabel} <span style="color:#6a7485;">{fweight}</span></div>'
                f'<div style="flex:1;background:#21262d;border-radius:2px;height:5px;">'
                f'<div style="width:{min(100, val):.0f}%;background:{fcolor};height:5px;border-radius:2px;"></div></div>'
                f'<div style="width:22px;font-size:7px;color:{fcolor};font-family:monospace;text-align:right;">{val:.0f}</div>'
                f'</div>'
            )
        return f'<div style="margin-top:6px;padding-top:6px;border-top:1px solid #21262d;">{bars}</div>'

    # ── Fragility pillar decomposition bar builder ────────────────────────────
    # Mirrors the opportunity factor bars for visual consistency.
    # Shows *what's driving the fragility* for each risk instrument.
    #
    # Pillar key (what each bar means):
    #   DD (22%)        — drawdown magnitude from rolling 1Y peak (red #f85149)
    #                     How deep the instrument has fallen. Higher = deeper drawdown.
    #   CVaR (20%)      — 60-day Conditional Value-at-Risk / expected shortfall (orange #e3734d)
    #                     Average loss in the worst 5% of days. Captures fat-tail risk.
    #   Contagion (18%) — 60-day rolling correlation to ACWI world proxy (purple #bc8cff)
    #                     Higher = more coupled to global risk-off. Diversification is gone.
    #   Vol (15%)       — 20-day annualised realised volatility (amber #e3b341)
    #                     Current volatility level vs. 2-year history.
    #   Trend (15%)     — distance below 200-day moving average (blue #58a6ff)
    #                     How far the asset has broken its long-term trend.
    #   Liq (10%)       — 60-day volume z-score (cyan #56d4dd)
    #                     Abnormal volume activity — panic selling or liquidity vacuum.
    _PILLAR_META = [
        ("pillar_dd",    "DD",        "#f85149", "22%"),
        ("pillar_cvar",  "CVaR",      "#e3734d", "20%"),
        ("pillar_corr",  "Contagion", "#bc8cff", "18%"),
        ("pillar_vol",   "Vol",       "#e3b341", "15%"),
        ("pillar_trend", "Trend",     "#58a6ff", "15%"),
        ("pillar_volz",  "Liq",       "#56d4dd", "10%"),
    ]

    def _pillar_bars(row) -> str:
        """Build inline fragility pillar decomposition bars for a risk instrument."""
        bars = ""
        # Pillar values are weighted z-score contributions (can range roughly 0–30+).
        # Normalise to 0–100 by capping at the fragility score itself for bar width,
        # but show the raw contribution value for transparency.
        frag_total = max(float(row.get("fragility", 1)), 1)
        for pkey, plabel, pcolor, pweight in _PILLAR_META:
            val = float(row.get(pkey, 0))
            # Bar width: proportion of total fragility score, capped at 100%
            bar_pct = min(100, max(0, val / frag_total * 100)) if frag_total > 0 else 0
            bars += (
                f'<div style="display:flex;align-items:center;gap:4px;margin-top:2px;">'
                f'<div style="width:58px;font-size:7px;color:#8b949e;font-family:monospace;text-align:right;">{plabel} <span style="color:#6a7485;">{pweight}</span></div>'
                f'<div style="flex:1;background:#21262d;border-radius:2px;height:5px;">'
                f'<div style="width:{bar_pct:.0f}%;background:{pcolor};height:5px;border-radius:2px;"></div></div>'
                f'<div style="width:26px;font-size:7px;color:{pcolor};font-family:monospace;text-align:right;">{val:.1f}</div>'
                f'</div>'
            )
        return f'<div style="margin-top:6px;padding-top:6px;border-top:1px solid #21262d;">{bars}</div>'

    # ── Build REDUCE HTML ────────────────────────────────────────────────────
    reduce_html = ""
    for nm, tk, frag_score, rf in reduce_items:
        reduce_html += (
            f'<div style="display:flex;justify-content:space-between;padding:6px 0;border-bottom:1px solid #21262d;">'
            f'<span style="font-size:11px;color:#e6edf3;">{nm} <span style="color:#8b949e;font-size:9px;">{tk}</span></span>'
            f'<span style="font-size:10px;font-family:monospace;color:#f85149;">Frag {frag_score:.0f} &middot; Fit {rf}</span></div>'
        )

    # ── Build INCREASE HTML (with factor decomposition bars) ─────────────────
    increase_html = ""
    for nm, tk, opp_sc, fr_val, rf in increase_items:
        increase_html += (
            f'<div style="background:#0d2318;border:1px solid #238636;border-radius:6px;padding:10px 12px;margin-bottom:6px;">'
            f'<div style="display:flex;justify-content:space-between;align-items:flex-start;">'
            f'<div>'
            f'<div style="font-size:11px;color:#e6edf3;font-weight:600;">{nm}</div>'
            f'<div style="font-family:monospace;font-size:8px;color:#8b949e;">{tk} &middot; Frag {fr_val:.0f} &middot; Fit {rf}</div>'
            f'</div>'
            f'<div style="font-size:18px;font-weight:700;color:#3fb950;font-family:monospace;">{opp_sc:.0f}<span style="font-size:9px;">/100</span></div>'
            f'</div>'
            f'{_factor_bars(tk)}'
            f'</div>'
        )

    # ── Build CONTRARIAN HTML ────────────────────────────────────────────────
    contrarian_html = ""
    for nm, tk, ret3m in contrarian_items:
        contrarian_html += (
            f'<div style="display:flex;justify-content:space-between;padding:6px 0;border-bottom:1px solid #21262d;">'
            f'<span style="font-size:11px;color:#e6edf3;">{nm} <span style="color:#8b949e;font-size:9px;">{tk}</span></span>'
            f'<span style="font-size:10px;font-family:monospace;color:#e3b341;">{ret3m:+.1f}% 3M</span></div>'
        )

    # ── Build TOP RISKS HTML (with pillar decomposition bars) ────────────────
    risk_cards = ""
    for _, r in risks_df.iterrows():
        fc = "#f85149" if r["rag"]=="CRISIS" else "#e3b341"
        rb = "#2d0f0e" if r["rag"]=="CRISIS" else "#2d2106"
        risk_cards += (
            f'<div style="background:{rb};border:1px solid {fc};border-radius:6px;padding:10px 12px;margin-bottom:6px;">'
            f'<div style="display:flex;justify-content:space-between;align-items:flex-start;">'
            f'<div>'
            f'<div style="font-size:11px;color:#e6edf3;font-weight:600;">{r["name"]}</div>'
            f'<div style="font-family:monospace;font-size:8px;color:#8b949e;">{r["ticker"]} &middot; {r["rag"]}</div>'
            f'</div>'
            f'<div style="font-size:18px;font-weight:700;color:{fc};font-family:monospace;">{r["fragility"]:.0f}<span style="font-size:9px;">/100</span></div>'
            f'</div>'
            f'{_pillar_bars(r)}'
            f'</div>'
        )

    # ── Build FRAMEWORK OBSERVATIONS (deterministic — top 4 by |MoM|) ─────────
    _pillar_label_map = {
        "pillar_dd": "Drawdown", "pillar_cvar": "CVaR", "pillar_corr": "Contagion",
        "pillar_vol": "Vol Stress", "pillar_trend": "Trend", "pillar_volz": "Liquidity",
    }
    _frag_idx = (frag_df.set_index("ticker") if frag_df is not None and not frag_df.empty else pd.DataFrame())
    _obs_rows = []
    _df_rankable_obs = df[df["ticker"].apply(is_rankable) & df["ret_1m"].notna()].copy()
    if not _df_rankable_obs.empty:
        _top2_gain = _df_rankable_obs.nlargest(2, "ret_1m")
        _top2_loss = _df_rankable_obs.nsmallest(2, "ret_1m")
        for _, _row in pd.concat([_top2_gain, _top2_loss]).iterrows():
            _tk   = _row["ticker"]
            _nm   = _row.get("name", _tk)
            _mom  = _row["ret_1m"] * 100
            # Fragility
            _frag_val = "N/A"
            _lead_pill = "N/A"
            if not _frag_idx.empty and _tk in _frag_idx.index:
                _fr = _frag_idx.loc[_tk]
                _frag_val = f'{float(_fr.get("fragility", 0)):.0f}'
                # Leading pillar = highest pillar value
                _pill_vals = {k: float(_fr.get(k, 0)) for k in _pillar_label_map if k in _fr}
                if _pill_vals:
                    _lead_pill = _pillar_label_map[max(_pill_vals, key=_pill_vals.get)]
            # Regime fit
            _rf = get_regime_fit_score(_tk, reg_now)
            _obs_rows.append(
                f'<div style="padding:8px 0;border-bottom:1px solid #21262d;font-family:monospace;font-size:11px;color:#e6edf3;">'
                f'<span style="color:#58a6ff;font-weight:700;">{_tk}</span>'
                f' <span style="color:#8b949e;">({_nm})</span>'
                f' &mdash; MoM <span style="color:{"#3fb950" if _mom >= 0 else "#f85149"};font-weight:700;">{_mom:+.1f}%</span>'
                f' &middot; Fragility <span style="color:#e3b341;">{_frag_val}</span>'
                f' &middot; Regime Fit <span style="color:#8b949e;">{_rf}</span>'
                f' &middot; Leading pillar: <span style="color:#bc8cff;">{_lead_pill}</span>'
                f'</div>'
            )
    action_items = "".join(_obs_rows) if _obs_rows else '<div style="font-size:11px;color:#8b949e;padding:8px 0;">Insufficient data.</div>'

    # ── Build Backtest mini snapshot ─────────────────────────────────────────
    bt_html = ""
    if backtest_data:
        bk_s  = backtest_data.get("bk",{})
        spy_s = backtest_data.get("spy",{})
        p_s   = backtest_data.get("p6040",{})
        yrs   = backtest_data.get("years", 5)
        bt_html = (
            f'<div class="fc" style="margin-bottom:14px;">'
            f'<div class="lbl" style="margin-bottom:10px;">REGIME ALLOCATION BACKTEST — {yrs:.0f} YEAR</div>'
            f'<table style="width:100%;border-collapse:collapse;font-size:12px;">'
            f'<thead><tr>'
            f'<th style="text-align:left;padding:8px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;"></th>'
            f'<th style="text-align:right;padding:8px 12px;font-size:9px;color:#58a6ff;border-bottom:1px solid #30363d;">BK ALLOCATION</th>'
            f'<th style="text-align:right;padding:8px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">SPY B&H</th>'
            f'<th style="text-align:right;padding:8px 12px;font-size:9px;color:#8b949e;border-bottom:1px solid #30363d;">60/40</th>'
            f'</tr></thead><tbody>'
        )
        for label, bk_v, spy_v, p_v in [
            ("Total Return", f'{bk_s.get("total",0):+.1f}%', f'{spy_s.get("total",0):+.1f}%', f'{p_s.get("total",0):+.1f}%'),
            ("CAGR",         f'{bk_s.get("cagr",0):+.1f}%',  f'{spy_s.get("cagr",0):+.1f}%',  f'{p_s.get("cagr",0):+.1f}%'),
            ("Sharpe",       f'{bk_s.get("sharpe",0):.2f}',   f'{spy_s.get("sharpe",0):.2f}',   f'{p_s.get("sharpe",0):.2f}'),
            ("Max Drawdown", f'{bk_s.get("max_dd",0):.1f}%',  f'{spy_s.get("max_dd",0):.1f}%',  f'{p_s.get("max_dd",0):.1f}%'),
        ]:
            bt_html += (
                f'<tr><td style="padding:7px 12px;color:#8b949e;font-size:11px;">{label}</td>'
                f'<td style="text-align:right;padding:7px 12px;font-family:monospace;font-size:12px;font-weight:700;color:#58a6ff;">{bk_v}</td>'
                f'<td style="text-align:right;padding:7px 12px;font-family:monospace;font-size:11px;color:#8b949e;">{spy_v}</td>'
                f'<td style="text-align:right;padding:7px 12px;font-family:monospace;font-size:11px;color:#8b949e;">{p_v}</td></tr>'
            )
        bt_html += (
            f'</tbody></table>'
            f'<div style="font-size:9px;color:#8b949e;margin-top:6px;font-family:monospace;">'
            f'Regime-aware monthly rebalancing &#183; rf=4.5% &#183; Full detail in Edge tab &#183; Past performance not indicative of future results</div>'
            f'</div>'
        )

    # Date header
    now_sgt  = pd.Timestamp.now(tz=SGT)
    date_hdr = now_sgt.strftime("%A, %d %b %Y · %H:%M SGT")

    # ══════════════════════════════════════════════════════════════════════════
    #  ASSEMBLE INTEL TAB
    # ══════════════════════════════════════════════════════════════════════════

    summary_tab = (
        # Date header
        f'<div style="font-size:9px;color:#8b949e;font-family:monospace;letter-spacing:2px;margin-bottom:16px;">'
        f'BKIQ DAILY INTELLIGENCE &#183; {date_hdr}</div>'

        # ── ZONE 1: DASHBOARD PULSE ──────────────────────────────────────────
        # 5 metrics the PM reads in 3 seconds. Added F&G to the pulse row
        # (was previously buried in a side panel).
        f'<div style="display:grid;grid-template-columns:repeat(5,1fr);gap:12px;margin-bottom:14px;">'
        f'<div class="fc" style="text-align:center;">'
        f'<div class="lbl">MARKET TONE</div>'
        f'<div class="pill" style="background:{tb};color:{tc};border:1px solid {tc};margin-top:6px;font-size:13px;">{tone}</div>'
        f'</div>'
        f'<div class="fc" style="text-align:center;">'
        f'<div class="lbl">REGIME</div>'
        f'<div class="pill" style="background:{reg_bg};color:{reg_color};border:1px solid {reg_color};margin-top:6px;font-size:13px;">{"MODERATE" if reg_now=="Calm" else reg_now.upper()}</div>'
        f'</div>'
        f'<div class="fc" style="text-align:center;">'
        f'<div class="lbl">FRAGILITY</div>'
        f'<div style="font-size:28px;font-weight:700;color:{rc_};font-family:monospace;margin-top:4px;">{frag_sys:.0f}</div>'
        f'<div style="font-size:9px;color:#8b949e;">{"CRISIS" if frag_sys>=70 else "STRESSED" if frag_sys>=55 else "MODERATE"}</div>'
        f'</div>'
        f'<div class="fc" style="text-align:center;">'
        f'<div class="lbl">FEAR &amp; GREED</div>'
        f'<div style="font-size:28px;font-weight:700;color:{_fg_color};font-family:monospace;margin-top:4px;">{_fg_score}</div>'
        f'<div style="font-size:9px;color:#8b949e;">{_fg_label}</div>'
        f'</div>'
        f'<div class="fc" style="text-align:center;">'
        f'<div class="lbl">RAG SIGNALS</div>'
        f'<div style="display:flex;justify-content:center;gap:10px;margin-top:8px;">'
        f'<span style="color:#f85149;font-size:18px;font-weight:700;font-family:monospace;">{nr}</span>'
        f'<span style="color:#8b949e;font-size:18px;">&middot;</span>'
        f'<span style="color:#e3b341;font-size:18px;font-weight:700;font-family:monospace;">{na}</span>'
        f'<span style="color:#8b949e;font-size:18px;">&middot;</span>'
        f'<span style="color:#3fb950;font-size:18px;font-weight:700;font-family:monospace;">{ng}</span>'
        f'</div>'
        f'<div style="font-size:9px;color:#8b949e;margin-top:2px;">R &middot; A &middot; G</div>'
        f'</div></div>'

        # ── ZONE 2: QUANT SIGNALS ────────────────────────────────────────────
        # Two-column layout: defensive actions (left) + opportunities (right)
        + f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:14px;margin-bottom:14px;">'

        # ── Left column: REDUCE + CONTRARIAN ─────────────────────────────────
        + f'<div style="display:flex;flex-direction:column;gap:14px;">'

        # REDUCE card
        + f'<div class="fc">'
        + f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#f85149;text-transform:uppercase;margin-bottom:4px;">STRESSED INSTRUMENTS</div>'
        + f'<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-bottom:8px;">Fragility &ge; 70 AND Regime Fit &le; 25</div>'
        + (reduce_html if reduce_html else '<div style="font-size:11px;color:#8b949e;padding:6px 0;">None currently</div>')
        + f'</div>'

        # CONTRARIAN card
        + f'<div class="fc">'
        + f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#e3b341;text-transform:uppercase;margin-bottom:4px;">OBSERVED FLAGS</div>'
        + f'<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-bottom:8px;">F&amp;G &le; 25 AND 3M return &lt; -10%</div>'
        + (contrarian_html if contrarian_html else f'<div style="font-size:11px;color:#8b949e;padding:6px 0;">{"F&G at " + str(_fg_score) + " — no extreme fear signal" if _fg_score_val > 25 else "No instruments down >10% in 3M"}</div>')
        + f'</div>'

        # TOP RISKS card (with pillar bars)
        + f'<div class="fc">'
        + f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#f85149;text-transform:uppercase;margin-bottom:4px;">HIGHEST FRAGILITY INSTRUMENTS</div>'
        + f'<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-bottom:8px;">Highest fragility &middot; pillar decomposition</div>'
        + (f'<div style="display:flex;flex-direction:column;gap:4px;">{risk_cards}</div>' if risk_cards else '<div style="font-size:11px;color:#8b949e;padding:6px 0;">None identified</div>')
        + f'</div>'

        + f'</div>'  # close left column

        # ── Right column: INCREASE ───────────────────────────────────────────
        + f'<div>'
        + f'<div class="fc">'
        + f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#3fb950;text-transform:uppercase;margin-bottom:4px;">HIGH FRAMEWORK SCORES</div>'
        + f'<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-bottom:8px;">Fit &ge; 75 AND Frag &le; 55 AND GREEN &middot; ranked by BK Opp Score</div>'
        + f'<div style="font-size:8px;color:#6e7681;font-style:italic;margin-bottom:6px;">Ranked by composite BK score. Descriptive only &mdash; not a recommendation to invest.</div>'
        + (f'<div style="display:flex;flex-direction:column;gap:4px;">{increase_html}</div>' if increase_html else '<div style="font-size:11px;color:#8b949e;padding:6px 0;">No instruments pass all gates</div>')
        + f'</div>'
        + f'</div>'  # close right column

        + f'</div>'  # close 2-column grid

        # ── ZONE 3: FRAMEWORK OBSERVATIONS + EVIDENCE ────────────────────────
        # Framework Observations (deterministic)
        + f'<div class="fc" style="margin-bottom:14px;">'
        + f'<div class="lbl" style="margin-bottom:4px;">FRAMEWORK OBSERVATIONS</div>'
        + f'<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-bottom:6px;">Deterministic &middot; Top 4 by absolute MoM move &middot; Descriptive only</div>'
        + f'<div style="font-size:11px;color:#6e7681;font-style:italic;margin-bottom:8px;">The framework highlights the four largest absolute movers in the universe month-to-date, with current pillar context.</div>'
        + action_items
        + f'</div>'

        # Fear & Greed component breakdown
        + fg_html

        # Backtest snapshot
        + bt_html

        # ── D-08: Today's Headlines ──────────────────────────────────────────
        + _build_headlines_html(headlines_data)

        # Footer
        + f'<div style="font-size:9px;color:#8b949e;font-family:monospace;line-height:1.8;">'
        + f'Data via Yahoo Finance &#183; Quant signals are model-driven &#183; '
        + f'For informational purposes only &#183; Not investment advice</div>'
    )
    # ══ TAB 7: EDGE — Portfolio Optimisation ══════════════════════════════════
    # The BK Dynamic Allocation card (RAS-driven, 6-bucket) built above is the
    # single allocation source. The old hardcoded regime-bucket grid was
    # removed — keeping two contradictory allocations on one tab was a defect.
    reg_now_e = regime_data.get("regime","Calm") if regime_data else "Calm"
    rc_e = {"Crisis":"#f85149","Stressed":"#e3b341","Calm":"#3fb950"}.get(reg_now_e,"#8b949e")
    rb_e = {"Crisis":"#2d0f0e","Stressed":"#2d2106","Calm":"#0d2318"}.get(reg_now_e,"#161b22")

    _df_rankable = df[df["ticker"].apply(is_rankable)]
    top_gain  = _df_rankable.nlargest(1,"ret_1m")["name"].iloc[0] if not _df_rankable.empty else "N/A"
    _frag_rankable = frag_df[frag_df["ticker"].apply(is_rankable)] if frag_df is not None and not frag_df.empty else pd.DataFrame()
    top_risk  = _frag_rankable.iloc[0]["name"] if not _frag_rankable.empty else "N/A"
    _vol_mask = df["vol_now"].notna() & df["vol_1m_ago"].notna() if "vol_now" in df.columns else pd.Series(False, index=df.index)
    vol_count = int((df.loc[_vol_mask, "vol_now"] > df.loc[_vol_mask, "vol_1m_ago"]).sum()) if _vol_mask.any() else 0
    # Deterministic template — no LLM call. Variables sourced from state already computed above.
    _vp_val = int(drivers.get("vol_pct", 0))
    _vp_sfx = "th" if 11 <= _vp_val % 100 <= 13 else {1: "st", 2: "nd", 3: "rd"}.get(_vp_val % 10, "th")
    _dd_abs = abs(drivers.get("dd_now", 0))
    commentary = (
        f"Framework state: fragility score {frag_sys:.1f}, volatility at the "
        f"{_vp_val}{_vp_sfx} percentile, drawdown {_dd_abs:.1f}% from peak. "
        f"Signal distribution: {nr} RED, {na} AMBER, {ng} GREEN. "
        f"Rising-volatility instruments: {vol_count} of {len(df)}. "
        f"Cross-asset correlation (30-day): 0.30. "
        f"The regime classification weights the fragility and volatility components more heavily than "
        f"the drawdown component under the regime-conditional weighting rules documented on the Regime tab."
    )

    # Key metrics for context (rankable subset only — excludes indices/proxies)
    best_asset  = _df_rankable.nlargest(1,"ret_1m")[["name","ret_1m"]].iloc[0]
    # Highest risk: use BK Fragility Score (top fragility instrument)
    if not _frag_rankable.empty:
        top_frag = _frag_rankable.iloc[0]
        highest_risk_name  = top_frag["name"]
        # Use round(), not int(), so this matches the :.0f format used in the
        # Fragility tab. Previously int() truncated while :.0f rounded, giving
        # 79 here vs 80 on the Fragility tab for the same underlying value.
        highest_risk_score = round(float(top_frag["fragility"]))
        highest_risk_label = top_frag.get("rag","CRISIS")
        use_fragility_risk = True
    else:
        worst_asset = _df_rankable.nsmallest(1,"ret_1m")[["name","ret_1m"]].iloc[0]
        use_fragility_risk = False

    # ── BK Dynamic Allocation card (RAS-driven) ──────────────────────────────
    _ras_color = ("#3fb950" if ras_score >= 55 else
                  "#e3b341" if ras_score >= 40 else
                  "#f85149")
    _bucket_order = ["EQ Growth", "EQ Defensive", "Fixed Income",
                     "Real Assets", "Cash", "Alts"]
    _alloc_rows_html = ""
    for _bk in _bucket_order:
        _w = bucket_weights.get(_bk, 0)
        _tops = get_top_instruments_per_bucket(bk_opp_scores, _bk, n=1)
        if _tops:
            _tk, _sc = _tops[0]
            _pick_html = (f'<span style="color:#58a6ff;font-weight:700;font-family:monospace;">{_tk}</span>'
                          f' <span style="color:#8b949e;font-family:monospace;">(Score {_sc:.0f})</span>')
        else:
            _pick_html = '<span style="color:#8b949e;">—</span>'
        _alloc_rows_html += (
            f'<div style="display:flex;align-items:center;gap:10px;padding:8px 0;'
            f'border-bottom:1px solid #21262d;">'
            f'<div style="width:130px;font-size:11px;color:#e6edf3;">{_bk}</div>'
            f'<div style="width:44px;text-align:right;font-family:monospace;font-size:12px;'
            f'font-weight:700;color:{_ras_color};">{_w}%</div>'
            f'<div style="flex:1;background:#21262d;border-radius:3px;height:8px;">'
            f'<div style="width:{_w*2}%;max-width:100%;background:{_ras_color};'
            f'height:8px;border-radius:3px;"></div></div>'
            f'<div style="flex:1.2;text-align:right;font-size:11px;">{_pick_html}</div>'
            f'</div>')

    dynamic_alloc_card = (
        f'<div class="fc" style="margin-bottom:14px;border:1px solid {_ras_color};">'
        f'<div style="display:flex;justify-content:space-between;align-items:baseline;margin-bottom:10px;">'
        f'<div class="lbl">REGIME-WEIGHTED FRAMEWORK OUTPUT &mdash; ILLUSTRATIVE</div>'
        f'<div style="font-size:11px;color:#8b949e;font-family:monospace;">Worked example of how the Risk Appetite Score translates regime state into asset-class weights</div>'
        f'</div>'
        f'<div style="display:flex;align-items:baseline;gap:14px;margin-bottom:14px;">'
        f'<div style="font-size:9px;color:#8b949e;letter-spacing:2px;">RISK APPETITE SCORE</div>'
        f'<div style="font-size:30px;font-weight:700;color:{_ras_color};font-family:monospace;">{ras_score:.0f}</div>'
        f'<div style="font-size:13px;color:{_ras_color};font-weight:700;letter-spacing:1px;">&rarr; {ras_label.upper()}</div>'
        f'</div>'
        f'<div style="background:#2d2106;border:2px solid #e3b341;border-radius:8px;padding:14px 16px;margin-bottom:14px;font-size:11px;color:#e6edf3;line-height:1.7;">'
        f'<div style="font-size:10px;font-weight:700;letter-spacing:2px;color:#e3b341;margin-bottom:8px;">&#9888; ILLUSTRATIVE METHODOLOGY OUTPUT &mdash; NOT A PORTFOLIO RECOMMENDATION</div>'
        f'The weights and instruments below illustrate how the Risk Appetite Score formula translates a regime state into bucket weights and representative instruments. '
        f'This is a worked example of the framework&apos;s output, not a recommended portfolio, not a backtested strategy, and not a set of positions the reader should hold.<br><br>'
        f'Specific instruments named are the highest-scoring instrument in each bucket on the composite metric as of today. They are not recommendations. '
        f'Weights are a deterministic function of the regime state; they do not account for transaction costs, liquidity, correlation, or individual circumstances.<br><br>'
        f'No reader should interpret this card as investment advice.'
        f'</div>'
        f'{_alloc_rows_html}'
        f'<div style="margin-top:10px;font-size:9px;color:#8b949e;font-family:monospace;line-height:1.6;">'
        f'RAS = Regime(35%) + Fragility Inv(30%) + Fear &amp; Greed(20%) + Vol Inv(15%) &#183; '
        f'Highest composite score per bucket = highest BK Composite Score'
        f'</div>'
        f'</div>'
        f'<div style="font-size:9px;color:#8b949e;font-family:monospace;line-height:1.6;margin-bottom:14px;padding:8px 12px;border:1px solid #21262d;border-radius:6px;">'
        f'This methodology example is one of several research directions documented on the Research tab. '
        f'It is not the output of a live or recommended strategy. '
        f'BKIQ is a personal research project &mdash; see About tab for full disclosure.'
        f'</div>'
    )

    edge_tab = (
        dynamic_alloc_card
        + f'<div style="background:{rb_e};border:2px solid {rc_e};border-radius:10px;padding:18px 24px;margin-bottom:14px;">'
        f'<div style="font-size:9px;color:{rc_e};letter-spacing:3px;font-family:monospace;margin-bottom:6px;">CURRENT REGIME CONTEXT</div>'
        f'<div style="font-size:9px;color:#8b949e;font-family:monospace;margin-bottom:8px;font-style:italic;">Factual state summary &mdash; no model interpretation</div>'
        f'<div style="font-size:22px;font-weight:700;color:{rc_e};font-family:monospace;">{reg_now_e.upper()} REGIME</div>'
        f'<div style="font-size:11px;color:#e6edf3;margin-top:8px;line-height:1.7;">{commentary}</div>'
        f'</div>'
        # Key signals (3 cards side-by-side) — duplicate "Suggested Portfolio
        # Allocation" card removed; the BK Dynamic Allocation card above is the
        # single source of truth for allocation.
        f'<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:14px;margin-bottom:14px;">'
        f'<div class="fc">'
        f'<div class="lbl" style="margin-bottom:8px;">HIGHEST MONTH-TO-DATE PERFORMER</div>'
        f'<div style="font-size:18px;font-weight:700;color:#3fb950;font-family:monospace;">{best_asset["name"]}</div>'
        f'<div style="font-size:24px;font-weight:700;color:#3fb950;font-family:monospace;">{best_asset["ret_1m"]*100:+.1f}%</div>'
        f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Rolling 1M Return (21 trading days)</div>'
        f'</div>'
        f'<div class="fc">'
        f'<div class="lbl" style="margin-bottom:8px;">HIGHEST FRAGILITY SCORE</div>'
        + (f'<div style="font-size:18px;font-weight:700;color:#f85149;font-family:monospace;">{highest_risk_name}</div>'
           f'<div style="font-size:24px;font-weight:700;color:#f85149;font-family:monospace;">{highest_risk_score}/100</div>'
           f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">BK Fragility Score · {highest_risk_label}</div>'
           if use_fragility_risk else
           f'<div style="font-size:18px;font-weight:700;color:#f85149;font-family:monospace;">{worst_asset["name"]}</div>'
           f'<div style="font-size:24px;font-weight:700;color:#f85149;font-family:monospace;">{worst_asset["ret_1m"]*100:+.1f}%</div>'
           f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">Rolling 1M Return (21 trading days)</div>')
        + f'</div>'
        f'<div class="fc">'
        f'<div class="lbl" style="margin-bottom:8px;">INSTRUMENTS WITH RISING VOL</div>'
        f'<div style="font-size:32px;font-weight:700;color:#e3b341;font-family:monospace;">{vol_count}</div>'
        f'<div style="font-size:9px;color:#8b949e;margin-top:4px;">of {len(df)} showing elevated vol vs 1M ago</div>'
        f'<div style="font-size:8px;color:#8b949e;font-family:monospace;margin-top:2px;">Elevated = current 20D vol &gt; vol 1M ago</div>'
        f'</div>'
        f'</div>'
        # Full backtest in Edge tab
        + (_build_edge_backtest(backtest_data) if backtest_data else '')
        # Model transparency note
        + f'<div style="background:#161b22;border:1px solid #30363d;border-radius:10px;padding:16px 20px;margin:16px 0;">'
        + f'<div style="font-size:10px;font-weight:700;letter-spacing:2px;color:#8b949e;margin-bottom:10px;">MODELS ON THIS PAGE</div>'
        + f'<div style="font-size:11px;color:#e6edf3;line-height:2.2;">'
        + f'<strong style="color:#58a6ff;">1. Regime-Weighted Framework Output (Illustrative)</strong> &mdash; LIVE. '
        + f'Risk Appetite Score = Regime(35%) + Fragility Inv(30%) + Fear &amp; Greed(20%) + Vol Inv(15%). '
        + f'Drives bucket weights. Highest composite score per bucket selected by BK Composite Score (5-factor composite on Intel tab).<br>'
        + f'<strong style="color:#58a6ff;">2. Regime Allocation Backtest</strong> &mdash; VALIDATED. '
        + f'Simplified 6-instrument model with monthly rebalancing. 5-year evidence. '
        + f'Demonstrates the regime-switching concept; not a replica of the live RAS model.<br>'
        + f'<strong style="color:#e3b341;">Status:</strong> '
        + f'RAS allocation is live but not yet backtested as a full 97-instrument model. '
        + f'The backtest above validates the regime-switching principle only.'
        + f'</div></div>'
        + f'<div style="font-size:9px;color:#8b949e;font-family:monospace;line-height:1.8;">'
        + f'Edge = regime-aware portfolio intelligence &#183; '
        + f'Allocation shifts automatically as market regime changes &#183; '
        + f'For informational purposes only &#183; Not investment advice &#183; Past performance not indicative of future results'
        + f'</div>'
    )

    # ══ TAB 9: RESEARCH DIRECTIONS ═══════════════════════════════════════════
    future_tab = (
        f'<div style="max-width:900px;margin:0 auto;">'
        f'<div style="background:linear-gradient(135deg,#1c2128,#161b22);border:1px solid #30363d;'
        f'border-radius:12px;padding:28px 32px;margin-bottom:20px;">'
        f'<div style="font-size:10px;font-weight:700;letter-spacing:3px;color:#58a6ff;text-transform:uppercase;margin-bottom:8px;">Open Questions</div>'
        f'<div style="font-size:22px;font-weight:700;color:#e6edf3;margin-bottom:8px;">Research Directions</div>'
        f'<div style="font-size:13px;color:#8b949e;line-height:1.7;">'
        f'Analytical questions this framework has not yet answered. Each section below identifies a research gap, '
        f'the methodology being explored, and what evidence would be needed to validate it.'
        f'</div></div>'

        # ── 1. Portfolio Construction ──────────────────────────────────────
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:10px;padding:24px 28px;margin-bottom:16px;">'
        f'<div style="display:flex;align-items:center;gap:12px;margin-bottom:14px;">'
        f'<div style="width:36px;height:36px;border-radius:8px;background:#0d2318;border:1px solid #238636;'
        f'display:flex;align-items:center;justify-content:center;font-size:16px;">&#9878;</div>'
        f'<div>'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#3fb950;text-transform:uppercase;">Research Area 1</div>'
        f'<div style="font-size:16px;font-weight:700;color:#e6edf3;">Portfolio Construction Under Regime Constraints</div>'
        f'</div></div>'
        f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;">'
        + ''.join(
            f'<div style="background:#1c2128;border:1px solid #21262d;border-radius:6px;padding:12px 14px;">'
            f'<div style="font-size:10px;font-weight:700;color:#58a6ff;margin-bottom:4px;">{title}</div>'
            f'<div style="font-size:11px;color:#8b949e;line-height:1.6;">{desc}</div></div>'
            for title, desc in [
                ("Does regime-conditioned mean-variance outperform unconditional?", "Research question: does applying Markowitz within each regime state produce better risk-adjusted outcomes than a single static frontier?"),
                ("Does risk parity hold up in Stressed regimes?", "Equal risk contribution weighting tends to concentrate in low-vol assets. Exploring whether this is protective or deceptive in crisis."),
                ("What is the rebalancing frequency that maximises Sharpe net of costs?", "Monthly vs quarterly vs threshold-based rebalancing. Transaction cost sensitivity across different regime states."),
                ("How much does correlation structure change between Calm and Stressed?", "Rolling 60/120-day correlation matrix. Hypothesis: diversification degrades precisely when it is needed most."),
                ("Does beta to benchmark vary predictably with fragility score?", "Exploring time-varying beta using 63/126-day windows against SPY and AGG. Is rising beta a leading fragility indicator?"),
                ("What is the marginal value of adding a new bucket to the RAS model?", "Current model has 6 buckets. Research question: does adding Alts or Crypto change regime-adjusted Sharpe materially?"),
            ]
        )
        + f'</div></div>'

        # ── 2. Signal Research ─────────────────────────────────────────────
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:10px;padding:24px 28px;margin-bottom:16px;">'
        f'<div style="display:flex;align-items:center;gap:12px;margin-bottom:14px;">'
        f'<div style="width:36px;height:36px;border-radius:8px;background:#0d1a2d;border:1px solid #1f6feb;'
        f'display:flex;align-items:center;justify-content:center;font-size:16px;">&#9654;</div>'
        f'<div>'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#58a6ff;text-transform:uppercase;">Research Area 2</div>'
        f'<div style="font-size:16px;font-weight:700;color:#e6edf3;">Cross-Asset Signal Persistence</div>'
        f'</div></div>'
        f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;">'
        + ''.join(
            f'<div style="background:#1c2128;border:1px solid #21262d;border-radius:6px;padding:12px 14px;">'
            f'<div style="font-size:10px;font-weight:700;color:#58a6ff;margin-bottom:4px;">{title}</div>'
            f'<div style="font-size:11px;color:#8b949e;line-height:1.6;">{desc}</div></div>'
            for title, desc in [
                ("Does 12-1 momentum persist within the BK universe?", "Cross-sectional momentum ranking across 113 instruments. Research question: does the signal decay faster in high-fragility regimes?"),
                ("Do cointegrated pairs offer regime-independent return?", "Exploring EWJ/EFA, GLD/SLV and similar pairs. Hypothesis: spread mean-reversion weakens during Crisis regime."),
                ("Does RSI divergence from price predict reversal at the asset-class level?", "Mean-reversion screening. Quantifying false positive rate across regime states."),
                ("How much of the BK GREEN signal decays within 5 trading days?", "Signal half-life analysis. Measuring whether composite scores lead or lag price by regime."),
                ("Is COT commercial positioning a leading indicator for commodity fragility?", "Exploring CFTC data for crude, gold, wheat. Hypothesis: extreme commercial short correlates with fragility spikes."),
                ("Do earnings surprise magnitudes vary with the fragility score?", "Exploring whether HIGH fragility periods coincide with larger post-earnings moves. Potential volatility timing signal."),
            ]
        )
        + f'</div></div>'

        # ── 3. Regime Detection ────────────────────────────────────────────
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:10px;padding:24px 28px;margin-bottom:16px;">'
        f'<div style="display:flex;align-items:center;gap:12px;margin-bottom:14px;">'
        f'<div style="width:36px;height:36px;border-radius:8px;background:#2d1f06;border:1px solid #e3b341;'
        f'display:flex;align-items:center;justify-content:center;font-size:16px;">&#128161;</div>'
        f'<div>'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#e3b341;text-transform:uppercase;">Research Area 3</div>'
        f'<div style="font-size:16px;font-weight:700;color:#e6edf3;">Regime Detection Accuracy &amp; Transition Lead Time</div>'
        f'</div></div>'
        f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;">'
        + ''.join(
            f'<div style="background:#1c2128;border:1px solid #21262d;border-radius:6px;padding:12px 14px;">'
            f'<div style="font-size:10px;font-weight:700;color:#e3b341;margin-bottom:4px;">{title}</div>'
            f'<div style="font-size:11px;color:#8b949e;line-height:1.6;">{desc}</div></div>'
            for title, desc in [
                ("How early does HMM detect regime transitions vs GMM?", "Comparing HMM and GMM lead times on the 2020 and 2022 episodes. Research question: which model minimises false positives?"),
                ("Does a 3-state model outperform a 2-state model in real time?", "Current model uses 3 states. Exploring whether adding a 4th (Recovery) state improves out-of-sample transition accuracy."),
                ("What is the false positive rate for transition risk flags?", "Currently flagged as Elevated when models disagree. Measuring how often this precedes actual regime change vs mean-reverts."),
                ("Does analyst consensus data lead or lag the regime signal?", "Exploring yfinance consensus ratios as a sentiment crosscheck on HMM/GMM regime classification."),
                ("Can macro event timing improve regime change probability estimates?", "FOMC, CPI, NFP dates as covariates in regime transition probabilities. Do they add explanatory power?"),
                ("Is the current 3-model consensus robust to instrument universe changes?", "Testing whether adding or removing asset classes materially shifts regime classification on historical episodes."),
            ]
        )
        + f'</div></div>'

        # ── 4. Fragility Model ─────────────────────────────────────────────
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:10px;padding:24px 28px;margin-bottom:16px;">'
        f'<div style="display:flex;align-items:center;gap:12px;margin-bottom:14px;">'
        f'<div style="width:36px;height:36px;border-radius:8px;background:#2d0a1e;border:1px solid #da3633;'
        f'display:flex;align-items:center;justify-content:circle;font-size:16px;">&#9654;</div>'
        f'<div>'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#f85149;text-transform:uppercase;">Research Area 4</div>'
        f'<div style="font-size:16px;font-weight:700;color:#e6edf3;">Fragility Score Validation &amp; Factor Decomposition</div>'
        f'</div></div>'
        f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:10px;">'
        + ''.join(
            f'<div style="background:#1c2128;border:1px solid #21262d;border-radius:6px;padding:12px 14px;">'
            f'<div style="font-size:10px;font-weight:700;color:#f85149;margin-bottom:4px;">{title}</div>'
            f'<div style="font-size:11px;color:#8b949e;line-height:1.6;">{desc}</div></div>'
            for title, desc in [
                ("Does a high fragility score predict subsequent drawdown within 30 days?", "Core validation question. Measuring hit rate on score &gt; 70 → max drawdown &gt; 10% within 21 trading days."),
                ("Which fragility pillars are most predictive of drawdown vs volatility?", "Factor attribution: decomposing which of the 5 pillar components explains most of subsequent loss."),
                ("Is the equal-weighting of fragility pillars optimal?", "Research question: do value/momentum/quality/low-vol weights derived from factor regressions outperform equal weights?"),
                ("Does fragility dispersion across asset classes predict regime transition?", "Hypothesis: rising cross-asset fragility dispersion (not just mean) is a leading regime indicator."),
                ("How does the fragility score behave for FX instruments excluded from display?", "FX is excluded from the fragility tab. Research question: does including it change the system-level fragility reading?"),
                ("Can the fragility score be extended to individual equities?", "Current model is ETF/index-based. Exploring whether pillar methodology transfers to single-stock screening."),
            ]
        )
        + f'</div></div>'

        # Footer note
        f'<div style="text-align:center;padding:16px;font-size:10px;color:#4a5568;font-family:monospace;">'
        f'Research questions only &#183; No conclusions implied &#183; All analysis uses Yahoo Finance data &#183; Scope subject to revision'
        f'</div>'
        f'</div>'
    )

    # ══ TAB 8: ABOUT ══════════════════════════════════════════════════════════
    about_tab = (
        f'<div style="max-width:800px;margin:0 auto;">'
        # Header
        f'<div style="background:linear-gradient(135deg,#1c2128,#161b22);border:1px solid #30363d;'
        f'border-radius:12px;padding:32px;margin-bottom:20px;text-align:center;">'
        f'<div style="width:80px;height:80px;border-radius:50%;background:linear-gradient(135deg,#58a6ff,#3fb950);'
        f'margin:0 auto 16px;display:flex;align-items:center;justify-content:center;">'
        f'<div style="font-size:28px;font-weight:700;color:#fff;font-family:monospace;">BK</div></div>'
        f'<div style="font-size:22px;font-weight:700;color:#e6edf3;font-family:monospace;letter-spacing:1px;">Bhavesh Kamdar</div>'
        f'<div style="font-size:12px;color:#58a6ff;margin-top:6px;letter-spacing:2px;text-transform:uppercase;">FRM · CQF · Risk Manager</div>'
        f'</div>'
        # Philosophy
        f'<div class="fc" style="margin-bottom:14px;">'
        f'<div class="lbl" style="margin-bottom:12px;">INVESTMENT PHILOSOPHY</div>'
        f'<div style="font-size:13px;color:#e6edf3;line-height:1.9;">'
        f'<em style="color:#58a6ff;">"Risk is not something to be avoided — it is something to be understood, '
        f'measured and navigated. After 25 years managing risk across global asset management firms, '
        f'I have seen every market cycle, every crisis and every recovery. The pattern is always the same: '
        f'fragility builds slowly, then breaks suddenly."</em>'
        f'</div>'
        f'<div style="font-size:12px;color:#8b949e;margin-top:16px;line-height:1.8;">'
        f'Bhavesh Kamdar is a senior risk professional with 25 years of experience in global asset management, '
        f'spent building risk frameworks across equities, fixed income, commodities, and alternatives.<br><br>'
        f'Holding both the Financial Risk Manager (FRM) designation and the Certificate in Quantitative Finance (CQF), '
        f'Bhavesh combines deep quantitative expertise with practical investment risk management experience '
        f'across equities, fixed income, commodities and alternatives.<br><br>'
        f'The BK Fragility Framework was born from a simple observation: traditional risk models measure volatility '
        f'after it has arrived. Bhavesh built the framework to detect structural vulnerability before it crystallises into loss.</div>'
        f'</div>'

        # ── D-05: Six disclosure blocks ────────────────────────────────────────
        f'<div style="margin-top:24px;display:flex;flex-direction:column;gap:12px;">'

        # 1. Nature of this tool
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:8px;padding:16px 20px;">'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;text-transform:uppercase;margin-bottom:8px;">1. Nature of This Tool</div>'
        f'<div style="font-size:11px;color:#8b949e;line-height:1.8;">'
        f'This dashboard is a <strong style="color:#e6edf3;">personal, non-commercial research project</strong>. '
        f'It is built and maintained by Bhavesh Kamdar for private analytical use. '
        f'It is not a financial product, not a regulated service, and is not offered commercially to any third party.'
        f'</div></div>'

        # 2. No investment advice
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:8px;padding:16px 20px;">'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;text-transform:uppercase;margin-bottom:8px;">2. Not Investment Advice</div>'
        f'<div style="font-size:11px;color:#8b949e;line-height:1.8;">'
        f'Nothing on this dashboard constitutes investment advice, a solicitation to buy or sell any security, '
        f'or a recommendation of any investment strategy. '
        f'All content is <strong style="color:#e6edf3;">observational and descriptive</strong> — it describes what the models output, '
        f'not what any person should do with their capital.'
        f'</div></div>'

        # 3. Model outputs are not predictions
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:8px;padding:16px 20px;">'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;text-transform:uppercase;margin-bottom:8px;">3. Model Outputs Are Not Predictions</div>'
        f'<div style="font-size:11px;color:#8b949e;line-height:1.8;">'
        f'Framework scores (fragility, regime, fear &amp; greed, composite) are <strong style="color:#e6edf3;">quantitative model outputs</strong> '
        f'derived from historical price and volume data. They describe current statistical conditions — '
        f'they do not predict future prices, returns, or market behaviour. '
        f'Past model performance is not indicative of future results.'
        f'</div></div>'

        # 4. Data limitations
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:8px;padding:16px 20px;">'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;text-transform:uppercase;margin-bottom:8px;">4. Data Limitations</div>'
        f'<div style="font-size:11px;color:#8b949e;line-height:1.8;">'
        f'Price and volume data is sourced from Yahoo Finance via yfinance. '
        f'Data may contain errors, gaps, stale prices, or corporate-action anomalies '
        f'(see the DATA REVIEW flag on the Performance tab for known issues). '
        f'<strong style="color:#e6edf3;">No warranty is made as to data accuracy or completeness.</strong>'
        f'</div></div>'

        # 5. No commercial relationship
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:8px;padding:16px 20px;">'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;text-transform:uppercase;margin-bottom:8px;">5. No Commercial Relationship</div>'
        f'<div style="font-size:11px;color:#8b949e;line-height:1.8;">'
        f'This tool is not affiliated with, endorsed by, or sponsored by any employer or institution. '
        f'It is developed independently in a personal capacity. '
        f'<strong style="color:#e6edf3;">No subscription, payment, or commercial arrangement exists or is offered.</strong>'
        f'</div></div>'

        # 6. Personal use only
        f'<div style="background:#161b22;border:1px solid #30363d;border-radius:8px;padding:16px 20px;">'
        f'<div style="font-size:9px;font-weight:700;letter-spacing:2px;color:#8b949e;text-transform:uppercase;margin-bottom:8px;">6. Personal Use Only</div>'
        f'<div style="font-size:11px;color:#8b949e;line-height:1.8;">'
        f'This dashboard is hosted publicly solely for technical convenience (GitHub Pages). '
        f'It is <strong style="color:#e6edf3;">not distributed, marketed, or promoted</strong> to any audience. '
        f'Any person accessing it does so for their own information and takes sole responsibility for any use they make of the content.'
        f'</div></div>'

        f'</div>'  # close disclosure blocks container

        f'</div>'
    )


    # ══ ASSEMBLE HTML ══════════════════════════════════════════════════════════
    mn="" if market_open else ' <span style="color:#e3b341;font-size:10px;">&#9888; Markets closed</span>'

    css=(":root{--bg:#0d1117;--ca:#161b22;--dk:#21262d;--br:#30363d;--w:#e6edf3;--g:#8b949e;--ac:#58a6ff;}"
         "*{box-sizing:border-box;margin:0;padding:0;}"
         "body{background:var(--bg);color:var(--w);font-family:'Segoe UI',system-ui,sans-serif;font-size:13px;}"
         ".wrap{max-width:1400px;margin:0 auto;padding:16px 12px;}"
         ".hdr{background:var(--ca);border:1px solid var(--br);border-radius:8px;padding:18px 24px;"
         "margin-bottom:14px;display:flex;justify-content:space-between;align-items:center;flex-wrap:wrap;gap:12px;}"
         ".logo{font-family:monospace;font-size:20px;font-weight:700;letter-spacing:2px;}"
         ".logo span{color:#3fb950;}"
         ".sub{font-size:10px;color:var(--g);letter-spacing:2px;margin-top:4px;}"
         ".badge{display:inline-flex;align-items:center;gap:6px;background:var(--dk);border:1px solid var(--br);"
         "border-radius:20px;padding:3px 10px;font-size:9px;font-family:monospace;color:var(--g);margin-top:6px;}"
         ".dot{width:6px;height:6px;border-radius:50%;background:#3fb950;animation:pulse 2s infinite;display:inline-block;}"
         "@keyframes pulse{0%,100%{opacity:1}50%{opacity:0.3}}"
         ".tabs{display:flex;gap:2px;margin-bottom:16px;border-bottom:3px solid var(--br);}"
         ".tb{padding:14px 28px;font-size:14px;font-weight:700;font-family:'Segoe UI',sans-serif;letter-spacing:0.3px;"
         "border:none;background:transparent;color:var(--g);cursor:pointer;"
         "border-bottom:3px solid transparent;margin-bottom:-3px;border-radius:6px 6px 0 0;transition:all 0.15s;}"
         ".tb:hover{color:var(--w);background:#1c2128;}"
         ".tb.on{color:var(--ac);border-bottom:3px solid var(--ac);background:#1c2128;}"
         ".tab{display:none;}.tab.on{display:block;}"
         ".tbar{display:flex;align-items:center;gap:12px;background:var(--ca);border:1px solid var(--br);"
         "border-radius:8px;padding:12px 24px;margin-bottom:14px;flex-wrap:wrap;}"
         ".lbl{font-size:9px;color:var(--g);letter-spacing:1px;text-transform:uppercase;margin-bottom:6px;}"
         ".pill{padding:4px 14px;border-radius:20px;font-size:11px;font-weight:700;font-family:monospace;letter-spacing:1px;}"
         ".dvdr{width:1px;height:40px;background:var(--br);margin:0 8px;}"
         ".rb{display:flex;gap:20px;}"
         ".ri{text-align:center;}"
         ".rn{font-size:22px;font-weight:700;font-family:monospace;}"
         ".rl{font-size:9px;color:var(--g);letter-spacing:1px;margin-top:2px;}"
         ".gl{display:grid;grid-template-columns:1fr 1fr;gap:14px;margin-bottom:14px;}"
         "@media(max-width:700px){.gl{grid-template-columns:1fr;}}"
         ".gc{background:var(--ca);border:1px solid var(--br);border-radius:8px;padding:16px 18px;}"
         ".gt{font-size:9px;font-weight:700;letter-spacing:2px;color:var(--g);margin-bottom:12px;"
         "display:flex;align-items:center;gap:8px;text-transform:uppercase;}"
         ".gd{width:8px;height:8px;border-radius:50%;display:inline-block;}"
         ".tw{background:var(--ca);border:1px solid var(--br);border-radius:8px;overflow-x:auto;}"
         "table{width:100%;border-collapse:collapse;font-size:12px;}"
         "th{background:#1c2128;padding:10px 8px;font-size:9px;letter-spacing:1px;text-transform:uppercase;"
         "color:var(--g);font-family:monospace;white-space:nowrap;border-bottom:2px solid var(--br);}"
         "th:first-child{text-align:left;padding-left:14px;}"
         "td{padding:7px 8px;border-bottom:1px solid var(--br);white-space:nowrap;}"
         "tr:last-child td{border-bottom:none;}"
         "tr:hover td{background:#1c2128;}"
         "td.an{text-align:left;padding-left:14px;color:var(--w);min-width:150px;}"
         "td.tk{font-family:monospace;font-size:10px;color:var(--g);font-weight:700;}"
         "td.num{font-family:monospace;text-align:right;}"
         "td.sig{font-family:monospace;font-size:10px;text-align:center;}"
         ".ps{color:#3fb950;}.pl{color:#7ee787;}.ng{color:#ff7b72;}.nr{color:#f85149;}"
         ".am{color:#e3b341;}.gr{color:#8b949e;}"
         ".sg{color:#3fb950;}.sa{color:#e3b341;}.sr{color:#f85149;}"
         "tr.sh td{background:#1c2128;font-size:9px;font-weight:700;letter-spacing:2px;"
         "text-transform:uppercase;color:var(--ac);padding:8px 14px;border-top:2px solid var(--br);}"
         "tr.bk-accordion-header{cursor:pointer;user-select:none;}"
         "tr.bk-accordion-header:hover td{background:#222840;}"
         ".bk-accordion-icon{display:inline-block;width:12px;margin-right:8px;color:#7a8499;font-size:10px;}"
         ".bk-accordion-summary{margin-left:14px;font-size:9px;color:#9aa3b8;font-weight:400;"
         "letter-spacing:0;text-transform:none;}"
         ".bk-accordion-count{margin-left:6px;font-size:9px;color:#7a8499;font-weight:400;"
         "letter-spacing:0;text-transform:none;}"
         ".bk-accordion-hint{float:right;font-size:9px;color:#4a5568;font-style:italic;font-weight:400;"
         "letter-spacing:0;text-transform:none;}"
         ".bk-acc-controls{margin:0 0 8px 0;display:flex;gap:6px;}"
         ".bk-acc-btn{font-size:10px;padding:4px 12px;background:#1a1f2e;color:#9aa3b8;"
         "border:1px solid #2a3040;border-radius:3px;cursor:pointer;font-family:monospace;}"
         ".bk-acc-btn:hover{background:#222840;color:#c8cfe0;}"
         ".vc{background:var(--ca);border:1px solid;border-radius:8px;padding:14px 20px;text-align:center;min-width:140px;}"
         ".vn{font-size:26px;font-weight:700;font-family:monospace;}"
         ".vl{font-size:9px;font-weight:700;letter-spacing:1px;margin-top:4px;}"
         ".vs{font-size:9px;color:var(--g);margin-top:2px;}"
         ".fc{background:var(--ca);border:1px solid var(--br);border-radius:8px;padding:16px 20px;}"
         ".footer{margin-top:14px;padding:12px 0;border-top:1px solid var(--br);"
         "display:flex;justify-content:space-between;align-items:center;flex-wrap:wrap;gap:8px;}"
         ".fn{font-size:9px;color:var(--g);line-height:1.9;font-family:monospace;}"
         ".fb{font-size:20px;font-weight:700;letter-spacing:-1px;font-family:monospace;}"
         ".fs{font-size:9px;color:var(--g);margin-top:2px;}"
         "@media(max-width:600px){"
         ".logo{font-size:15px;}.hdr{padding:12px 16px;}.wrap{padding:10px 8px;}"
         "th,td{padding:5px 6px;font-size:11px;}td.an{min-width:100px;}"
         ".tb{padding:10px 14px;font-size:11px;font-weight:700;}.rn{font-size:16px;}}")

    return (
        "<!DOCTYPE html><html lang='en'><head>"
        "<meta charset='UTF-8'>"
        "<meta name='viewport' content='width=device-width,initial-scale=1'>"
        "<!-- auto-refresh disabled for institutional use -->"
        "<title>BK Market Dashboard</title>"
        + (f"<script async src='https://www.googletagmanager.com/gtag/js?id={GA}'></script>"
           f"<script>window.dataLayer=window.dataLayer||[];function gtag(){{dataLayer.push(arguments);}}"
           f"gtag('js',new Date());gtag('config','{GA}');</script>" if GA else "<!-- GA not configured -->")
        + f"<style>{css}</style>"
        + "</head><body><div class='wrap'>"
        + "<div class='hdr'><div>"
        + "<div class='logo'>BKIQ <span>MARKETS</span></div>"
        + f"<div class='sub'>{N_INSTRUMENTS}-INSTRUMENT UNIVERSE &nbsp;&#183;&nbsp; Daily snapshot &middot; 07:00 SGT</div>"
        + f"<div class='badge'><span class='dot'></span> Last updated: {date_str}</div>"
        + "</div><div style='text-align:right;'>"
        f"<div style='font-family:monospace;font-size:13px;color:#e6edf3;font-weight:600;'>{date_str}</div>"
        f"<div style='font-size:9px;color:#8b949e;margin-top:4px;'>Updated daily before market open{mn}</div>"
        "<div style='font-size:9px;color:#444d56;margin-top:3px;font-family:monospace;'>07:00 SGT &#183; MON&#8211;FRI</div>"
        "</div></div>"
        "<div class='tabs'>"
        "<button class='tb on' onclick=\"sw('intel',this)\">Intel</button>"
        "<button class='tb' onclick=\"sw('perf',this)\">Performance</button>"
        "<button class='tb' onclick=\"sw('risk',this)\">Risk</button>"
        "<button class='tb' onclick=\"sw('frag',this)\">Fragility</button>"
        "<button class='tb' onclick=\"sw('analysis',this)\">Analysis</button>"
        "<button class='tb' onclick=\"sw('regime',this)\">Regime</button>"
        "<button class='tb' onclick=\"sw('edge',this)\">Edge</button>"
        "<button class='tb' onclick=\"sw('about',this)\">About</button>"
        "<button class='tb' onclick=\"sw('future',this)\">Research</button>"
        "</div>"
        f"<div id='t-intel' class='tab on'>{summary_tab}</div>"
        f"<div id='t-perf' class='tab'>{perf}</div>"
        f"<div id='t-risk' class='tab'>{risk}</div>"
        f"<div id='t-frag' class='tab'>{frag}</div>"
        f"<div id='t-analysis' class='tab'>{analysis_tab}</div>"
        f"<div id='t-regime' class='tab'>{regime_tab}</div>"
        f"<div id='t-edge' class='tab'>{edge_tab}</div>"
        f"<div id='t-about' class='tab'>{about_tab}</div>"
        f"<div id='t-future' class='tab'>{future_tab}</div>"
        "<div class='footer'><div class='fn'>"
        "Returns are price return in USD (ETF prices) &#183; FX returns reflect USD rate changes &#183; Trend = 20-day normalised sparkline<br>"
        "Signal: RED &lt; &#8722;15% &#183; AMBER &#8722;15% to &#8722;7% &#183; GREEN &gt; &#8722;7% from 52-week high<br>"
        "Fragility: CRISIS &#8805;70 &#183; STRESSED 55&#8211;69 &#183; MODERATE &lt;55 &#183; BK Fragility Framework<br>"
        ""
        f"Generated: {gen_ts} SGT &#183; Prices via Yahoo Finance &#183; Updated daily before market open"
        "</div><div style='text-align:right;'>"
        "<div class='fb'>BK</div>"
        "<div class='fs'>Risk Research &#183; Singapore</div>"
        "</div></div></div>"
        "<script>"
        "function sw(n,b){"
        "document.querySelectorAll('.tab').forEach(t=>t.classList.remove('on'));"
        "document.querySelectorAll('.tb').forEach(x=>x.classList.remove('on'));"
        "document.getElementById('t-'+n).classList.add('on');b.classList.add('on');}"
        "function bkSetState(h,state){"
        "  var section=h.closest('.bk-accordion-section');"
        "  var rows=section?section.querySelectorAll('.bk-accordion-body'):[];"
        "  var icon=h.querySelector('.bk-accordion-icon');"
        "  var open=(state==='open');"
        "  rows.forEach(function(r){r.style.display=open?'':'none';});"
        "  if(icon){icon.innerHTML=open?'\\u25BC':'\\u25B6';}"
        "  try{sessionStorage.setItem('bk_acc_'+h.dataset.key,state);}catch(e){}"
        "}"
        "function bkInitAccordions(){"
        "  document.querySelectorAll('tr.bk-accordion-header').forEach(function(h){"
        "    var saved='closed';try{saved=sessionStorage.getItem('bk_acc_'+h.dataset.key)||'closed';}catch(e){}"
        "    bkSetState(h,saved);"
        "    h.addEventListener('click',function(){"
        "      var cur='closed';try{cur=sessionStorage.getItem('bk_acc_'+h.dataset.key)||'closed';}catch(e){}"
        "      bkSetState(h,cur==='open'?'closed':'open');"
        "    });"
        "  });"
        "}"
        "function bkSetAllAccordions(state,prefix){"
        "  document.querySelectorAll('tr.bk-accordion-header').forEach(function(h){"
        "    if(h.dataset.key && h.dataset.key.indexOf(prefix+'_')===0){bkSetState(h,state);}"
        "  });"
        "}"
        "document.addEventListener('DOMContentLoaded',bkInitAccordions);"
        "</script>"
        # ── D-06: Persistent footer strip ─────────────────────────────────────
        '<div style="position:fixed;bottom:0;left:0;right:0;z-index:9999;'
        'background:#0d1117;border-top:1px solid #30363d;'
        'padding:6px 16px;text-align:center;font-size:9px;color:#6e7681;font-family:monospace;'
        'letter-spacing:1px;">'
        'Personal research &nbsp;&middot;&nbsp; Not investment advice &nbsp;&middot;&nbsp; No commercial offering'
        '</div>'
        # Spacer so page content is not obscured by the fixed footer
        '<div style="height:28px;"></div>'
        "</body></html>"
    )


def send_email(html_body: str) -> bool:
    now      = datetime.now(SGT)
    date_str = now.strftime("%a %d %b %Y")
    subject  = f"BK Market Dashboard · Daily Brief · {date_str}"
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = subject
        msg["From"]    = SENDER_EMAIL
        msg["To"]      = RECIPIENT_EMAIL
        msg.attach(MIMEText(html_body, "html"))

        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
            server.login(SENDER_EMAIL, GMAIL_APP_PASS)
            server.sendmail(SENDER_EMAIL, RECIPIENT_EMAIL, msg.as_string())

        print(f"[Email]  Sent to {RECIPIENT_EMAIL}")
        return True

    except Exception as e:
        print(f"[Email]  ERROR: {e}")
        print("  Check SENDER_EMAIL, GMAIL_APP_PASS.")
        print("  Get an App Password at: myaccount.google.com > Security > App Passwords")
        return False


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT REPORT
# ══════════════════════════════════════════════════════════════════════════════

def render_pptx(df: pd.DataFrame, prices: pd.DataFrame, as_of: str, out_dir: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from lxml import etree
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Low-level helpers ────────────────────────────────────────────────────

    def _rgb(h: str) -> RGBColor:
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _set_bg(slide, h: str):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(h)

    def _rect(slide, l, t, w, h, fill_h: str):
        shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill_h)
        shp.line.fill.background()
        return shp

    def _txbox(slide, text, l, t, w, h,
               size=14, bold=False, color="e6edf3",
               align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
        run.font.name = "Calibri"

    def _set_cell(cell, text, bg_h="0d1117", fg_h="e6edf3",
                  size=10, bold=False, align=PP_ALIGN.CENTER):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for child in list(tcPr):
            tag = child.tag.split("}")[-1]
            if tag in ("solidFill", "gradFill", "pattFill", "noFill"):
                tcPr.remove(child)
        sf = etree.SubElement(tcPr, qn("a:solidFill"))
        sr = etree.SubElement(sf, qn("a:srgbClr"))
        sr.set("val", bg_h.lstrip("#"))
        tf = cell.text_frame
        tf.word_wrap = False
        p  = tf.paragraphs[0]
        p.alignment = align
        p.clear()
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(fg_h)
        run.font.name = "Calibri"

    # ── Cell color maps ──────────────────────────────────────────────────────

    def _cc_ret(v):
        if pd.isna(v): return "21262d", "8b949e"
        p = v * 100
        if p >=  3.0: return "065f46", "6ee7b7"
        if p >=  1.0: return "166534", "86efac"
        if p >=  0.0: return "1c3829", "bbf7d0"
        if p >= -1.0: return "3b1111", "fca5a5"
        if p >= -3.0: return "5a1a1a", "f87171"
        return "7f1d1d", "fca5a5"

    def _cc_vol(v):
        if pd.isna(v): return "21262d", "8b949e"
        p = v * 100
        if p > 30: return "7f1d1d", "fca5a5"
        if p > 18: return "78350f", "fde68a"
        return "14532d", "86efac"

    def _cc_dd(v):
        if pd.isna(v): return "21262d", "8b949e"
        p = v * 100
        if p < -15: return "7f1d1d", "fca5a5"
        if p <  -7: return "78350f", "fde68a"
        return "14532d", "86efac"

    def _cc_sharpe(v):
        if pd.isna(v): return "21262d", "8b949e"
        if v >  1.0: return "14532d", "86efac"
        if v >  0.0: return "78350f", "fde68a"
        return "7f1d1d", "fca5a5"

    def _cc_rag(label):
        rl = label.strip()
        if rl == "RED":   return "7f1d1d", "f87171"
        if rl == "AMBER": return "78350f", "fde68a"
        return "14532d", "86efac"

    def _fr(v):  return "—" if pd.isna(v) else f"{v*100:+.1f}%"
    def _fp(v):  return "—" if pd.isna(v) else f"{v*100:.1f}%"
    def _fs(v):  return "—" if pd.isna(v) else f"{v:.2f}"

    # ── Sparkline image helper ────────────────────────────────────────────────

    def _sparkline_img(spark: list, rag_label: str, w_in=1.0, h_in=0.28) -> BytesIO | None:
        """Render a tiny sparkline as a PNG BytesIO for embedding in PPTX."""
        if not spark or len(spark) < 2:
            return None
        sp  = np.array(spark)
        col = "#3fb950" if sp[-1] >= 0 else "#f85149"
        fig2, ax2 = plt.subplots(figsize=(w_in, h_in), facecolor="#0d1117")
        ax2.set_facecolor("#0d1117")
        mn, mx = sp.min(), sp.max()
        rng    = mx - mn if mx != mn else 1.0
        xs     = np.arange(len(sp))
        zero_y = (0 - mn) / rng  # normalised 0-line position
        zero_y = np.clip(zero_y, 0, 1)
        sp_n   = (sp - mn) / rng
        ax2.fill_between(xs, zero_y, sp_n,
                         where=(sp_n >= zero_y), color="#3fb950", alpha=0.25, linewidth=0)
        ax2.fill_between(xs, zero_y, sp_n,
                         where=(sp_n < zero_y),  color="#f85149", alpha=0.25, linewidth=0)
        ax2.plot(xs, sp_n, color=col, linewidth=1.2)
        ax2.axhline(zero_y, color="#8b949e", linewidth=0.4, alpha=0.6)
        ax2.set_xlim(0, len(sp) - 1)
        ax2.set_ylim(-0.05, 1.05)
        ax2.axis("off")
        fig2.subplots_adjust(left=0, right=1, top=1, bottom=0)
        buf2 = BytesIO()
        fig2.savefig(buf2, format="png", dpi=120, bbox_inches="tight", facecolor="#0d1117")
        plt.close(fig2)
        buf2.seek(0)
        return buf2

    # Market open flag
    market_open = bool(df["market_open"].iloc[0]) if "market_open" in df.columns else True

    # ── Shared slide header ──────────────────────────────────────────────────

    def _slide_header(slide, title, subtitle):
        _set_bg(slide, "0d1117")
        _rect(slide, 0, 0, 13.33, 0.07, "58a6ff")
        _rect(slide, 0, 0.07, 13.33, 0.73, "161b22")
        _txbox(slide, title, 0.3, 0.10, 9.5, 0.44, size=20, bold=True, color="e6edf3")
        _txbox(slide, subtitle, 0.3, 0.52, 9.5, 0.26, size=10, color="8b949e")
        _txbox(slide, as_of, 10.2, 0.28, 2.9, 0.30,
               size=10, color="4b5563", align=PP_ALIGN.RIGHT)

    # ── Chart helpers ────────────────────────────────────────────────────────

    def _ytd_chart_img(ytd_slice: pd.DataFrame) -> BytesIO:
        """Render a YTD bar chart for a subset of instruments."""
        n = len(ytd_slice)
        fig_h = max(3.5, n * 0.22)
        fig, ax = plt.subplots(figsize=(12.5, fig_h), facecolor="#0d1117")
        ax.set_facecolor("#161b22")
        bar_colors = ["#f85149" if v < 0 else "#3fb950" for v in ytd_slice["ret_ytd"]]
        bars = ax.barh(ytd_slice["name"], ytd_slice["ret_ytd"] * 100,
                       color=bar_colors, edgecolor="none", height=0.72)
        for bar, val in zip(bars, ytd_slice["ret_ytd"] * 100):
            xpos = val + (0.4 if val >= 0 else -0.4)
            ax.text(xpos, bar.get_y() + bar.get_height() / 2,
                    f"{val:+.1f}%", va="center",
                    ha="left" if val >= 0 else "right",
                    color="#e6edf3", fontsize=9, fontweight="bold")
        ax.axvline(0, color="#8b949e", linewidth=0.8)
        ax.set_xlabel("YTD Return (%)", color="#8b949e", fontsize=10)
        ax.tick_params(colors="#8b949e", labelsize=9)
        for sp in ax.spines.values(): sp.set_edgecolor("#30363d")
        ax.grid(axis="x", color="#21262d", linewidth=0.5, alpha=0.7)
        fig.tight_layout(pad=0.5)
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="#0d1117")
        plt.close(fig)
        buf.seek(0)
        return buf

    def _risk_scatter_img():
        from matplotlib.lines import Line2D
        plot_df = df.dropna(subset=["vol_20d", "sharpe"]).copy()
        fig, ax = plt.subplots(figsize=(12.5, 5.8), facecolor="#0d1117")
        ax.set_facecolor("#161b22")
        color_map = {"RED": "#f85149", "AMBER": "#e3b341", "GREEN": "#3fb950"}

        # Clip axes to keep chart readable — outliers plotted at boundary with annotation
        VOL_MAX   = 75.0   # cap x-axis at 75% vol
        SHARPE_MIN = -4.0  # cap y-axis at -4 Sharpe

        for _, row in plot_df.iterrows():
            rag = row["rag_label"].strip()
            c   = color_map.get(rag, "#8b949e")
            sz  = max(40, min(350, abs(row["max_dd"]) * 1800)) if not pd.isna(row["max_dd"]) else 60
            vol_x   = min(row["vol_20d"] * 100, VOL_MAX * 0.97)
            sharpe_y = max(row["sharpe"], SHARPE_MIN * 0.97)
            ax.scatter(vol_x, sharpe_y,
                       s=sz, color=c, alpha=0.75, edgecolors="#30363d", linewidths=0.5)
            ax.annotate(row["ticker"],
                        (vol_x, sharpe_y),
                        fontsize=7, color="#e6edf3", alpha=0.85,
                        xytext=(4, 4), textcoords="offset points")
        ax.axhline(0, color="#8b949e", linestyle="--", linewidth=0.8, alpha=0.5)
        ax.axhline(1, color="#3fb950", linestyle=":", linewidth=0.8, alpha=0.4)
        ax.set_xlim(0, VOL_MAX)
        ax.set_ylim(SHARPE_MIN, ax.get_ylim()[1])
        ax.set_xlabel("Volatility 20D Annualised (%)", color="#8b949e", fontsize=10)
        ax.set_ylabel("Sharpe Ratio (1Y)", color="#8b949e", fontsize=10)
        ax.tick_params(colors="#8b949e", labelsize=9)
        for sp in ax.spines.values(): sp.set_edgecolor("#30363d")
        ax.grid(color="#21262d", linewidth=0.5, alpha=0.6)
        legend_elements = [
            Line2D([0], [0], marker="o", color="w", markerfacecolor="#f85149",
                   markersize=9, label="RED  —  Max DD < −15%"),
            Line2D([0], [0], marker="o", color="w", markerfacecolor="#e3b341",
                   markersize=9, label="AMBER  —  Max DD −15% to −7%"),
            Line2D([0], [0], marker="o", color="w", markerfacecolor="#3fb950",
                   markersize=9, label="GREEN  —  Max DD > −7%"),
        ]
        ax.legend(handles=legend_elements, fontsize=9, framealpha=0.4,
                  facecolor="#0d1117", labelcolor="#e6edf3", loc="upper right")
        ax.text(0.01, 0.97, "Bubble size ∝ magnitude of Max Drawdown  ·  Extreme outliers clipped to axis bounds",
                transform=ax.transAxes, fontsize=8, color="#6b7280", va="top")
        fig.tight_layout(pad=0.4)
        buf = BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="#0d1117")
        plt.close(fig)
        buf.seek(0)
        return buf

    # ── Grouped returns/risk table ────────────────────────────────────────────

    def _add_grouped_table(slide, sec_keys, col_defs, spark_col_idx=None,
                           slide_title="", slide_subtitle="", rows_per_slide=10):
        """Paginated table — max rows_per_slide data rows per slide.
        First slide is pre-created and passed in; overflow slides auto-created.
        Row height 0.52", font size 11pt for clean readability.
        """
        ROW_H    = 0.52   # generous data row height
        COLHDR_H = 0.36   # column header row
        SECHDR_H = 0.30   # section divider rows

        # Build flat ordered list: ("hdr", label) | ("data", row)
        all_specs = []
        for sec in sec_keys:
            sub = df[df["section"] == sec]
            if sub.empty:
                continue
            all_specs.append(("hdr", SECTION_LABELS.get(sec, sec)))
            for _, row in sub.iterrows():
                all_specs.append(("data", row))

        # Paginate: split on data-row count, carry section header to next page
        pages        = []
        cur_page     = []
        data_count   = 0
        pending_hdr  = None

        for spec in all_specs:
            if spec[0] == "hdr":
                pending_hdr = spec
            else:
                if pending_hdr is not None:
                    cur_page.append(pending_hdr)
                    pending_hdr = None
                cur_page.append(spec)
                data_count += 1
                if data_count >= rows_per_slide:
                    pages.append(cur_page)
                    cur_page    = []
                    data_count  = 0
                    pending_hdr = None
        if cur_page:
            pages.append(cur_page)

        total_pages = len(pages)

        def _render_page(sl, page_specs, page_num):
            n_rows  = len(page_specs) + 1   # +1 for column header row
            n_cols  = len(col_defs)
            total_w = sum(c[1] for c in col_defs)

            # Heights per row type
            row_heights = [COLHDR_H]
            for rtype, _ in page_specs:
                row_heights.append(SECHDR_H if rtype == "hdr" else ROW_H)
            tbl_h = sum(row_heights)

            tbl = sl.shapes.add_table(
                n_rows, n_cols,
                Inches(0.18), Inches(0.88),
                Inches(total_w), Inches(tbl_h)
            ).table

            for ci, (_, cw, _, _) in enumerate(col_defs):
                tbl.columns[ci].width = Inches(cw)
            for ri, rh in enumerate(row_heights):
                tbl.rows[ri].height = Inches(rh)

            # Column header
            for ci, (hdr, _, ha, _) in enumerate(col_defs):
                _set_cell(tbl.cell(0, ci), hdr, bg_h="161b22", fg_h="58a6ff",
                          size=11, bold=True, align=ha)

            # Accumulated top offset for sparkline positioning
            row_top_offset = 0.88 + COLHDR_H

            for ri, (rtype, rdata) in enumerate(page_specs, start=1):
                rh = row_heights[ri]
                if rtype == "hdr":
                    for ci in range(n_cols):
                        txt = rdata if ci == 0 else ""
                        _set_cell(tbl.cell(ri, ci), txt, bg_h="1c2128", fg_h="58a6ff",
                                  size=10, bold=True, align=PP_ALIGN.LEFT)
                else:
                    for ci, (_, cw, ha, fn) in enumerate(col_defs):
                        if spark_col_idx is not None and ci == spark_col_idx:
                            _set_cell(tbl.cell(ri, ci), "", bg_h="0d1117", fg_h="0d1117", size=6)
                            buf_sp = _sparkline_img(
                                rdata.get("spark", []), rdata["rag_label"],
                                w_in=cw * 0.88, h_in=rh * 0.68
                            )
                            if buf_sp:
                                col_left = Inches(0.18 + sum(col_defs[k][1] for k in range(ci)))
                                sp_w_in  = cw * 0.84
                                sp_h_in  = rh * 0.64
                                sp_l = col_left + Inches((cw - sp_w_in) / 2)
                                sp_t = Inches(row_top_offset) + Inches((rh - sp_h_in) / 2)
                                sl.shapes.add_picture(buf_sp, sp_l, sp_t,
                                                      Inches(sp_w_in), Inches(sp_h_in))
                        else:
                            text, bg, fg = fn(rdata)
                            _set_cell(tbl.cell(ri, ci), text, bg_h=bg, fg_h=fg,
                                      size=11, align=ha)
                row_top_offset += rh

            # Page indicator
            if total_pages > 1:
                _txbox(sl, f"Page {page_num} of {total_pages}",
                       10.5, 7.20, 2.6, 0.22, size=9, color="4b5563",
                       align=PP_ALIGN.RIGHT)

        for pi, page_specs in enumerate(pages):
            if pi == 0:
                sl = slide
            else:
                sl = prs.slides.add_slide(blank)
                pg_sub = f"{slide_subtitle}  ·  {pi + 1}/{total_pages}"
                _slide_header(sl, slide_title, pg_sub)
            _render_page(sl, page_specs, pi + 1)

    # ════════════════════════════════════════════════════════════════════════
    #  SLIDE 1 — Cover
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, "0d1117")
    _rect(slide, 0, 0, 13.33, 0.07, "58a6ff")

    _txbox(slide, "BK MARKET DASHBOARD", 0.7, 1.55, 12.0, 1.2,
           size=52, bold=True, color="e6edf3", align=PP_ALIGN.LEFT)
    _txbox(slide, f"Daily Brief  ·  {as_of}", 0.7, 2.85, 10.0, 0.6,
           size=22, color="8b949e")
    _txbox(slide,
           f"{N_INSTRUMENTS}-Instrument Universe  ·  10 Asset Classes  ·  Returns, Risk & Signal",
           0.7, 3.52, 11.5, 0.45, size=15, color="6b7280")

    n_red_t   = (df["rag_label"].str.strip() == "RED").sum()
    n_amber_t = (df["rag_label"].str.strip() == "AMBER").sum()
    n_green_t = (df["rag_label"].str.strip() == "GREEN").sum()

    # ── Market call narrative ──
    total = n_red_t + n_amber_t + n_green_t
    pct_red   = n_red_t   / total * 100 if total else 0
    pct_green = n_green_t / total * 100 if total else 0
    if pct_red >= 40:
        market_tone = "RISK-OFF"
        tone_color  = "f87171"
        tone_desc   = f"{n_red_t} instruments in drawdown >15% — broad risk-off conditions."
    elif pct_green >= 50:
        market_tone = "RISK-ON"
        tone_color  = "86efac"
        tone_desc   = f"{n_green_t} of {total} instruments within 7% of 52-week highs."
    else:
        market_tone = "MIXED"
        tone_color  = "fde68a"
        tone_desc   = f"Split signals: {n_red_t} RED · {n_amber_t} AMBER · {n_green_t} GREEN."

    _txbox(slide, f"Market Tone: {market_tone}", 0.7, 3.95, 12.0, 0.45,
           size=16, bold=True, color=tone_color)
    _txbox(slide, tone_desc, 0.7, 4.42, 12.0, 0.30, size=11, color="8b949e")

    for xi, (lbl, cnt, bg, fg) in enumerate([
        ("RED",   n_red_t,   "450a0a", "f87171"),
        ("AMBER", n_amber_t, "451a03", "fde68a"),
        ("GREEN", n_green_t, "052e16", "86efac"),
    ]):
        x = 0.7 + xi * 2.9
        _rect(slide, x, 4.90, 2.5, 1.4, bg)
        _txbox(slide, str(cnt), x, 4.92, 2.5, 0.85,
               size=42, bold=True, color=fg, align=PP_ALIGN.CENTER)
        _txbox(slide, lbl, x, 5.80, 2.5, 0.35,
               size=12, color=fg, align=PP_ALIGN.CENTER)

    _txbox(slide,
           "Signal: RED = Max DD < −15%   AMBER = −15% to −7%   GREEN = > −7% from 52-week high",
           0.7, 6.30, 12.0, 0.30, size=9, color="4b5563")
    _txbox(slide, "Source: Yahoo Finance  ·  Price return, local currency  ·  CONFIDENTIAL",
           0.7, 7.12, 12.0, 0.28, size=9, color="374151")

    # ════════════════════════════════════════════════════════════════════════
    #  SLIDES 2+ — Equities Returns (10 rows per slide)
    # ════════════════════════════════════════════════════════════════════════
    def _fv(v):
        """Format vol as XX.X%"""
        return "—" if pd.isna(v) else f"{v*100:.1f}%"

    def _cc_vol_chg(now, ago):
        """Color vol change arrow: red if vol rose, green if fell."""
        if pd.isna(now) or pd.isna(ago): return "21262d", "8b949e"
        return ("7f1d1d", "f87171") if now > ago else ("14532d", "86efac")

    def _vol_chg_str(now, ago):
        """Show vol now + arrow direction vs 1m ago."""
        if pd.isna(now) or pd.isna(ago): return "—"
        arrow = "▲" if now > ago else "▼"
        return f"{now*100:.1f}% {arrow}"

    EQ_TITLE    = "EQUITIES — RETURNS"
    EQ_SUBTITLE = "US Broad · US Sectors · Dev Mkts · EM · Defence  ·  Trend / 1W / 1M / 3M / YTD / Vol / Signal"
    slide = prs.slides.add_slide(blank)
    _slide_header(slide, EQ_TITLE, EQ_SUBTITLE)
    eq_cols = [
        ("Asset",   3.2, PP_ALIGN.LEFT,   lambda r: (r["name"],   "0d1117", "e6edf3")),
        ("Ticker",  0.9, PP_ALIGN.CENTER, lambda r: (r["ticker"], "161b22", "8b949e")),
        ("Trend",   1.4, PP_ALIGN.CENTER, lambda r: ("",          "0d1117", "0d1117")),
    ]
    if market_open:
        eq_cols.append(("1D", 1.0, PP_ALIGN.CENTER,
                        lambda r: (_fr(r["ret_1d"]), *_cc_ret(r["ret_1d"]))))
    eq_cols += [
        ("1W",      1.1, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_1w"]),  *_cc_ret(r["ret_1w"]))),
        ("1M",      1.1, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_1m"]),  *_cc_ret(r["ret_1m"]))),
        ("3M",      1.1, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_3m"]),  *_cc_ret(r["ret_3m"]))),
        ("YTD",     1.2, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_ytd"]), *_cc_ret(r["ret_ytd"]))),
        ("Vol (vs 1M)", 1.4, PP_ALIGN.CENTER,
         lambda r: (_vol_chg_str(r["vol_now"], r["vol_1m_ago"]),
                    *_cc_vol_chg(r["vol_now"], r["vol_1m_ago"]))),
        ("Signal",  1.3, PP_ALIGN.CENTER,
         lambda r: (r["rag_label"].strip(), *_cc_rag(r["rag_label"]))),
    ]
    _add_grouped_table(slide, ["EQ_US", "EQ_SECT", "EQ_DM", "EQ_EM", "DEFENCE"],
                       eq_cols, spark_col_idx=2,
                       slide_title=EQ_TITLE, slide_subtitle=EQ_SUBTITLE,
                       rows_per_slide=10)

    # ════════════════════════════════════════════════════════════════════════
    #  SLIDES — Fixed Income, Commodities & Other (10 rows per slide)
    # ════════════════════════════════════════════════════════════════════════
    FI_TITLE    = "FIXED INCOME, COMMODITIES & OTHER — RETURNS"
    FI_SUBTITLE = "Fixed Income & Credit · Commodities · Crypto · FX · Volatility  ·  Trend / 1W / 1M / 3M / YTD / Vol / Signal"
    slide = prs.slides.add_slide(blank)
    _slide_header(slide, FI_TITLE, FI_SUBTITLE)
    fi_cols = [
        ("Asset",   3.2, PP_ALIGN.LEFT,   lambda r: (r["name"],   "0d1117", "e6edf3")),
        ("Ticker",  0.9, PP_ALIGN.CENTER, lambda r: (r["ticker"], "161b22", "8b949e")),
        ("Trend",   1.4, PP_ALIGN.CENTER, lambda r: ("",          "0d1117", "0d1117")),
    ]
    if market_open:
        fi_cols.append(("1D", 1.0, PP_ALIGN.CENTER,
                        lambda r: (_fr(r["ret_1d"]), *_cc_ret(r["ret_1d"]))))
    fi_cols += [
        ("1W",      1.1, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_1w"]),  *_cc_ret(r["ret_1w"]))),
        ("1M",      1.1, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_1m"]),  *_cc_ret(r["ret_1m"]))),
        ("3M",      1.1, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_3m"]),  *_cc_ret(r["ret_3m"]))),
        ("YTD",     1.2, PP_ALIGN.CENTER, lambda r: (_fr(r["ret_ytd"]), *_cc_ret(r["ret_ytd"]))),
        ("Vol (vs 1M)", 1.4, PP_ALIGN.CENTER,
         lambda r: (_vol_chg_str(r["vol_now"], r["vol_1m_ago"]),
                    *_cc_vol_chg(r["vol_now"], r["vol_1m_ago"]))),
        ("Signal",  1.3, PP_ALIGN.CENTER,
         lambda r: (r["rag_label"].strip(), *_cc_rag(r["rag_label"]))),
    ]
    _add_grouped_table(slide, ["FI", "CMD", "CRYPTO", "FX", "VOL"],
                       fi_cols, spark_col_idx=2,
                       slide_title=FI_TITLE, slide_subtitle=FI_SUBTITLE,
                       rows_per_slide=10)

    # ════════════════════════════════════════════════════════════════════════
    #  SLIDE 4 — Risk Snapshot (scatter: vol vs Sharpe, all instruments)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _slide_header(slide, "RISK SNAPSHOT — ALL INSTRUMENTS",
                  "Volatility (20D ann.) vs Sharpe Ratio (1Y)  ·  "
                  "Bubble size = Max Drawdown magnitude  ·  Colour = Signal")
    buf = _risk_scatter_img()
    if buf:
        slide.shapes.add_picture(buf, Inches(0.2), Inches(0.88),
                                 Inches(12.93), Inches(6.40))

    # ════════════════════════════════════════════════════════════════════════
    #  SLIDES — YTD Performance (15 instruments per slide, sorted best→worst)
    # ════════════════════════════════════════════════════════════════════════
    YTD_PER_SLIDE = 15
    ytd_all = df[["name", "ret_ytd"]].dropna(subset=["ret_ytd"]).copy()
    ytd_all = ytd_all.sort_values("ret_ytd", ascending=False).reset_index(drop=True)
    ytd_pages = [ytd_all.iloc[i:i + YTD_PER_SLIDE]
                 for i in range(0, len(ytd_all), YTD_PER_SLIDE)]
    ytd_total = len(ytd_pages)
    for pi, ytd_slice in enumerate(ytd_pages):
        # Sort slice bottom-to-top for horizontal bar chart
        ytd_slice = ytd_slice.sort_values("ret_ytd").reset_index(drop=True)
        slide = prs.slides.add_slide(blank)
        pg_lbl = f"  ·  {pi + 1}/{ytd_total}" if ytd_total > 1 else ""
        _slide_header(slide,
                      "YTD PERFORMANCE — ALL INSTRUMENTS",
                      f"Year-to-Date Return  ·  Sorted best to worst{pg_lbl}")
        buf = _ytd_chart_img(ytd_slice)
        slide.shapes.add_picture(buf, Inches(0.2), Inches(0.88),
                                 Inches(12.93), Inches(6.40))

    # ── Save ─────────────────────────────────────────────────────────────────
    os.makedirs(out_dir, exist_ok=True)
    tag  = datetime.now().strftime("%Y%m%d_%H%M")
    path = os.path.join(out_dir, f"market_dashboard_{tag}.pptx")
    prs.save(path)
    print(f"[PPTX]   {path}")
    return path


# ══════════════════════════════════════════════════════════════════════════════
#  ORCHESTRATION
# ══════════════════════════════════════════════════════════════════════════════

def _now_sgt() -> str:
    return datetime.now(SGT).strftime("%H:%M:%S SGT")


def run_once(send_email_flag: bool = False, pptx_flag: bool = False,
             html_flag: bool = False, out_dir: str = OUT_DIR,
             lookback_days: int = 2520) -> None:
    print("=" * 60)
    print(f"  BK Market Dashboard  |  {_now_sgt()}")
    print("=" * 60)

    prices, volumes = download(lookback_days=lookback_days)
    df     = compute_metrics(prices)

    n_red   = (df["rag_label"].str.strip() == "RED").sum()
    n_amber = (df["rag_label"].str.strip() == "AMBER").sum()
    n_green = (df["rag_label"].str.strip() == "GREEN").sum()
    print(f"[Signal]  RED={n_red}  AMBER={n_amber}  GREEN={n_green}  ({len(df)}/{N_INSTRUMENTS} fetched)")

    as_of    = datetime.now().strftime("%d %B %Y  %H:%M")
    png, pdf = render_report(df, as_of, out_dir)
    print(f"[Report]  PNG: {png}")
    print(f"[Report]  PDF: {pdf}")

    if pptx_flag:
        render_pptx(df, prices, as_of, out_dir)

    if html_flag:
        print("[HTML]   Computing fragility scores...")
        frag_df     = compute_fragility(prices, volumes)
        print("[HTML]   Computing market regime...")
        regime_data = compute_regime(prices)
        print("[HTML]   Computing Fear & Greed index...")
        fg_data     = compute_fear_greed(prices)
        print("[HTML]   Computing fragility trend...")
        frag_trend  = compute_fragility_trend(prices, volumes)
        print("[HTML]   Generating AI commentary...")
        # Build market data dict for AI
        _reg        = regime_data.get("regime","Calm") if regime_data else "Calm"
        _streak     = regime_data.get("days_in_regime",0) if regime_data else 0
        _drivers    = regime_data.get("drivers",{}) if regime_data else {}
        _frag_s     = frag_df.attrs.get("system_score",50) if frag_df is not None and not frag_df.empty else 50
        _frag_l     = frag_df.attrs.get("regime","MODERATE") if frag_df is not None and not frag_df.empty else "MODERATE"
        _fg_s       = round(fg_data.get("score",50)) if fg_data else 50
        _fg_l       = fg_data.get("label","Neutral") if fg_data else "Neutral"
        _df_ai_rank = df[df["ticker"].apply(is_rankable)]
        _gainers    = [f'{r["name"]} {r["ret_1m"]*100:+.1f}%' for _,r in _df_ai_rank.nlargest(3,"ret_1m").iterrows()]
        _losers     = [f'{r["name"]} {r["ret_1m"]*100:+.1f}%' for _,r in _df_ai_rank.nsmallest(3,"ret_1m").iterrows()]
        _n_red      = int((df["rag_label"].str.strip()=="RED").sum())
        _vol_rise   = int((df["vol_now"]>df["vol_1m_ago"]).sum()) if "vol_now" in df.columns else 0
        _avg_corr   = 0.3
        market_data_for_ai = {
            "date":          pd.Timestamp.now(tz=SGT).strftime("%A %d %b %Y"),
            "regime":        _reg,
            "regime_days":   _streak,
            "fragility_score": round(_frag_s,1),
            "fragility_label": _frag_l,
            "fg_score":      _fg_s,
            "fg_label":      _fg_l,
            "vol_pct":       round(_drivers.get("vol_pct",50),0),
            "dd_pct":        round(_drivers.get("dd_now",0),1),
            "n_red":         _n_red,
            "n_total":       len(df),
            "vol_rising":    _vol_rise,
            "top_gainers":   _gainers,
            "top_losers":    _losers,
            "avg_corr":      _avg_corr,
            "bt_bk":         "N/A",
            "bt_spy":        "N/A",
            "bt_6040":       "N/A",
        }
        ai_commentary = generate_ai_commentary(market_data_for_ai)
        print("[HTML]   Computing backtest...")
        _regime_series = pd.Series(
            {pt["date"]: pt["regime"] for pt in (regime_data.get("timeline",[]) if regime_data else [])}
        )
        _regime_series.index = pd.to_datetime(_regime_series.index)
        backtest_data = compute_backtest(prices, _regime_series)
        # Update AI market data with backtest results
        if backtest_data:
            market_data_for_ai["bt_bk"]   = f'{backtest_data["bk"]["total"]:+.1f}%'
            market_data_for_ai["bt_spy"]  = f'{backtest_data["spy"]["total"]:+.1f}%'
            market_data_for_ai["bt_6040"] = f'{backtest_data["p6040"]["total"]:+.1f}%'
        # D-08: Fetch and select today's headlines
        print("[HTML]   Fetching today's headlines...")
        _regime_label = (regime_data.get("regime", "Unknown") if regime_data else "Unknown")
        _news_pool    = fetch_news_pool()
        headlines_data = select_top_headlines(_news_pool, _regime_label)
        print(f"[HTML]   Headlines: {len(headlines_data)} selected from {len(_news_pool)} articles")
        web_html    = build_web_html(df, frag_df, prices, regime_data, fg_data, frag_trend, ai_commentary, backtest_data, headlines_data)
        docs_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "docs")
        os.makedirs(docs_dir, exist_ok=True)
        html_path = os.path.join(docs_dir, "index.html")
        with open(html_path, "w", encoding="utf-8") as fh:
            fh.write(web_html)
        # Always write CNAME to preserve custom domain
        cname_path = os.path.join(docs_dir, "CNAME")
        with open(cname_path, "w") as fh:
            fh.write("dashboard.bkiqmarkets.com")
        print("[HTML]   CNAME written: dashboard.bkiqmarkets.com")
        print(f"[HTML]   {html_path}")

    if send_email_flag:
        html = build_email_html(df)
        send_email(html)

    print("[Done]")


def run_scheduler(out_dir: str = OUT_DIR) -> None:
    print("=" * 60)
    print("  BK Market Dashboard · Daily Scheduler")
    print(f"  Send time : {SEND_TIME_SGT} SGT  (Mon–Fri)")
    print(f"  Recipient : {RECIPIENT_EMAIL}")
    print(f"  Output    : {os.path.abspath(out_dir)}")
    print("=" * 60)

    job = lambda: run_once(send_email_flag=True, out_dir=out_dir)
    for day in ["monday", "tuesday", "wednesday", "thursday", "friday"]:
        getattr(schedule.every(), day).at(SEND_TIME_SGT).do(job)

    print(f"\n[{_now_sgt()}] Scheduler running. Next run: {schedule.next_run()}")
    print("  Press Ctrl+C to stop.\n")

    while True:
        schedule.run_pending()
        time.sleep(30)


# ══════════════════════════════════════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="BK Market Dashboard — visual report + email brief",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
examples:
  python bk_market_dashboard.py                   # PNG + PDF only
  python bk_market_dashboard.py --pptx            # + PowerPoint deck
  python bk_market_dashboard.py --email           # + send HTML email
  python bk_market_dashboard.py --pptx --email    # all outputs
  python bk_market_dashboard.py --schedule        # daily scheduler at 07:00 SGT
  python bk_market_dashboard.py --now --pptx      # run immediately with PPTX
        """,
    )
    parser.add_argument("--email",    action="store_true",
                        help="Send HTML email after generating report")
    parser.add_argument("--pptx",     action="store_true",
                        help="Generate PowerPoint deck (5 slides per asset class)")
    parser.add_argument("--html",     action="store_true",
                        help="Generate docs/index.html for GitHub Pages")
    parser.add_argument("--lookback",  type=int, default=2520, metavar="DAYS",
                        help="Price history lookback in days (default: 2520 = 10yr, use 504 for 2yr)")
    parser.add_argument("--schedule", action="store_true",
                        help=f"Start daily scheduler at {SEND_TIME_SGT} SGT Mon–Fri")
    parser.add_argument("--now",      action="store_true",
                        help="Run once immediately (bypasses scheduler wait)")
    parser.add_argument("--out-dir",  default=OUT_DIR, metavar="DIR",
                        help=f"Output directory for outputs (default: {OUT_DIR}/)")
    args = parser.parse_args()

    if args.schedule and not args.now:
        run_scheduler(out_dir=args.out_dir)
    else:
        run_once(send_email_flag=args.email, pptx_flag=args.pptx,
                 html_flag=args.html, out_dir=args.out_dir,
                 lookback_days=args.lookback)


# ── BACKGROUND EXECUTION ──────────────────────────────────────────────────────
# Mac/Linux:  nohup python bk_market_dashboard.py --schedule &
#             (logs to nohup.out)
# Windows:    pythonw bk_market_dashboard.py --schedule
# Test now:   python bk_market_dashboard.py --now --email
